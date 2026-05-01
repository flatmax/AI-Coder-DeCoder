"""DocConvert RPC service — Layer 4.5.

The DocConvert service is the backend half of the Doc Convert tab.
It scans the repository for convertible documents (`.docx`, `.pdf`,
`.pptx`, `.xlsx`, `.csv`, `.rtf`, `.odt`, `.odp`), classifies each by
status (new / current / stale / conflict), and — on explicit user
request — converts them to markdown via a pipeline that dispatches
by file type through several optional dependencies (markitdown,
LibreOffice, PyMuPDF, python-pptx, openpyxl).

Current scope — **Pass A (foundation) + Pass A2 (markitdown for
simple formats)**. What's delivered:

- `DocConvert` class registered via `server.add_class(doc_convert)`.
- `is_available()` — probes every optional dependency and reports
  each via the specs4-specified shape `{available, libreoffice,
  pymupdf, pdf_pipeline}`.
- `scan_convertible_files()` — walks the repo, classifies each
  source file via the provenance-header round-trip.
- `convert_files(paths)` — Pass A2: dispatches by extension.
  `.docx`, `.rtf`, `.odt` route through markitdown with full
  provenance-header writing, data-URI image extraction, DOCX
  truncated-URI workaround, and orphan-image cleanup on
  re-conversion. Other extensions (`.pdf`, `.pptx`, `.xlsx`,
  `.csv`, `.odp`) return per-file "not yet supported" results
  — caller sees specific failures rather than a blanket
  NotImplementedError.
- Clean-tree gate — refuses conversion if the repo has
  uncommitted changes. Matches the Code Review prerequisite in
  specs4/4-features/doc-convert.md.

Pass A3 will add xlsx colour extraction via openpyxl. Pass A4
adds pptx fallback via python-pptx. Pass A5 adds the full PDF
pipeline (LibreOffice → PyMuPDF → text extraction + SVG export
+ image externalization).

Governing spec: ``specs4/4-features/doc-convert.md``.
Restriction pattern: ``specs4/1-foundation/communication-layer.md#restricted-operations``.

Design decisions pinned here:

- **Availability reports are authoritative.** The frontend decides
  whether to show the tab based on what this returns. A dependency
  that imports but doesn't actually work (e.g., PyMuPDF installed
  without its bundled binaries) should show as unavailable — we
  don't pre-import to verify, because import itself is cheap but
  running the converter isn't. Imports are the best cheap signal;
  true runtime availability surfaces via per-file conversion errors
  when the user actually clicks Convert.

- **LibreOffice probe uses PATH lookup, not subprocess.** Running
  `soffice --version` would be slow (LibreOffice starts a UNO
  listener on every invocation on some platforms) and unnecessary.
  `shutil.which("soffice")` returning non-None is a reliable
  indicator that the binary is installed; failures at conversion
  time produce clean error messages for the affected files.

- **Status classification priority.** For each source file, check:
  1. No output file exists → `new`
  2. Output exists, no docuvert header → `conflict`
  3. Output exists, header present, source hash matches → `current`
  4. Output exists, header present, source hash doesn't match → `stale`

  The `conflict` status is deliberately distinct from `stale` —
  specs4 is explicit that overwriting a file the user manually
  authored (or converted with a different tool) requires explicit
  opt-in. The UI highlights `conflict` files with a warning icon.

- **Directory exclusions mirror the indexers.** `.git`, `.ac-dc`,
  `node_modules`, `__pycache__`, `.venv`, `venv`, `dist`, `build`,
  `.egg-info`, and hidden directories (except `.github`). Imported
  duplicates would be brittle — the doc-convert scan has the same
  philosophy as the indexers (skip build artefacts, skip tooling
  directories) but is a separate code path, so we redefine the set
  rather than import.

- **Output path is the source stem + `.md` as a sibling.** `docs/
  architecture.docx` → `docs/architecture.md`. The per-source
  assets subdirectory (`docs/architecture/`) is created during
  conversion, not scan. Scan only cares about the `.md` output;
  assets subdirectory is an implementation detail of conversion.

- **Source size filter is advisory.** Files larger than
  `max_source_size_mb` appear in the scan with an over-size flag
  so the UI can show a warning badge. They're also rejected at
  conversion time. Including them in the scan (rather than
  excluding outright) means the user sees WHY a file isn't being
  converted, rather than wondering where it went.

- **Config is read-through, not snapshot.** The `enabled` flag,
  `extensions` list, and `max_source_size_mb` all come from
  `config.doc_convert_config` on every call. Hot-reloaded config
  values take effect immediately — useful during development and
  matches the pattern other services use.

- **`is_available()` is always callable.** Even when `enabled=false`
  in config, the availability probe runs — the frontend uses it to
  decide whether to show the tab at all, which is separate from the
  enabled flag. Enabled/disabled is a user opt-out; availability is
  a capability query.
"""

from __future__ import annotations

import base64
import hashlib
import logging
import os
import re
import shutil
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from ac_dc.config import ConfigManager

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Module constants
# ---------------------------------------------------------------------------


# Default extensions recognised as convertible. The config can
# override via `doc_convert.extensions`; this tuple is the fallback
# when config is absent or malformed. Matches specs4 exactly.
_DEFAULT_EXTENSIONS: tuple[str, ...] = (
    ".docx",
    ".pdf",
    ".pptx",
    ".xlsx",
    ".csv",
    ".rtf",
    ".odt",
    ".odp",
)


# Directories we never walk. Mirrors the indexers' exclusion list —
# rebuilt here rather than imported because the doc-convert scan is a
# separate code path from indexing, and coupling them would make the
# indexer refactorable-but-only-if-you-also-update-convert.
_EXCLUDED_DIRS: frozenset[str] = frozenset({
    ".git",
    ".ac-dc4",
    "node_modules",
    "__pycache__",
    ".venv",
    "venv",
    "dist",
    "build",
})


# Provenance-header regex. Matches the marker comment at the top of a
# converted file. Captures the whole body so we can parse the
# space-separated key=value pairs ourselves — re-using capture groups
# for each possible field would miss unknown ones, and we want to be
# forward-compatible with future field additions.
_PROVENANCE_RE = re.compile(
    r"<!--\s*docuvert:\s*([^>]+?)\s*-->",
    re.IGNORECASE,
)


# Pattern for one key=value pair inside the provenance body. Values
# are either bare tokens or comma-separated lists. We don't enforce
# field names here — the parser keeps unknown keys in the dict so
# future schema additions don't silently lose data on an older
# version reading a newer file.
_PROV_FIELD_RE = re.compile(r"(\w+)=([^\s]+)")


# How much of each output file we scan when looking for the
# provenance header. The header lives on the first line; a few
# hundred bytes is plenty, even accounting for a stray blank line
# or two preceding it.
_PROVENANCE_PROBE_BYTES = 2048


# Extensions handled by the markitdown path. `.csv` is included
# because markitdown produces clean markdown tables for simple
# CSVs — good enough until a dedicated pass lands if we ever
# need colour-aware CSV (unlikely; CSVs don't carry formatting).
_MARKITDOWN_EXTENSIONS: frozenset[str] = frozenset({
    ".docx",
    ".rtf",
    ".odt",
    ".csv",
})


# Extensions handled by the openpyxl-based colour-aware xlsx
# pipeline. Separated from the markitdown path because xlsx
# conversion preserves cell background colours as emoji markers
# — markitdown ignores formatting entirely, which loses
# information for spreadsheets used to track status (red =
# blocked, green = done, etc.).
_XLSX_EXTENSIONS: frozenset[str] = frozenset({
    ".xlsx",
})


# RGB values treated as "effectively no fill" during xlsx
# extraction. Near-white and near-black fills are almost always
# defaults (unformatted cells, borders) rather than meaningful
# status markers — emitting an emoji for every such cell would
# overwhelm the output. The threshold is a per-channel distance.
_IGNORE_NEAR_WHITE_THRESHOLD = 20  # per-channel delta from 255
_IGNORE_NEAR_BLACK_THRESHOLD = 20  # per-channel delta from 0

# RGB Euclidean distance below which two colours are treated as
# the same cluster for the fallback marker assignment. Tuned so
# three visibly-distinct shades of brown each get their own
# marker, but slight rendering variations of the same "red"
# collapse together.
_COLOUR_CLUSTER_DISTANCE = 40.0

# Well-known hue markers — named colours that users reach for
# first when adding cell fills. Order doesn't matter; lookup is
# by closest match within the named-hue set.
_NAMED_COLOURS: tuple[tuple[str, tuple[int, int, int], str], ...] = (
    ("red",    (255,   0,   0), "🔴"),
    ("green",  (  0, 200,   0), "🟢"),
    ("yellow", (255, 230,   0), "🟡"),
    ("blue",   (  0, 100, 255), "🔵"),
    ("orange", (255, 140,   0), "🟠"),
    ("purple", (150,   0, 200), "🟣"),
    ("pink",   (255, 130, 200), "🩷"),
    ("brown",  (139,  69,  19), "🟤"),
)

# Distance threshold for matching against named colours. An
# unknown colour closer than this to a named colour is assigned
# the named marker; otherwise it falls through to the fallback
# clustering. Larger than the cluster distance because named
# colours are allowed to absorb a wider range of shades (every
# "pinkish red" should get 🔴 rather than proliferating fallback
# markers).
_NAMED_COLOUR_DISTANCE = 80.0

# Fallback markers assigned in order to unique colour clusters
# that don't match any named hue. Distinct enough visually to
# differentiate three shades of brown without confusing the
# reader. More than eight clusters is rare in practice.
_FALLBACK_MARKERS: tuple[str, ...] = (
    "⬛", "◆", "▲", "●", "■", "★", "◉", "◈",
)


# Extensions handled by the python-pptx fallback pipeline. The
# primary path for presentations is LibreOffice + PyMuPDF
# (Pass A5) which produces text+SVG hybrid output; this
# fallback runs when those dependencies aren't installed and
# renders each slide as a full SVG via python-pptx.
_PPTX_EXTENSIONS: frozenset[str] = frozenset({
    ".pptx",
})


# Extensions handled directly by the PyMuPDF pipeline (no
# LibreOffice conversion needed — PyMuPDF reads PDFs natively).
# pptx and odp use PyMuPDF too but go through LibreOffice first
# to produce an intermediate PDF (handled by Pass A5b).
_PDF_EXTENSIONS: frozenset[str] = frozenset({
    ".pdf",
})


# Extensions routed through the LibreOffice → PDF → PyMuPDF
# pipeline when available. Falls back to format-specific paths
# when LibreOffice or PyMuPDF is missing:
#   .pptx → python-pptx fallback
#   .odp  → markitdown fallback
_LIBREOFFICE_EXTENSIONS: frozenset[str] = frozenset({
    ".pptx",
    ".odp",
})


# Timeout (seconds) for the `soffice --headless --convert-to
# pdf` subprocess. LibreOffice launches its UNO listener lazily
# and can take several seconds on first invocation; subsequent
# invocations are faster. 120 seconds is generous but bounded —
# prevents hung conversions from wedging the executor.
_LIBREOFFICE_TIMEOUT_SECONDS = 120


# Minimum number of "significant" drawings on a page before we
# trigger SVG export alongside text extraction. Below this
# threshold, the page is treated as text-only — no SVG produced.
# Tuned so pages with just borders or table rules (which every
# PDF generator emits for layout) don't bloat the output.
_PAGE_GRAPHICS_THRESHOLD = 3


# Minimum segment counts for a vector drawing to count as
# "significant". Simple rectangles and straight lines don't
# qualify; curves and multi-segment paths do.
_PATH_SIGNIFICANT_SEGMENTS = 4
_POLYGON_SIGNIFICANT_SEGMENTS = 2


# EMU (English Metric Units) — python-pptx's native unit for
# all dimensions. 914400 EMU per inch.
_EMU_PER_INCH = 914400

# Screen DPI for SVG viewBox dimensions. SVG's default user
# unit is 1px at 96 DPI; using 96 here produces SVGs that
# render at the same visual size as the original slide when
# displayed in a 1:1 viewer. User units in SVG scale cleanly,
# so this is just a reference point — nothing in the pipeline
# actually depends on the literal pixel values.
_SVG_DPI = 96

# Conversion factor EMU → pixels.
_EMU_TO_PX = _SVG_DPI / _EMU_PER_INCH

# Default slide dimensions in EMU when python-pptx reports None.
# Standard 4:3 slide at 10" x 7.5" — the pptx default. Rarely
# encountered (real files always have a slide size) but keeps
# the pipeline robust against corrupted or exotic templates.
_DEFAULT_SLIDE_WIDTH_EMU = 9144000   # 10 inches
_DEFAULT_SLIDE_HEIGHT_EMU = 6858000  # 7.5 inches

# Fallback font size in points when python-pptx reports None.
# Matches PowerPoint's default body text size.
_DEFAULT_FONT_SIZE_PT = 18

# Fallback font colour when python-pptx reports None.
# Black reads correctly against the default white slide
# background; slides with dark themes will need future work
# to resolve background-colour-aware defaults.
_DEFAULT_FONT_COLOR = "#000000"

# Pixels per font-size point for SVG font-size attribute.
# SVG font-size is in user units (pixels at 96 DPI); 1pt =
# 1/72 inch, so 1pt * 96/72 = 4/3 pixels.
_PT_TO_PX = 96 / 72

# Zero-padding width for slide filenames. Two digits covers
# presentations up to 99 slides; longer decks pad to three.
# Using a fixed width per deck keeps the file listing sorted
# correctly in every tool.
_SLIDE_NUMBER_MIN_WIDTH = 2


# Regex matching `![alt](data:mime;base64,payload)` — the shape
# markitdown emits for embedded images. Group 1 is the whole
# `data:...` URL; group 2 is the MIME subtype (e.g. `png`);
# group 3 is the base64 payload OR the literal `...` for the
# DOCX truncated case.
#
# The payload allows `=` (base64 padding) and `/`, `+`
# (standard alphabet). We use `[^)]*` rather than a strict
# base64 charset because some markitdown outputs include stray
# whitespace or newlines that a strict pattern would break on
# — decoding catches real errors, and we don't want regex
# strictness to drop legitimately-encoded payloads.
_DATA_URI_IMAGE_RE = re.compile(
    r"!\[([^\]]*)\]\((data:image/([^;]+);base64,([^)]+))\)",
    re.IGNORECASE,
)


# Regex matching the DOCX truncated-URI shape:
# `data:image/png;base64...` (literal ellipsis, no closing paren
# fence because markitdown sometimes emits these without one).
# Group 1 is the MIME subtype. Handles both the wrapped form
# `![alt](data:image/png;base64...)` and the bare reference.
_TRUNCATED_URI_RE = re.compile(
    r"data:image/([^;]+);base64\.{3}",
    re.IGNORECASE,
)


# Per-image MIME-to-extension map for image extraction. Covers
# the formats DOCX / ODT / RTF realistically embed. Unknown
# MIMEs fall through to `.bin` so the file still lands on disk
# and the provenance header still records it — the user can
# rename if needed, but we never silently drop.
_MIME_TO_EXT: dict[str, str] = {
    "png": ".png",
    "jpeg": ".jpg",
    "jpg": ".jpg",
    "gif": ".gif",
    "webp": ".webp",
    "bmp": ".bmp",
    "tiff": ".tif",
    "svg+xml": ".svg",
}


# ---------------------------------------------------------------------------
# Data shapes
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class ProvenanceHeader:
    """Parsed docuvert provenance comment.

    Fields match the spec-defined header format. All fields
    optional except `source` and `sha256` — those are the minimum
    to classify status.

    Frozen because the parser returns these to the scanner and the
    scanner should never mutate them; immutability is a small
    safety win.
    """

    source: str
    sha256: str
    images: tuple[str, ...] = ()
    extra: dict[str, str] | None = None


# ---------------------------------------------------------------------------
# DocConvert
# ---------------------------------------------------------------------------


class DocConvert:
    """Document-to-markdown conversion service.

    Construct with a :class:`ConfigManager` and optional :class:`Repo`.
    The repo argument lets the scanner report paths relative to the
    repository root — without it, paths are resolved against CWD,
    which matches the convention used by tests that don't want to
    construct a full repo.

    Register via ``server.add_class(doc_convert)`` alongside the other
    services. In collab mode, `main.py` sets `_collab` on the
    instance after construction; in single-user mode the attribute
    stays None and every caller is treated as localhost.
    """

    def __init__(
        self,
        config: "ConfigManager",
        repo: Any = None,
        event_callback: Any = None,
    ) -> None:
        """Construct the service.

        Parameters
        ----------
        config:
            The config manager — used for reading `doc_convert_config`
            and (in future passes) for getting the working directory
            when the PDF pipeline needs temp space.
        repo:
            Optional :class:`Repo` instance. When provided, scans
            resolve against the repo root; when None, scans resolve
            against CWD. Kept as `Any` to avoid the circular import
            between Layer 1 (Repo) and Layer 4 (DocConvert).
        event_callback:
            Optional async callable ``(event_name, *args) ->
            awaitable`` for pushing progress events to the browser.
            Wired by `main.py` to the shared event dispatcher so
            `docConvertProgress` events reach the webapp. When None
            (tests, standalone CLI use), events are silently dropped
            and `convert_files` falls back to fully synchronous
            operation returning inline results.
        """
        self._config = config
        self._repo = repo
        self._event_callback = event_callback
        # Collab reference, set by main.py when collab mode is active.
        # None in single-user mode, which means every caller is
        # treated as localhost. Matches the pattern on Repo,
        # LLMService, and Settings.
        self._collab: Any = None

    # ------------------------------------------------------------------
    # Localhost-only guard
    # ------------------------------------------------------------------

    def _check_localhost_only(self) -> dict[str, Any] | None:
        """Return a restricted-error dict when caller is non-localhost.

        Same contract as the other services' guards. Fails closed on
        collab-check exceptions — a raising collab denies the call
        rather than letting it through. The frontend's RpcMixin
        surfaces the restricted shape as a distinct error class.
        """
        if self._collab is None:
            return None
        try:
            is_local = self._collab.is_caller_localhost()
        except Exception as exc:
            logger.warning(
                "Collab localhost check raised: %s; denying",
                exc,
            )
            return {
                "error": "restricted",
                "reason": (
                    "Internal error checking caller identity"
                ),
            }
        if is_local:
            return None
        return {
            "error": "restricted",
            "reason": (
                "Participants cannot perform this action"
            ),
        }

    # ------------------------------------------------------------------
    # Configuration readers
    # ------------------------------------------------------------------

    @property
    def _enabled(self) -> bool:
        """Read-through to `config.doc_convert_config['enabled']`."""
        return bool(self._config.doc_convert_config.get("enabled", True))

    @property
    def _extensions(self) -> tuple[str, ...]:
        """Lowercased configured extensions, with defaults fallback.

        Normalises to a tuple so the result is hashable and
        comparable. Lowercases the entries because the extension
        matching in `_is_convertible` compares lowercased suffixes —
        a config entry of `.DOCX` should still work.
        """
        configured = self._config.doc_convert_config.get("extensions")
        if not configured or not isinstance(configured, (list, tuple)):
            return _DEFAULT_EXTENSIONS
        return tuple(str(ext).lower() for ext in configured)

    @property
    def _max_size_bytes(self) -> int:
        """Size threshold in bytes (config stores it in MB)."""
        mb = int(
            self._config.doc_convert_config.get("max_source_size_mb", 50)
        )
        return mb * 1024 * 1024

    # ------------------------------------------------------------------
    # Repo-root resolution
    # ------------------------------------------------------------------

    def _root(self) -> Path:
        """Return the root directory for scans.

        Uses `repo.root` when a Repo is attached; falls back to CWD
        when not. The fallback is for tests that don't construct a
        full repo, not for production — production always passes
        `repo`.
        """
        if self._repo is not None:
            # Repo exposes `.root` as a Path property.
            root = getattr(self._repo, "root", None)
            if root is not None:
                return Path(root)
        return Path.cwd()

    # ------------------------------------------------------------------
    # is_available — optional dependency probe
    # ------------------------------------------------------------------

    def is_available(self) -> dict[str, bool]:
        """Probe every optional dependency and report status.

        Returns
        -------
        dict
            - ``available`` — True when markitdown is importable.
              Without it, no conversion is possible at all (even the
              PDF path ultimately produces markdown output, though
              via PyMuPDF text extraction rather than markitdown).
              The frontend uses this to decide whether to render the
              Doc Convert tab at all.
            - ``libreoffice`` — True when `soffice` is on PATH.
              Enables the pptx/odp → PDF conversion step.
            - ``pymupdf`` — True when `fitz` is importable. Enables
              the PDF-to-markdown pipeline (text extraction + image
              detection + SVG export).
            - ``pdf_pipeline`` — True only when BOTH LibreOffice and
              PyMuPDF are available. The pptx/odp route needs both;
              pdf alone needs only PyMuPDF, but the frontend shows
              a single "PDF pipeline" capability to keep the UI
              simple.

        Safe to call frequently — each probe is a single import or
        PATH lookup. No subprocess launches, no network calls.
        Always callable regardless of the `enabled` config flag
        (availability is a capability query, not a feature toggle).
        """
        markitdown_ok = self._probe_import("markitdown")
        pymupdf_ok = self._probe_import("fitz")
        libreoffice_ok = shutil.which("soffice") is not None
        return {
            "available": markitdown_ok,
            "libreoffice": libreoffice_ok,
            "pymupdf": pymupdf_ok,
            "pdf_pipeline": libreoffice_ok and pymupdf_ok,
        }

    @staticmethod
    def _probe_import(module_name: str) -> bool:
        """Return True when `module_name` can be imported.

        Broad exception catch — a module that installs but fails on
        import (version mismatch, missing binary dependency, corrupted
        install) should show as unavailable rather than propagating
        the exception into what is meant to be a cheap probe.
        """
        import importlib
        try:
            importlib.import_module(module_name)
            return True
        except Exception as exc:
            logger.debug(
                "DocConvert probe: %s not available (%s)",
                module_name, exc,
            )
            return False

    # ------------------------------------------------------------------
    # scan_convertible_files — repository walk + status classification
    # ------------------------------------------------------------------

    def scan_convertible_files(self) -> list[dict[str, Any]]:
        """Walk the repo and classify every convertible source file.

        Each entry is a dict with:

        - ``path`` — source path relative to the root
        - ``name`` — basename
        - ``size`` — source file size in bytes
        - ``status`` — one of `"new"`, `"current"`, `"stale"`,
          `"conflict"`
        - ``output_path`` — prospective output path, relative to root
        - ``over_size`` — True when size exceeds `max_source_size_mb`
          (advisory — the file still appears in the scan so the UI
          can show a warning badge; conversion will refuse it)

        Returns an empty list when `enabled=false` in config. The
        result is stable-sorted by path for deterministic frontend
        rendering.

        Never raises — filesystem errors on individual files produce
        a debug log and that file is skipped. A pathological repo
        (weird symlinks, permission denied on a directory) should not
        prevent the rest of the scan from succeeding.
        """
        if not self._enabled:
            return []

        root = self._root()
        if not root.is_dir():
            logger.warning(
                "DocConvert scan root is not a directory: %s", root
            )
            return []

        extensions = self._extensions
        max_bytes = self._max_size_bytes

        entries: list[dict[str, Any]] = []
        for source_abs in self._iter_candidates(root, extensions):
            try:
                rel_path = source_abs.relative_to(root)
            except ValueError:
                # Shouldn't happen — the walker is rooted at `root` —
                # but be defensive.
                continue
            try:
                size = source_abs.stat().st_size
            except OSError as exc:
                logger.debug(
                    "DocConvert scan: stat failed for %s: %s",
                    source_abs, exc,
                )
                continue

            output_abs = source_abs.with_suffix(".md")
            try:
                output_rel = output_abs.relative_to(root)
            except ValueError:
                continue

            status = self._classify_status(source_abs, output_abs)

            entries.append({
                "path": str(rel_path).replace("\\", "/"),
                "name": source_abs.name,
                "size": size,
                "status": status,
                "output_path": str(output_rel).replace("\\", "/"),
                "over_size": size > max_bytes,
            })

        # Deterministic sort by source path so repeated scans produce
        # byte-identical output. Frontend relies on stable ordering
        # for UI state (selected checkboxes, scroll position).
        entries.sort(key=lambda e: e["path"])
        return entries

    def _iter_candidates(
        self,
        root: Path,
        extensions: tuple[str, ...],
    ) -> Any:
        """Walk `root`, yielding files with matching extensions.

        Skips every directory in `_EXCLUDED_DIRS` plus hidden
        directories (except `.github` which some repos use for CI
        config that might contain docs worth converting).

        Uses `os.walk` rather than `Path.rglob` because it's easier
        to mutate the directory list in place (the standard way to
        prune a walk). `Path.rglob` has no equivalent prune hook;
        filtering after the fact would still descend into
        `node_modules`.
        """
        extensions_set = set(extensions)
        for dirpath, dirnames, filenames in os.walk(root):
            # Prune excluded and hidden dirs in place. The walker
            # respects in-place mutation of `dirnames`.
            dirnames[:] = [
                d for d in dirnames
                if d not in _EXCLUDED_DIRS
                and (not d.startswith(".") or d == ".github")
            ]
            dir_path = Path(dirpath)
            for filename in filenames:
                suffix = Path(filename).suffix.lower()
                if suffix not in extensions_set:
                    continue
                yield dir_path / filename

    def _classify_status(
        self,
        source_abs: Path,
        output_abs: Path,
    ) -> str:
        """Classify the source file's conversion status.

        Order of checks matters:

        1. No output file exists → `new`
        2. Output exists but has no docuvert header → `conflict`
           (manually authored or externally converted)
        3. Header present, source hash matches → `current`
        4. Header present, source hash doesn't match → `stale`

        Defensive — any I/O failure reading the output file or
        hashing the source is treated as `new`, so a permissions
        hiccup doesn't silently show `current` for a file the user
        can't actually read.
        """
        if not output_abs.is_file():
            return "new"

        header = self._read_provenance_header(output_abs)
        if header is None:
            return "conflict"

        try:
            current_hash = self._hash_file(source_abs)
        except OSError as exc:
            logger.debug(
                "DocConvert scan: hash failed for %s: %s",
                source_abs, exc,
            )
            return "new"

        if header.sha256 == current_hash:
            return "current"
        return "stale"

    # ------------------------------------------------------------------
    # Provenance header parsing
    # ------------------------------------------------------------------

    def _read_provenance_header(
        self,
        output_abs: Path,
    ) -> ProvenanceHeader | None:
        """Read the provenance header from a converted output file.

        Returns None when the file has no header (manually authored),
        when the header is malformed, or when required fields
        (source, sha256) are missing. Lenient on everything else:
        unknown fields land in the `extra` dict for forward
        compatibility with future header additions.

        Reads only the first few KB of the file — the header lives
        on the first line. Reading the whole file would slow
        scanning on repos with large converted outputs.
        """
        try:
            with output_abs.open("rb") as fh:
                probe = fh.read(_PROVENANCE_PROBE_BYTES)
        except OSError as exc:
            logger.debug(
                "DocConvert: failed to read header from %s: %s",
                output_abs, exc,
            )
            return None

        text = probe.decode("utf-8", errors="replace")
        match = _PROVENANCE_RE.search(text)
        if match is None:
            return None

        return self.parse_provenance_body(match.group(1))

    @staticmethod
    def parse_provenance_body(body: str) -> ProvenanceHeader | None:
        """Parse the key=value pairs inside a docuvert header.

        Returns None when required fields are missing. Exposed as a
        static method so tests and future utilities can exercise the
        parser directly without needing a DocConvert instance.

        Recognised fields:

        - ``source`` — source filename (required)
        - ``sha256`` — source content hash (required)
        - ``images`` — comma-separated list of extracted image
          filenames (optional)

        Any other key=value pairs are captured in `extra` for
        forward compatibility. A future release adding a new
        field (e.g., ``tool_version``) won't cause older clients to
        fail — they just won't display the new field.
        """
        fields: dict[str, str] = {
            match.group(1).lower(): match.group(2)
            for match in _PROV_FIELD_RE.finditer(body)
        }

        source = fields.pop("source", None)
        sha256 = fields.pop("sha256", None)
        if not source or not sha256:
            return None

        images_raw = fields.pop("images", "")
        images = tuple(
            name.strip() for name in images_raw.split(",")
            if name.strip()
        )

        return ProvenanceHeader(
            source=source,
            sha256=sha256,
            images=images,
            extra=fields or None,
        )

    # ------------------------------------------------------------------
    # Source hashing
    # ------------------------------------------------------------------

    @staticmethod
    def _hash_file(path: Path) -> str:
        """SHA-256 hex digest of a file's content.

        Streams in 64 KB chunks so large files don't need to fit in
        memory. The hex output is what goes into the provenance
        header (shorter prefix would risk collisions across a repo
        with thousands of source files).
        """
        h = hashlib.sha256()
        with path.open("rb") as fh:
            while True:
                chunk = fh.read(65536)
                if not chunk:
                    break
                h.update(chunk)
        return h.hexdigest()

    # ------------------------------------------------------------------
    # convert_files — Pass A2 markitdown path
    # ------------------------------------------------------------------

    def convert_files(
        self,
        paths: list[str],
    ) -> dict[str, Any]:
        """Convert the named source files to markdown.

        Dispatches by extension. Pass A2 (`.docx`/`.rtf`/`.odt`/
        `.csv`) uses markitdown; A3 (`.xlsx`) uses openpyxl;
        A4 (`.pptx` fallback) uses python-pptx; A5 (`.pdf`,
        `.pptx`, `.odp`) uses PyMuPDF with LibreOffice for
        presentations.

        Runs the clean-tree gate first: refuses if the repo has
        uncommitted changes. Without a git working tree to diff
        against, the converted output wouldn't be reviewable
        before commit. The gate requires a repo with an
        `is_clean()` method; when constructed without a repo
        (tests, standalone CLI use), the gate is skipped —
        caller accepts the risk.

        Two execution modes:

        - **Background (async)** — when an asyncio event loop is
          running AND an event callback is wired, launches
          per-file conversion as a background task and returns
          ``{"status": "started", "count": N}`` immediately.
          Per-file progress and the final summary arrive via
          ``docConvertProgress`` server-push events.
        - **Inline (sync)** — when no event loop is running
          (tests, CLI use without websocket) or no event
          callback is wired, runs conversions inline and
          returns ``{"status": "ok", "results": [...]}``
          directly. Matches the pre-async contract so tests
          don't need event-loop setup.

        The guard and gate run identically in both modes, so
        restricted/dirty-tree errors surface the same way.

        Parameters
        ----------
        paths:
            Repo-relative source file paths.

        Returns
        -------
        dict
            Restricted-error on non-localhost callers, dirty-tree
            error on dirty working tree, ``{"status": "started",
            "count": N}`` when running in background mode, or
            ``{"status": "ok", "results": [...]}`` when running
            inline.
        """
        restricted = self._check_localhost_only()
        if restricted is not None:
            return restricted

        # Decide execution mode. Background requires both a
        # running loop AND a wired callback — if either is
        # missing, fall back to inline so tests and CLI callers
        # get results directly.
        import asyncio
        loop: asyncio.AbstractEventLoop | None = None
        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = None

        if loop is not None and self._event_callback is not None:
            asyncio.ensure_future(
                self._convert_files_background(paths)
            )
            return {
                "status": "started",
                "count": len(paths),
            }

        # Inline fallback — synchronous execution.
        return self._convert_files_sync(paths)

    def _convert_files_sync(
        self,
        paths: list[str],
    ) -> dict[str, Any]:
        """Synchronous per-file conversion. Returns results inline.

        Used by tests and by the CLI-without-websocket path. The
        background path uses the same per-file loop but emits
        progress events between files.
        """
        root = self._root()
        results: list[dict[str, Any]] = []
        for rel_path in paths:
            result = self._convert_one(root, rel_path)
            results.append(result)
        return {"status": "ok", "results": results}

    async def _convert_files_background(
        self,
        paths: list[str],
    ) -> None:
        """Background conversion task with progress events.

        Each file is converted in the default thread executor so
        the event loop stays responsive — individual file
        conversions (especially PDFs via PyMuPDF) routinely block
        for many seconds, and running them directly on the loop
        would stall the websocket and prevent the ``start`` and
        ``file`` events from reaching the browser until the
        whole batch finished. That produced a "Converting 0 of
        N…" hang with no progress events visible in the UI.

        Emits three event stages:

        - ``start`` — before the first file, carrying the total
          count so the UI can size its progress display
        - ``file`` — per-file, after conversion completes,
          carrying the result dict and the running index
        - ``complete`` — after all files, carrying the full
          results list so the UI has the final summary even
          if it missed some per-file events

        Event failures are swallowed so a broken frontend
        subscriber can't abort an in-flight conversion batch.
        """
        import asyncio

        root = self._root()
        total = len(paths)
        await self._send_convert_event({
            "stage": "start",
            "count": total,
        })

        loop = asyncio.get_running_loop()
        results: list[dict[str, Any]] = []
        for index, rel_path in enumerate(paths):
            try:
                # Run the blocking conversion off-loop. The
                # default executor is a ThreadPoolExecutor, so
                # libraries that release the GIL (PyMuPDF's C
                # core does) actually run in parallel with the
                # loop — not that we need parallelism here; we
                # just need the loop to stay free to deliver
                # the progress events queued above and below.
                result = await loop.run_in_executor(
                    None,
                    self._convert_one,
                    root,
                    rel_path,
                )
            except Exception as exc:
                # Defensive — a bug in a per-file conversion
                # shouldn't abort the whole batch. Surface as an
                # error result so the UI shows the specific file
                # that failed, not a silent drop.
                logger.exception(
                    "DocConvert: unhandled error converting %s",
                    rel_path,
                )
                result = self._fail(
                    rel_path,
                    f"Internal error: {exc}",
                )
            results.append(result)
            await self._send_convert_event({
                "stage": "file",
                "index": index,
                "total": total,
                "result": result,
            })

        await self._send_convert_event({
            "stage": "complete",
            "results": results,
        })

    async def _send_convert_event(
        self,
        data: dict[str, Any],
    ) -> None:
        """Dispatch a docConvertProgress event. Swallows failures.

        Wrapped so a throwing callback (unmounted frontend,
        websocket error) doesn't abort the conversion loop.
        The event is lost for that frontend client; the batch
        continues and the final `complete` event carries the
        full results list so a reconnecting client can recover
        the state.
        """
        if self._event_callback is None:
            return
        try:
            await self._event_callback(
                "docConvertProgress",
                data,
            )
        except Exception as exc:
            logger.debug(
                "DocConvertProgress event dispatch failed: %s",
                exc,
            )

    def _convert_one(
        self,
        root: Path,
        rel_path: str,
    ) -> dict[str, Any]:
        """Convert a single file. Returns a per-file result dict.

        Entry point for the per-file dispatch. Handles validation
        (path inside root, file exists, size within budget),
        extension routing, and wraps every failure in a result
        dict rather than propagating. One failing file shouldn't
        abort the whole batch — the user may have mixed supported
        and unsupported files in one selection.
        """
        # Resolve and validate the path is inside the scan root.
        # A traversal attempt (`../foo.docx`) would be caught by
        # relative_to below, but we normalise first so the error
        # message names the requested path, not the resolved one.
        try:
            source_abs = (root / rel_path).resolve()
            source_abs.relative_to(root.resolve())
        except (OSError, ValueError):
            return self._fail(
                rel_path,
                "Path must be within repository root",
            )

        if not source_abs.is_file():
            return self._fail(rel_path, "File not found")

        # Size check — mirrors the scan's over_size flag. A file
        # over the budget produces a skip result with an
        # explanatory message so the UI can render a warning
        # rather than silently dropping.
        try:
            size = source_abs.stat().st_size
        except OSError as exc:
            return self._fail(rel_path, f"Stat failed: {exc}")
        if size > self._max_size_bytes:
            return self._skip(
                rel_path,
                f"File exceeds {self._max_size_bytes // (1024 * 1024)}MB limit",
            )

        suffix = source_abs.suffix.lower()

        # markitdown handles the simple formats (.docx, .rtf,
        # .odt, .csv).
        if suffix in _MARKITDOWN_EXTENSIONS:
            return self._convert_via_markitdown(root, source_abs, rel_path)

        # xlsx uses a dedicated openpyxl-based pipeline so cell
        # background colours survive as emoji markers. Falls back
        # to markitdown if openpyxl is unavailable or fails.
        if suffix in _XLSX_EXTENSIONS:
            return self._convert_via_openpyxl(root, source_abs, rel_path)

        # pptx / odp route through LibreOffice + PyMuPDF when
        # available — gives proper text extraction plus per-page
        # SVGs for diagrams. Falls back to format-specific paths
        # when dependencies are missing.
        if suffix in _LIBREOFFICE_EXTENSIONS:
            return self._convert_via_libreoffice(
                root, source_abs, rel_path
            )

        # pptx uses the python-pptx fallback pipeline — renders
        # each slide as an SVG with embedded text, images, and
        # tables. Only reached when LibreOffice isn't available
        # (the libreoffice path would have handled .pptx above).
        if suffix in _PPTX_EXTENSIONS:
            return self._convert_via_python_pptx(
                root, source_abs, rel_path
            )

        # pdf goes directly to PyMuPDF — hybrid text + SVG
        # output. Text extracted into markdown paragraphs;
        # pages with significant graphics also get companion
        # SVGs with glyphs stripped (text already in markdown).
        if suffix in _PDF_EXTENSIONS:
            return self._convert_via_pymupdf(
                root, source_abs, rel_path
            )

        # Other supported extensions — explicit "not yet
        # supported" rather than a silent skip. Users see exactly
        # which files the current release can convert.
        if suffix in self._extensions:
            return self._skip(
                rel_path,
                (
                    f"{suffix} conversion is not yet implemented; "
                    "ships in a later pass"
                ),
            )

        # Extension isn't recognised at all. Should be unreachable
        # if the caller used `scan_convertible_files` as its source
        # of truth, but handle defensively.
        return self._fail(
            rel_path,
            f"Unsupported extension: {suffix}",
        )

    def _convert_via_markitdown(
        self,
        root: Path,
        source_abs: Path,
        rel_path: str,
    ) -> dict[str, Any]:
        """Convert via markitdown — the `.docx`/`.rtf`/`.odt` path.

        Orchestrates the full per-file pipeline:

        1. Compute source hash (for provenance header)
        2. Read prior provenance (for orphan-image cleanup on stale)
        3. Call markitdown to produce markdown text
        4. For DOCX: pre-extract images from zip `word/media/`,
           substitute truncated URIs
        5. Extract data-URI images, save to assets subdir,
           rewrite markdown references
        6. Clean up orphan images from prior conversion
        7. Remove empty assets subdir if no images were saved
        8. Write markdown with provenance header
        """
        # Lazy import — markitdown is an optional dependency.
        # Surfacing the error here rather than at module load
        # means the rest of DocConvert (scan, is_available) works
        # in stripped-down releases.
        try:
            from markitdown import MarkItDown
        except ImportError:
            return self._fail(
                rel_path,
                (
                    "markitdown is not installed. Install with: "
                    "pip install 'ac-dc[docs]'"
                ),
            )

        output_abs = source_abs.with_suffix(".md")
        try:
            output_rel = output_abs.relative_to(root)
        except ValueError:
            return self._fail(
                rel_path,
                "Output path escapes repository root",
            )

        # Hash the source — written into the provenance header and
        # used by future scans to classify this output as `current`.
        try:
            source_hash = self._hash_file(source_abs)
        except OSError as exc:
            return self._fail(rel_path, f"Source hash failed: {exc}")

        # Read prior provenance — we'll diff its image list against
        # the new one to identify orphans. If there's no prior
        # header (new or conflict), there are no orphans to clean.
        prior_images: tuple[str, ...] = ()
        if output_abs.is_file():
            prior_header = self._read_provenance_header(output_abs)
            if prior_header is not None:
                prior_images = prior_header.images

        # Run markitdown. Broad exception catch — library errors
        # vary wildly (CorruptedFileError, InvalidArgumentError,
        # NotImplementedError for unsupported variants). Wrap
        # into a per-file error result rather than propagating.
        try:
            md = MarkItDown()
            result = md.convert(str(source_abs))
            markdown_text = result.text_content or ""
        except Exception as exc:
            return self._fail(
                rel_path,
                f"markitdown conversion failed: {exc}",
            )

        # Assets subdirectory — `docs/architecture.docx` produces
        # `docs/architecture/` for image storage. Created on
        # demand; removed at the end if no images were saved.
        assets_dir = source_abs.with_suffix("")
        stem = source_abs.stem

        # DOCX: unconditionally extract images from the zip's
        # ``word/media/`` directory BEFORE running the data-URI
        # pipeline. markitdown's behaviour with DOCX images is
        # unreliable — for some files it emits truncated
        # ``data:image/...base64...`` placeholders (handled by
        # `_replace_docx_truncated_uris` below), for others it
        # drops the reference entirely. The zip is the
        # authoritative source of "what images does this .docx
        # contain?", matching how the old AC-DC system worked.
        #
        # The data-URI pipeline runs afterwards to handle any
        # real inline data URIs markitdown did successfully
        # emit (small images sometimes survive intact). That
        # pass uses a disjoint filename range — its counter
        # starts at ``len(zip_extracted) + 1`` via the
        # ``start_index`` parameter — so filenames stay unique
        # across the two sources.
        zip_extracted: list[str] = []
        if source_abs.suffix.lower() == ".docx":
            zip_extracted = self._save_docx_zip_images(
                source_abs, assets_dir, stem
            )
            markdown_text = self._replace_docx_truncated_uris(
                markdown_text,
                zip_extracted,
                assets_dir.name,
            )

        # Extract data-URI images and rewrite the markdown.
        # Start numbering after the zip-extracted images so
        # filenames don't collide when both paths produce
        # output for the same source.
        markdown_text, data_uri_saved = self._extract_data_uri_images(
            markdown_text, assets_dir, stem,
            start_index=len(zip_extracted) + 1,
        )
        # Merge both sources in deterministic order: zip
        # images first (document-order within the archive),
        # then any data-URI extras. This list feeds the
        # provenance header and the orphan-cleanup diff.
        saved_images = tuple(zip_extracted) + data_uri_saved

        # Orphan cleanup — images listed in the prior provenance
        # header but NOT produced by this conversion are deleted.
        # Prevents the assets subdir accumulating stale files
        # across re-conversions of a changing source.
        new_image_set = set(saved_images)
        for orphan in prior_images:
            if orphan in new_image_set:
                continue
            orphan_path = assets_dir / orphan
            try:
                orphan_path.unlink()
            except OSError as exc:
                # Non-fatal — log and continue. A leftover orphan
                # is cosmetic; a failed conversion from a missing
                # file isn't.
                logger.debug(
                    "Failed to remove orphan image %s: %s",
                    orphan_path, exc,
                )

        # Remove the assets subdir if empty (no images extracted
        # AND the dir is empty — which could be the case if orphan
        # cleanup just emptied it).
        if assets_dir.is_dir():
            try:
                # `iterdir` is cheaper than listing; we just need
                # to know if anything remains.
                if not any(assets_dir.iterdir()):
                    assets_dir.rmdir()
            except OSError as exc:
                logger.debug(
                    "Failed to remove empty assets dir %s: %s",
                    assets_dir, exc,
                )

        # Build the final output: provenance header + markdown.
        provenance_line = self._build_provenance_header(
            source_name=source_abs.name,
            source_hash=source_hash,
            images=saved_images,
        )
        final_content = (
            provenance_line + "\n\n" + markdown_text.lstrip("\n")
        )

        # Write output. Parent dir already exists (source is
        # there), but be defensive for symlink edge cases.
        try:
            output_abs.parent.mkdir(parents=True, exist_ok=True)
            output_abs.write_text(final_content, encoding="utf-8")
        except OSError as exc:
            return self._fail(
                rel_path,
                f"Failed to write output: {exc}",
            )

        return {
            "path": rel_path,
            "status": "ok",
            "output_path": str(output_rel).replace("\\", "/"),
            "images": list(saved_images),
        }

    # ------------------------------------------------------------------
    # DOCX image pre-extraction (zip-based, unconditional)
    # ------------------------------------------------------------------

    def _save_docx_zip_images(
        self,
        source_abs: Path,
        assets_dir: Path,
        stem: str,
    ) -> list[str]:
        """Unconditionally extract images from ``word/media/``.

        DOCX files are zip archives; embedded images live
        under ``word/media/``. markitdown does not reliably
        surface these in the markdown output — for some
        files it emits truncated ``data:image/...base64...``
        placeholders (the ellipsis case handled below), for
        others it drops the reference entirely, and for small
        inline images it may inline a real data URI.

        The old AC-DC system extracted every media file
        unconditionally and then reconciled with the markdown
        separately. That's what we do here too: the source of
        truth for "what images does this .docx contain?" is
        the zip, not markitdown's output.

        Returns the list of saved filenames in document order
        (as emitted by ``zipfile.namelist()``), with names of
        the form ``{stem}_img{N}{ext}`` matching the
        data-URI pipeline's convention so the provenance
        header lists both sources with the same naming.

        Empty list on any failure (not a zip, no media dir,
        unreadable entries) — caller treats as "no images"
        and proceeds.
        """
        saved: list[str] = []
        try:
            with zipfile.ZipFile(source_abs, "r") as zf:
                # Sort media entries so the output is stable
                # across extractions of the same file —
                # zipfile.namelist preserves storage order
                # which for well-formed docx matches document
                # order, but a belt-and-braces sort costs
                # nothing and protects against oddly-assembled
                # archives.
                media_names = sorted(
                    name for name in zf.namelist()
                    if name.startswith("word/media/")
                )
                if not media_names:
                    return []
                for name in media_names:
                    ext = Path(name).suffix.lower()
                    if not ext:
                        continue
                    # Normalise `.jpeg` → `.jpg` to match the
                    # old system's convention; leaves other
                    # extensions untouched.
                    if ext == ".jpeg":
                        ext = ".jpg"
                    try:
                        raw = zf.read(name)
                    except (KeyError, RuntimeError) as exc:
                        logger.debug(
                            "DOCX media read failed for %s: %s",
                            name, exc,
                        )
                        continue
                    # Create assets dir lazily — only if we
                    # have at least one byte to write, so
                    # image-free docs still skip the dir.
                    try:
                        assets_dir.mkdir(
                            parents=True, exist_ok=True,
                        )
                    except OSError as exc:
                        logger.debug(
                            "Assets dir create failed %s: %s",
                            assets_dir, exc,
                        )
                        return saved
                    image_index = len(saved) + 1
                    image_name = f"{stem}_img{image_index}{ext}"
                    image_path = assets_dir / image_name
                    try:
                        image_path.write_bytes(raw)
                    except OSError as exc:
                        logger.debug(
                            "Image write failed %s: %s",
                            image_path, exc,
                        )
                        continue
                    saved.append(image_name)
        except zipfile.BadZipFile:
            logger.debug(
                "Not a valid zip: %s (docx extraction skipped)",
                source_abs,
            )
        except OSError as exc:
            logger.debug(
                "DOCX media extraction failed for %s: %s",
                source_abs, exc,
            )
        return saved

    def _replace_docx_truncated_uris(
        self,
        markdown_text: str,
        extracted_names: list[str],
        assets_dir_name: str,
    ) -> str:
        """Substitute truncated ``data:image`` placeholders.

        markitdown emits ``data:image/png;base64...`` (literal
        ellipsis, no payload) for large embedded images.
        Now that we've already extracted the real images from
        the zip in :meth:`_save_docx_zip_images`, simply
        rewrite each truncated reference to point at the
        corresponding extracted file.

        Order-of-appearance matching is the best we can do
        without docx relationship parsing — markitdown
        doesn't surface rIds. In practice the zip's document
        order and markitdown's emission order align for
        well-formed files.

        Returns the markdown unchanged when no truncated
        URIs are present.
        """
        if not extracted_names:
            return markdown_text
        truncated_matches = list(
            _TRUNCATED_URI_RE.finditer(markdown_text)
        )
        if not truncated_matches:
            return markdown_text
        count = min(len(truncated_matches), len(extracted_names))
        for i in range(count):
            rel_ref = f"{assets_dir_name}/{extracted_names[i]}"
            # Substitute one at a time so each truncated
            # reference gets its own filename, not all
            # replaced with the same one.
            markdown_text = _TRUNCATED_URI_RE.sub(
                rel_ref, markdown_text, count=1,
            )
        return markdown_text

    # ------------------------------------------------------------------
    # Data-URI image extraction
    # ------------------------------------------------------------------

    def _extract_data_uri_images(
        self,
        markdown_text: str,
        assets_dir: Path,
        stem: str,
        start_index: int = 1,
    ) -> tuple[str, tuple[str, ...]]:
        """Decode data-URI images, save to disk, rewrite references.

        Parameters
        ----------
        markdown_text:
            Output from markitdown, potentially containing
            `![alt](data:image/...;base64,...)` references.
        assets_dir:
            Per-source assets directory. Created on demand.
            Names follow `{stem}_img{N}{ext}` so repeated
            conversions produce stable filenames — matters for
            orphan detection and git diffs.
        stem:
            Source file stem for naming extracted images
            (e.g. `architecture` → `architecture_img1.png`).
        start_index:
            First image counter value. DOCX callers pass
            ``len(zip_extracted) + 1`` so data-URI images
            numbered after the zip-extracted images and
            filenames stay unique across the two sources.
            Defaults to 1 for non-DOCX callers.

        Returns
        -------
        tuple
            `(rewritten_markdown, image_filenames)`. The
            filename list is ordered by appearance; caller uses
            it for the provenance header and orphan cleanup.
            Empty tuple when no images were extracted.

        Images that fail decoding are left as-is in the
        markdown — the broken reference is preferable to
        silently dropping the image.
        """
        matches = list(_DATA_URI_IMAGE_RE.finditer(markdown_text))
        if not matches:
            return markdown_text, ()

        saved: list[str] = []
        # We iterate matches and build the new text by splicing
        # — simpler than re.sub with a callback because we need
        # to increment the counter AND create the assets dir
        # lazily.
        output_parts: list[str] = []
        cursor = 0
        for match in matches:
            alt_text = match.group(1)
            mime_sub = match.group(3).lower()
            payload = match.group(4)

            # Append text before the match verbatim.
            output_parts.append(markdown_text[cursor:match.start()])

            # Decode the payload. Failure (invalid base64,
            # truncated `...` slipping through) leaves the
            # original data URI in place so the markdown renders
            # a broken image rather than nothing.
            try:
                # Strip whitespace — markitdown sometimes wraps
                # long payloads. Decoding is tolerant of `=`
                # padding errors when we strip them first; we
                # don't bother because valid base64 decoders
                # handle missing padding gracefully in `b64decode`
                # with `validate=False`.
                image_bytes = base64.b64decode(payload, validate=False)
            except (ValueError, Exception) as exc:
                logger.debug(
                    "Data URI decode failed: %s", exc,
                )
                output_parts.append(match.group(0))
                cursor = match.end()
                continue

            # Decoded successfully. Create assets dir on first
            # successful extraction only — avoids empty dirs.
            try:
                assets_dir.mkdir(parents=True, exist_ok=True)
            except OSError as exc:
                logger.debug(
                    "Assets dir create failed %s: %s",
                    assets_dir, exc,
                )
                output_parts.append(match.group(0))
                cursor = match.end()
                continue

            # Name the image. `{stem}_img{N}{ext}` — `N` is
            # 1-indexed for user-friendliness. DOCX callers
            # offset past zip-extracted images via
            # ``start_index``; other callers get the default
            # of 1.
            ext = _MIME_TO_EXT.get(mime_sub, ".bin")
            image_index = len(saved) + start_index
            image_name = f"{stem}_img{image_index}{ext}"
            image_path = assets_dir / image_name

            try:
                image_path.write_bytes(image_bytes)
            except OSError as exc:
                logger.debug(
                    "Image write failed %s: %s",
                    image_path, exc,
                )
                output_parts.append(match.group(0))
                cursor = match.end()
                continue

            saved.append(image_name)
            # Rewrite the reference to point at the saved file.
            # Path is relative from the markdown file's location:
            # `{stem}/{image_name}` (the assets dir is a sibling
            # of the `.md` named after the source stem).
            rel_ref = f"{assets_dir.name}/{image_name}"
            output_parts.append(f"![{alt_text}]({rel_ref})")
            cursor = match.end()

        # Trailing text after the last match.
        output_parts.append(markdown_text[cursor:])
        return "".join(output_parts), tuple(saved)

    # ------------------------------------------------------------------
    # xlsx — colour-aware extraction via openpyxl
    # ------------------------------------------------------------------

    def _convert_via_openpyxl(
        self,
        root: Path,
        source_abs: Path,
        rel_path: str,
    ) -> dict[str, Any]:
        """Convert an .xlsx file, preserving cell background colours.

        Two-pass approach:

        1. Pass 1 — walk every cell in every sheet, collecting
           text values (normalised) and raw hex fill colours. The
           set of unique non-ignorable fills is built during this
           pass.
        2. Colour mapping — well-known hues (red, green, yellow,
           blue, purple, etc.) get named emoji markers. Remaining
           colours are clustered by Euclidean RGB distance and
           assigned fallback markers per cluster.
        3. Pass 2 — emit markdown tables sheet by sheet, cells
           prefixed with their colour marker.

        Empty columns and fully-empty rows are stripped. A
        legend mapping markers to colour names appears at the end.

        Falls back to markitdown on any openpyxl failure
        (ImportError, corrupt file, unexpected structure) — the
        user still gets SOMETHING rather than an error result.
        """
        # Lazy import — openpyxl is optional in stripped-down
        # releases. ImportError is expected in that case and
        # means "use markitdown instead", not "fail the
        # conversion".
        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.debug(
                "openpyxl not installed; falling back to "
                "markitdown for %s",
                rel_path,
            )
            return self._convert_via_markitdown(
                root, source_abs, rel_path
            )

        output_abs = source_abs.with_suffix(".md")
        try:
            output_rel = output_abs.relative_to(root)
        except ValueError:
            return self._fail(
                rel_path,
                "Output path escapes repository root",
            )

        # Hash source for provenance — matches the markitdown
        # path so status classification works uniformly.
        try:
            source_hash = self._hash_file(source_abs)
        except OSError as exc:
            return self._fail(
                rel_path, f"Source hash failed: {exc}"
            )

        # Open workbook in read-only mode for performance on
        # large files. Also pass data_only=False so we see
        # formulas as formulas (not cached values) — mostly
        # defensive; formula cells rarely have fills anyway.
        try:
            workbook = load_workbook(
                filename=str(source_abs),
                read_only=False,  # need .fill on cells
                data_only=True,
            )
        except Exception as exc:
            # Corrupt xlsx, password-protected, etc. — fall back
            # to markitdown rather than erroring. markitdown may
            # still extract something useful.
            logger.debug(
                "openpyxl failed to open %s: %s; "
                "falling back to markitdown",
                rel_path, exc,
            )
            return self._convert_via_markitdown(
                root, source_abs, rel_path
            )

        # Pass 1: collect cells and unique fills across all sheets.
        try:
            sheets_data, unique_fills = self._xlsx_pass1_collect(
                workbook
            )
        except Exception as exc:
            logger.debug(
                "openpyxl pass-1 failed for %s: %s; "
                "falling back to markitdown",
                rel_path, exc,
            )
            workbook.close()
            return self._convert_via_markitdown(
                root, source_abs, rel_path
            )
        finally:
            workbook.close()

        # Build colour → marker map.
        colour_map = self._xlsx_build_colour_map(unique_fills)

        # Pass 2: emit markdown per sheet.
        body_parts: list[str] = []
        for sheet_name, rows in sheets_data:
            sheet_md = self._xlsx_render_sheet(
                sheet_name, rows, colour_map
            )
            if sheet_md:
                body_parts.append(sheet_md)

        if not body_parts:
            # Workbook had no data in any sheet. Still produce
            # an output file so the scan classifies it as
            # `current`, but with an informative placeholder.
            body_parts.append("(empty spreadsheet)")

        # Append legend mapping each used marker to its colour name.
        legend = self._xlsx_render_legend(colour_map)
        if legend:
            body_parts.append(legend)

        markdown_text = "\n\n".join(body_parts) + "\n"

        # Write output with provenance header. xlsx path never
        # produces embedded images, so we skip the data-URI
        # extraction pipeline entirely.
        provenance_line = self._build_provenance_header(
            source_name=source_abs.name,
            source_hash=source_hash,
            images=(),
        )
        final_content = (
            provenance_line + "\n\n" + markdown_text.lstrip("\n")
        )

        try:
            output_abs.parent.mkdir(parents=True, exist_ok=True)
            output_abs.write_text(final_content, encoding="utf-8")
        except OSError as exc:
            return self._fail(
                rel_path,
                f"Failed to write output: {exc}",
            )

        return {
            "path": rel_path,
            "status": "ok",
            "output_path": str(output_rel).replace("\\", "/"),
            "images": [],
        }

    def _xlsx_pass1_collect(
        self,
        workbook: Any,
    ) -> tuple[
        list[tuple[str, list[list[tuple[str, str | None]]]]],
        set[str],
    ]:
        """First pass over an xlsx workbook.

        Returns a list of `(sheet_name, rows)` pairs and the set
        of unique non-ignorable hex fill colours across the
        workbook. Each row is a list of `(value, hex_fill)`
        tuples where `hex_fill` is None for ignorable fills.

        Cell values are normalised — None becomes empty string,
        "nan"/"none" (case-insensitive) become empty string, and
        all non-string values are stringified. Whitespace is
        preserved except that leading/trailing is trimmed.

        Exceptions from openpyxl propagate; the caller wraps
        with a fallback-to-markitdown path.
        """
        sheets_data: list[
            tuple[str, list[list[tuple[str, str | None]]]]
        ] = []
        unique_fills: set[str] = set()

        for sheet in workbook.worksheets:
            rows: list[list[tuple[str, str | None]]] = []
            # iter_rows with values_only=False gives us Cell
            # objects (needed for .fill). We walk the used range
            # only — iter_rows defaults to the sheet's dimension.
            for row in sheet.iter_rows():
                row_cells: list[tuple[str, str | None]] = []
                for cell in row:
                    value = self._normalise_cell_value(cell.value)
                    fill_hex = self._extract_cell_fill(cell)
                    if fill_hex is not None:
                        unique_fills.add(fill_hex)
                    row_cells.append((value, fill_hex))
                rows.append(row_cells)
            sheets_data.append((sheet.title, rows))

        return sheets_data, unique_fills

    @staticmethod
    def _normalise_cell_value(value: Any) -> str:
        """Normalise a raw cell value for markdown emission.

        - None → empty string
        - Pandas/numpy artifacts "nan" / "none" (case-insensitive)
          → empty string. These crop up when a spreadsheet was
          generated from a DataFrame with missing values.
        - Everything else is stringified and stripped of
          leading/trailing whitespace.

        Pipe characters are escaped as `\\|` since they would
        otherwise break the markdown table row.
        """
        if value is None:
            return ""
        text = str(value).strip()
        if text.lower() in ("nan", "none"):
            return ""
        # Escape pipes so they don't break table rows. Literal
        # backslashes in cell values are rare; we don't escape
        # those.
        return text.replace("|", r"\|")

    @staticmethod
    def _extract_cell_fill(cell: Any) -> str | None:
        """Return the cell's fill as a hex string, or None.

        Ignorable fills (near-white, near-black, no fill at all)
        return None so they don't produce emoji markers. The
        hex string is lowercase without a leading hash.

        openpyxl's fill model is verbose — cells with no explicit
        fill still have a PatternFill with fgColor set to a
        default theme colour. We filter those by checking the
        patternType and the raw RGB.
        """
        try:
            fill = cell.fill
            if fill is None:
                return None
            # Only solid fills carry meaningful colour info.
            # patternType "none" means no explicit fill.
            pattern_type = getattr(fill, "patternType", None)
            if pattern_type not in ("solid", "lightGrid", "darkGrid"):
                return None
            fg = fill.fgColor
            if fg is None:
                return None
            # fgColor.rgb is an 8-char hex string (AARRGGBB)
            # when set. Theme colours return None here; we can't
            # resolve them without the workbook's theme table,
            # which isn't worth the complexity for a diagnostic
            # marker.
            raw = getattr(fg, "rgb", None)
            if not raw or not isinstance(raw, str):
                return None
            # Some versions of openpyxl return the rgb as a
            # Value wrapper — coerce defensively.
            raw = str(raw).strip().lower()
            # Strip alpha channel if present.
            if len(raw) == 8:
                raw = raw[2:]
            if len(raw) != 6:
                return None
            # Filter ignorable colours.
            try:
                r = int(raw[0:2], 16)
                g = int(raw[2:4], 16)
                b = int(raw[4:6], 16)
            except ValueError:
                return None
            # Near-white → ignore.
            if (
                (255 - r) < _IGNORE_NEAR_WHITE_THRESHOLD
                and (255 - g) < _IGNORE_NEAR_WHITE_THRESHOLD
                and (255 - b) < _IGNORE_NEAR_WHITE_THRESHOLD
            ):
                return None
            # Near-black → ignore.
            if (
                r < _IGNORE_NEAR_BLACK_THRESHOLD
                and g < _IGNORE_NEAR_BLACK_THRESHOLD
                and b < _IGNORE_NEAR_BLACK_THRESHOLD
            ):
                return None
            return raw
        except Exception:
            # Anything unexpected — treat as no fill. Defensive
            # against openpyxl API drift.
            return None

    @staticmethod
    def _hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
        """Convert a 6-char hex string to an (r, g, b) tuple.

        Input is assumed valid (produced by `_extract_cell_fill`
        which already validates). No error handling on the
        integer parses — if they fail, the caller has a bug.
        """
        return (
            int(hex_str[0:2], 16),
            int(hex_str[2:4], 16),
            int(hex_str[4:6], 16),
        )

    @staticmethod
    def _colour_distance(
        a: tuple[int, int, int],
        b: tuple[int, int, int],
    ) -> float:
        """Euclidean RGB distance.

        Perceptually naive (doesn't weight green higher like
        proper colour-diff metrics) but fine for the "are these
        two reds the same colour?" question the clustering needs.
        """
        dr = a[0] - b[0]
        dg = a[1] - b[1]
        db = a[2] - b[2]
        return (dr * dr + dg * dg + db * db) ** 0.5

    def _xlsx_build_colour_map(
        self,
        unique_fills: set[str],
    ) -> dict[str, tuple[str, str]]:
        """Assign an emoji marker and colour name to each fill.

        Returns a dict mapping hex colour → `(marker, name)`
        tuple. The name is either a named-colour label (red,
        green, etc.) or a synthesised label for fallback
        clusters ("cluster-1", "cluster-2", …).

        Algorithm:

        1. For each unique fill, find the closest named colour
           within `_NAMED_COLOUR_DISTANCE`. If found, assign the
           named marker.
        2. For remaining fills (no named match), cluster by
           proximity — fills within `_COLOUR_CLUSTER_DISTANCE` of
           an existing cluster join it; otherwise start a new
           cluster. Assign fallback markers in order.

        Named colours can be shared by multiple fills (all
        "reddish" cells get 🔴 regardless of exact shade).
        Fallback clusters each get their own marker.
        """
        result: dict[str, tuple[str, str]] = {}

        # Sort fills for deterministic assignment — two runs on
        # the same workbook produce the same markers. Without
        # this, set iteration order would vary.
        sorted_fills = sorted(unique_fills)

        unmatched: list[tuple[str, tuple[int, int, int]]] = []

        # Step 1 — assign named colours where close enough.
        for hex_fill in sorted_fills:
            rgb = self._hex_to_rgb(hex_fill)
            best_name: str | None = None
            best_marker: str | None = None
            best_dist = _NAMED_COLOUR_DISTANCE
            for name, named_rgb, marker in _NAMED_COLOURS:
                dist = self._colour_distance(rgb, named_rgb)
                if dist < best_dist:
                    best_dist = dist
                    best_name = name
                    best_marker = marker
            if best_name is not None and best_marker is not None:
                result[hex_fill] = (best_marker, best_name)
            else:
                unmatched.append((hex_fill, rgb))

        # Step 2 — cluster remaining fills. Each cluster holds a
        # representative RGB and the set of hex strings assigned
        # to it.
        clusters: list[tuple[tuple[int, int, int], list[str]]] = []
        for hex_fill, rgb in unmatched:
            joined = False
            for i, (cluster_rgb, cluster_hexes) in enumerate(clusters):
                if self._colour_distance(rgb, cluster_rgb) < _COLOUR_CLUSTER_DISTANCE:
                    cluster_hexes.append(hex_fill)
                    joined = True
                    break
            if not joined:
                clusters.append((rgb, [hex_fill]))

        # Assign fallback markers. More clusters than we have
        # markers — cycle through and append an index to the
        # name so entries remain distinguishable in the legend.
        for i, (_rgb, cluster_hexes) in enumerate(clusters):
            marker = _FALLBACK_MARKERS[i % len(_FALLBACK_MARKERS)]
            name = f"cluster-{i + 1}"
            for hex_fill in cluster_hexes:
                result[hex_fill] = (marker, name)

        return result

    def _xlsx_render_sheet(
        self,
        sheet_name: str,
        rows: list[list[tuple[str, str | None]]],
        colour_map: dict[str, tuple[str, str]],
    ) -> str:
        """Render one sheet as a markdown section.

        - Empty columns (all values empty across every row) are
          dropped.
        - Fully-empty rows are dropped.
        - First non-empty row becomes the table header. If every
          cell in that row is a string with meaningful content,
          it's used as-is; otherwise synthetic headers
          (`col1`, `col2`, …) are generated.
        - Coloured cells get their marker prepended with a
          space separator.

        Returns an empty string when the sheet has no data after
        stripping — the caller skips those sheets.
        """
        # Drop fully-empty rows.
        non_empty_rows = [
            row for row in rows
            if any(value for value, _ in row)
        ]
        if not non_empty_rows:
            return ""

        # Find the widest row — column count is the max width
        # across all non-empty rows. Shorter rows are padded
        # with empty cells during render.
        max_width = max(len(row) for row in non_empty_rows)
        if max_width == 0:
            return ""

        # Drop fully-empty columns. A column is empty if every
        # non-empty row has an empty value at that position.
        keep_columns: list[int] = []
        for col_idx in range(max_width):
            for row in non_empty_rows:
                if col_idx < len(row) and row[col_idx][0]:
                    keep_columns.append(col_idx)
                    break

        if not keep_columns:
            return ""

        # Header row — use the first non-empty row's values if
        # they look like headers (all non-empty strings, no
        # colour markers). Otherwise synthesise column names.
        first_row = non_empty_rows[0]
        use_first_as_header = all(
            col_idx < len(first_row) and first_row[col_idx][0]
            for col_idx in keep_columns
        )

        header_cells: list[str]
        data_rows: list[list[tuple[str, str | None]]]
        if use_first_as_header:
            header_cells = [
                first_row[col_idx][0] for col_idx in keep_columns
            ]
            data_rows = non_empty_rows[1:]
        else:
            header_cells = [
                f"col{i + 1}" for i in range(len(keep_columns))
            ]
            data_rows = non_empty_rows

        # Build markdown.
        lines: list[str] = [f"## {sheet_name}", ""]
        lines.append(
            "| " + " | ".join(header_cells) + " |"
        )
        lines.append(
            "|" + "|".join("---" for _ in keep_columns) + "|"
        )
        for row in data_rows:
            rendered_cells: list[str] = []
            for col_idx in keep_columns:
                if col_idx < len(row):
                    value, fill_hex = row[col_idx]
                else:
                    value, fill_hex = "", None
                if fill_hex is not None and fill_hex in colour_map:
                    marker = colour_map[fill_hex][0]
                    rendered_cells.append(
                        f"{marker} {value}" if value else marker
                    )
                else:
                    rendered_cells.append(value)
            lines.append(
                "| " + " | ".join(rendered_cells) + " |"
            )

        return "\n".join(lines)

    @staticmethod
    def _xlsx_render_legend(
        colour_map: dict[str, tuple[str, str]],
    ) -> str:
        """Render the colour-marker legend at the end of the output.

        Lists each unique (marker, name) pair exactly once. The
        colour map may contain multiple hex values mapped to the
        same named colour (all reddish fills → 🔴 red); the
        legend shows the named entry once rather than repeating
        per-hex.

        Returns an empty string when no markers were used.
        """
        if not colour_map:
            return ""
        # Collect unique (marker, name) pairs.
        seen: set[tuple[str, str]] = set()
        ordered_entries: list[tuple[str, str]] = []
        for marker, name in colour_map.values():
            key = (marker, name)
            if key in seen:
                continue
            seen.add(key)
            ordered_entries.append(key)
        if not ordered_entries:
            return ""
        lines = ["## Legend", ""]
        for marker, name in ordered_entries:
            lines.append(f"- {marker} {name}")
        return "\n".join(lines)

    # ------------------------------------------------------------------
    # pptx / odp — LibreOffice → PDF → PyMuPDF pipeline (primary)
    # ------------------------------------------------------------------

    def _convert_via_libreoffice(
        self,
        root: Path,
        source_abs: Path,
        rel_path: str,
    ) -> dict[str, Any]:
        """Convert a pptx/odp via LibreOffice + PyMuPDF.

        Spawns ``soffice --headless --convert-to pdf`` to produce
        an intermediate PDF in a temp directory, then routes that
        PDF through :meth:`_convert_via_pymupdf` with overridden
        display name and hash source so the provenance header
        records the original filename and hash.

        Graceful fallback — when either LibreOffice or PyMuPDF is
        unavailable, or when the soffice invocation fails for any
        reason (timeout, non-zero exit, missing output), falls
        back to the format-specific path:

        - ``.pptx`` → :meth:`_convert_via_python_pptx`
        - ``.odp`` → :meth:`_convert_via_markitdown`

        Fallback rather than error because the user asked for a
        conversion; producing some output (even lower-fidelity)
        beats failing the whole file.

        Temp directory lifetime is bounded by the method call —
        ``TemporaryDirectory`` cleans up on exit regardless of
        which branch returns.
        """
        # Pre-flight — check both deps before spending subprocess
        # time. shutil.which is cheap and doesn't launch soffice.
        soffice_path = shutil.which("soffice")
        if soffice_path is None:
            return self._libreoffice_fallback(
                root, source_abs, rel_path,
                reason="LibreOffice (soffice) not on PATH",
            )
        if not self._probe_import("fitz"):
            return self._libreoffice_fallback(
                root, source_abs, rel_path,
                reason="PyMuPDF not installed",
            )

        # Run LibreOffice in a temp dir. The --outdir flag tells
        # soffice where to write the PDF; it picks the filename
        # from the source stem.
        import subprocess
        import tempfile

        with tempfile.TemporaryDirectory(
            prefix="ac-dc-libreoffice-"
        ) as tmpdir:
            tmp_path = Path(tmpdir)
            try:
                proc = subprocess.run(
                    [
                        soffice_path,
                        "--headless",
                        "--convert-to", "pdf",
                        "--outdir", str(tmp_path),
                        str(source_abs),
                    ],
                    capture_output=True,
                    timeout=_LIBREOFFICE_TIMEOUT_SECONDS,
                    text=True,
                )
            except subprocess.TimeoutExpired:
                logger.debug(
                    "LibreOffice timed out for %s; falling back",
                    rel_path,
                )
                return self._libreoffice_fallback(
                    root, source_abs, rel_path,
                    reason="LibreOffice timed out",
                )
            except (OSError, subprocess.SubprocessError) as exc:
                logger.debug(
                    "LibreOffice subprocess failed for %s: %s; "
                    "falling back",
                    rel_path, exc,
                )
                return self._libreoffice_fallback(
                    root, source_abs, rel_path,
                    reason=f"LibreOffice launch failed: {exc}",
                )

            if proc.returncode != 0:
                logger.debug(
                    "LibreOffice exited %d for %s (stderr: %s); "
                    "falling back",
                    proc.returncode, rel_path,
                    proc.stderr.strip() if proc.stderr else "",
                )
                return self._libreoffice_fallback(
                    root, source_abs, rel_path,
                    reason=(
                        f"LibreOffice exited with code "
                        f"{proc.returncode}"
                    ),
                )

            # soffice names the output as {source_stem}.pdf in
            # the --outdir. Find it rather than assuming.
            expected_pdf = tmp_path / (source_abs.stem + ".pdf")
            if not expected_pdf.is_file():
                # Some locale / template configs produce
                # differently-named output. Fall back to scanning
                # the tmp dir for any .pdf.
                candidates = list(tmp_path.glob("*.pdf"))
                if not candidates:
                    logger.debug(
                        "LibreOffice produced no PDF for %s; "
                        "falling back",
                        rel_path,
                    )
                    return self._libreoffice_fallback(
                        root, source_abs, rel_path,
                        reason="LibreOffice produced no output",
                    )
                expected_pdf = candidates[0]

            # Route through the PyMuPDF pipeline. source_abs
            # stays as the original (.pptx/.odp) so output lands
            # next to the original, not in the temp dir. The
            # display_name and hash_source overrides ensure the
            # provenance header records the original file.
            # strip_text_when_present=False disables the
            # direct-PDF text-dedup pass: presentation text
            # labels the diagram shapes, so stripping it would
            # leave meaningless coloured rectangles. See
            # specs-reference/4-features/doc-convert.md
            # § "SVG text preservation in PDF pipeline".
            return self._convert_via_pymupdf(
                root=root,
                source_abs=source_abs,
                rel_path=rel_path,
                pdf_source=expected_pdf,
                display_name=source_abs.name,
                hash_source=source_abs,
                strip_text_when_present=False,
            )

    def _libreoffice_fallback(
        self,
        root: Path,
        source_abs: Path,
        rel_path: str,
        reason: str,
    ) -> dict[str, Any]:
        """Route to the format-specific fallback path.

        ``.pptx`` falls back to python-pptx (per-slide SVG
        rendering). ``.odp`` falls back to markitdown (plain-text
        extraction). Both are lower-fidelity than the LibreOffice
        path but produce SOMETHING — better than failing the
        whole conversion.
        """
        logger.debug(
            "LibreOffice path unavailable for %s: %s; "
            "using format-specific fallback",
            rel_path, reason,
        )
        suffix = source_abs.suffix.lower()
        if suffix == ".pptx":
            return self._convert_via_python_pptx(
                root, source_abs, rel_path
            )
        if suffix == ".odp":
            return self._convert_via_markitdown(
                root, source_abs, rel_path
            )
        # Shouldn't happen — caller only dispatches extensions
        # in _LIBREOFFICE_EXTENSIONS. Defensive.
        return self._fail(
            rel_path,
            f"No fallback available for {suffix}",
        )

    # ------------------------------------------------------------------
    # pdf — PyMuPDF hybrid text + SVG pipeline
    # ------------------------------------------------------------------

    def _convert_via_pymupdf(
        self,
        root: Path,
        source_abs: Path,
        rel_path: str,
        *,
        pdf_source: Path | None = None,
        display_name: str | None = None,
        hash_source: Path | None = None,
        strip_text_when_present: bool = True,
    ) -> dict[str, Any]:
        """Convert a PDF via PyMuPDF's hybrid text + SVG pipeline.

        For each page:

        - Text is extracted into markdown paragraphs via
          ``page.get_text("dict")``. Paragraphs are separated by
          blank lines. Font info is captured but not currently
          rendered (heading detection is a future enhancement).
        - Images and vector drawings are detected. If the page
          has any raster images OR at least
          :data:`_PAGE_GRAPHICS_THRESHOLD` significant drawings,
          a companion SVG is exported for that page.
        - SVGs preserve ``<text>`` elements when the page has no
          extractable text (figure-only pages, where the text
          likely labels the figure itself) OR when the caller
          sets ``strip_text_when_present=False`` (the LibreOffice
          → PDF → PyMuPDF route for pptx/odp, where the text IS
          the diagram). On direct-PDF pages that DO have
          extractable text, the `<text>` and `<tspan>` elements
          are stripped after export so the same prose doesn't
          appear twice — once in the markdown paragraphs, once
          embedded in the SVG.
        - Embedded raster images in the SVG are externalised —
          base64 data URIs replaced with relative file refs.
        - Text-only pages produce no SVG. Pages with no text AND
          no detected images/drawings still get a full-page SVG
          as a fallback, so lightweight vector content isn't
          silently dropped.

        Output layout:

            docs/report.pdf                   ← source
            docs/report.md                     ← index + text
            docs/report/
                02_page.svg                    ← page 2 (had figures)
                05_page.svg                    ← page 5 (had charts)
                02_page_img01.png              ← externalized raster

        Fails with a per-file error when PyMuPDF isn't installed.
        Unlike xlsx (which falls back to markitdown), PyMuPDF is
        the only reliable PDF extractor — no fallback.

        Parameters
        ----------
        root:
            Repository root.
        source_abs:
            Absolute path used to compute the output location.
            Pass A5b note — when converting via LibreOffice,
            this remains the ORIGINAL source (.pptx/.odp) so the
            output markdown lands next to the original, not next
            to the intermediate PDF in the temp dir.
        rel_path:
            Relative path for per-file result reporting.
        pdf_source:
            Optional — when set, PyMuPDF opens this PDF instead
            of `source_abs`. Used by Pass A5b to route an
            intermediate PDF produced by LibreOffice through
            the pipeline while keeping output paths anchored to
            the original source.
        display_name:
            Optional — what appears in `source=` of the
            provenance header. Defaults to `source_abs.name`.
            Pass A5b uses this so a converted .pptx records
            `source=deck.pptx`, not `source=deck.pdf`.
        hash_source:
            Optional — file to hash for the provenance header.
            Defaults to `source_abs`. Pass A5b hashes the
            original .pptx so re-running against an unchanged
            source classifies as `current` regardless of whether
            LibreOffice produces byte-identical intermediate
            PDFs across runs (it doesn't — timestamps vary).
        strip_text_when_present:
            When True (the default, used by the direct-PDF
            path), pages that have extractable text get their
            ``<text>`` / ``<tspan>`` elements stripped from the
            generated SVG. The same text already appears in the
            companion markdown, so duplicating it in the SVG
            bloats output without benefit for real PDFs (papers,
            reports). When False (used by the LibreOffice route
            for pptx/odp), SVG text is preserved unconditionally
            — presentation text labels the diagram shapes, and
            stripping it would leave meaningless coloured
            rectangles.
        """
        # Lazy import — PyMuPDF is optional in stripped-down
        # releases. Clean error with install hint on ImportError.
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return self._fail(
                rel_path,
                (
                    "PyMuPDF is not installed. Install with: "
                    "pip install 'ac-dc[docs]'"
                ),
            )

        output_abs = source_abs.with_suffix(".md")
        try:
            output_rel = output_abs.relative_to(root)
        except ValueError:
            return self._fail(
                rel_path,
                "Output path escapes repository root",
            )

        # Hash source for provenance. Use hash_source when given
        # (A5b path) so the hash reflects the original file the
        # user actually edits, not the intermediate PDF.
        hash_target = hash_source if hash_source is not None else source_abs
        try:
            source_hash = self._hash_file(hash_target)
        except OSError as exc:
            return self._fail(
                rel_path, f"Source hash failed: {exc}"
            )

        # Resolve display name for the provenance header. Defaults
        # to the original source's basename so converted files
        # carry the user-recognisable name, not the intermediate
        # PDF's.
        resolved_display_name = (
            display_name if display_name is not None
            else source_abs.name
        )

        # Open the document. Broad catch — corrupt PDFs, wrong
        # version, encrypted without a password all produce
        # various PyMuPDF exception types. When pdf_source is
        # given, open that instead of the original source — the
        # A5b path hands us an intermediate PDF to process.
        open_target = pdf_source if pdf_source is not None else source_abs
        try:
            doc = fitz.open(str(open_target))
        except Exception as exc:
            return self._fail(
                rel_path,
                f"PyMuPDF failed to open: {exc}",
            )

        try:
            return self._process_pdf_document(
                doc=doc,
                root=root,
                source_abs=source_abs,
                output_abs=output_abs,
                output_rel=output_rel,
                source_hash=source_hash,
                rel_path=rel_path,
                display_name=resolved_display_name,
                strip_text_when_present=strip_text_when_present,
            )
        finally:
            # Always close — PyMuPDF holds file handles.
            try:
                doc.close()
            except Exception:
                pass

    def _process_pdf_document(
        self,
        doc: Any,
        root: Path,
        source_abs: Path,
        output_abs: Path,
        output_rel: Path,
        source_hash: str,
        rel_path: str,
        display_name: str | None = None,
        strip_text_when_present: bool = True,
    ) -> dict[str, Any]:
        """Walk the pages of an open PDF document and emit output.

        Split out from :meth:`_convert_via_pymupdf` so the
        `doc.close()` is guaranteed in the caller's finally
        block regardless of which branch we exit through.

        ``display_name`` defaults to the original source's
        basename when None — used by Pass A5b to override for
        converted pptx/odp (the provenance header shows the
        original filename, not the intermediate PDF's).

        ``strip_text_when_present`` is forwarded to
        :meth:`_process_pdf_page`; see :meth:`_convert_via_pymupdf`
        for the full rationale.
        """
        page_count = doc.page_count
        if page_count == 0:
            # Empty PDF — placeholder output so scan classifies
            # as `current`.
            return self._write_pdf_output(
                output_abs=output_abs,
                output_rel=output_rel,
                source_abs=source_abs,
                source_hash=source_hash,
                markdown_text="(empty PDF)\n",
                rel_path=rel_path,
                artefacts=(),
                display_name=display_name,
            )

        # Per-page filename width — 2 digits for small PDFs,
        # 3+ for larger ones. Matches the pptx path.
        pad_width = max(
            _SLIDE_NUMBER_MIN_WIDTH, len(str(page_count))
        )

        # Assets subdir — created lazily on the first page that
        # actually needs it. Avoids empty subdirs for text-only
        # PDFs.
        assets_dir = source_abs.with_suffix("")
        assets_created = False

        body_parts: list[str] = []
        artefacts: list[str] = []  # all files under assets_dir
        prior_artefacts = self._read_prior_images(output_abs)

        for page_index in range(page_count):
            try:
                page = doc.load_page(page_index)
            except Exception as exc:
                logger.debug(
                    "PDF page %d load failed for %s: %s",
                    page_index, rel_path, exc,
                )
                body_parts.append(
                    f"## Page {page_index + 1}\n\n"
                    "*(page load failed)*"
                )
                continue

            try:
                page_result = self._process_pdf_page(
                    page=page,
                    page_index=page_index,
                    pad_width=pad_width,
                    assets_dir=assets_dir,
                    assets_created=assets_created,
                    strip_text_when_present=strip_text_when_present,
                )
            except Exception as exc:
                logger.debug(
                    "PDF page %d of %s render failed: %s",
                    page_index + 1, rel_path, exc,
                )
                body_parts.append(
                    f"## Page {page_index + 1}\n\n"
                    "*(page rendering failed)*"
                )
                continue

            assets_created = assets_created or page_result["assets_created"]
            body_parts.append(page_result["markdown"])
            artefacts.extend(page_result["artefacts"])

        markdown_text = "\n\n".join(body_parts) + "\n"

        # Orphan cleanup — anything listed in the old provenance
        # header that we didn't re-produce this round gets
        # unlinked.
        if prior_artefacts and assets_dir.is_dir():
            new_set = set(artefacts)
            for orphan in prior_artefacts:
                if orphan in new_set:
                    continue
                orphan_path = assets_dir / orphan
                try:
                    orphan_path.unlink()
                except OSError as exc:
                    logger.debug(
                        "Failed to remove orphan artefact %s: %s",
                        orphan_path, exc,
                    )

        # If the assets dir was created but is now empty (every
        # artefact was an orphan removed above), clean it up.
        if assets_dir.is_dir():
            try:
                if not any(assets_dir.iterdir()):
                    assets_dir.rmdir()
            except OSError as exc:
                logger.debug(
                    "Failed to remove empty assets dir %s: %s",
                    assets_dir, exc,
                )

        return self._write_pdf_output(
            output_abs=output_abs,
            output_rel=output_rel,
            source_abs=source_abs,
            source_hash=source_hash,
            markdown_text=markdown_text,
            rel_path=rel_path,
            artefacts=tuple(artefacts),
            display_name=display_name,
        )

    def _process_pdf_page(
        self,
        page: Any,
        page_index: int,
        pad_width: int,
        assets_dir: Path,
        assets_created: bool,
        strip_text_when_present: bool = True,
    ) -> dict[str, Any]:
        """Emit markdown + optional SVG for one PDF page.

        Returns ``{"markdown": str, "artefacts": list[str],
        "assets_created": bool}`` where ``artefacts`` is the list
        of filenames (SVG plus externalized images) produced for
        this page.

        Page processing logic:

        1. Extract text — if non-empty, emit as markdown paragraphs
        2. Detect images and significant drawings
        3. If the page has ANY raster images, emit an SVG AND
           embed image refs in markdown (both places are visible
           to the LLM — markdown for grep, SVG for visual
           fidelity)
        4. Else if the page has significant drawings AND text,
           emit a companion SVG (for the graphics)
        5. Else if the page has NO text AND NO detected content,
           emit a full-page SVG as a safety net (lightweight
           vector graphics that don't reach the "significant"
           threshold)
        6. Text-only pages emit no SVG

        ``<text>`` preservation in the emitted SVG depends on
        ``strip_text_when_present`` AND whether the page has
        extractable text. Direct-PDF pages with text strip the
        SVG's ``<text>`` / ``<tspan>`` elements (the markdown
        already carries the paragraphs); LibreOffice-routed
        pptx/odp pages always keep text (it labels the diagram
        shapes); figure-only pages keep text regardless of the
        flag (the text likely labels the figure). See
        :meth:`_export_pdf_page_svg` for the strip implementation.
        """
        page_number = page_index + 1
        slide_name = f"{str(page_number).zfill(pad_width)}_page.svg"

        # Extract text.
        text_paragraphs = self._extract_pdf_text(page)
        has_text = bool(text_paragraphs)

        # Detect images and drawings.
        raster_images = self._count_pdf_raster_images(page)
        significant_drawings = self._count_significant_drawings(page)
        has_raster = raster_images > 0
        has_significant_graphics = (
            significant_drawings >= _PAGE_GRAPHICS_THRESHOLD
        )

        # Decide whether to emit an SVG for this page.
        emit_svg = (
            has_raster
            or has_significant_graphics
            or (not has_text and not has_raster)
        )

        markdown_parts: list[str] = [f"## Page {page_number}"]
        artefacts: list[str] = []

        if has_text:
            markdown_parts.extend(text_paragraphs)

        if emit_svg:
            # Ensure assets dir exists.
            if not assets_created:
                try:
                    assets_dir.mkdir(parents=True, exist_ok=True)
                    assets_created = True
                except OSError as exc:
                    logger.debug(
                        "Assets dir create failed %s: %s",
                        assets_dir, exc,
                    )
                    # Without the dir we can't write the SVG —
                    # skip the SVG for this page and continue.
                    return {
                        "markdown": "\n\n".join(markdown_parts),
                        "artefacts": [],
                        "assets_created": assets_created,
                    }

            # Decide whether to strip SVG text for this page.
            # Two conditions must both hold:
            #   1. The caller asked for stripping (direct-PDF
            #      path; LibreOffice route passes False to
            #      preserve diagram labels).
            #   2. The page has extractable text — otherwise
            #      the SVG's <text> elements probably ARE the
            #      figure labels and stripping would lose them.
            # See specs-reference/4-features/doc-convert.md
            # § "SVG text preservation in PDF pipeline".
            strip_svg_text = strip_text_when_present and has_text
            svg_text, image_files = self._export_pdf_page_svg(
                page=page,
                svg_name_stem=slide_name.removesuffix(".svg"),
                assets_dir=assets_dir,
                strip_text=strip_svg_text,
            )
            svg_path = assets_dir / slide_name
            try:
                svg_path.write_text(svg_text, encoding="utf-8")
                artefacts.append(slide_name)
                artefacts.extend(image_files)
            except OSError as exc:
                logger.debug(
                    "SVG write failed %s: %s", svg_path, exc
                )

            # Add markdown image ref so LLM sees the SVG link.
            # Two cases:
            # - Page has text + images: markdown has text AND a
            #   link to the SVG (visual fidelity is in the SVG).
            # - Page has no text: the SVG is the content.
            rel_ref = f"{assets_dir.name}/{slide_name}"
            markdown_parts.append(f"![Page {page_number}]({rel_ref})")

        if len(markdown_parts) == 1:
            # Only the heading — page had no text and no SVG.
            # This shouldn't happen given the fallback logic
            # but be defensive.
            markdown_parts.append("*(blank page)*")

        return {
            "markdown": "\n\n".join(markdown_parts),
            "artefacts": artefacts,
            "assets_created": assets_created,
        }

    def _extract_pdf_text(self, page: Any) -> list[str]:
        """Extract text from a page as markdown paragraphs.

        Uses ``page.get_text("dict")`` which returns structured
        data. Each text block becomes one paragraph; spans
        within lines are joined by spaces; lines within a block
        by spaces too. A future enhancement could detect
        heading levels from font sizes.

        Returns an empty list when the page has no extractable
        text.
        """
        try:
            data = page.get_text("dict")
        except Exception as exc:
            logger.debug("get_text failed: %s", exc)
            return []

        paragraphs: list[str] = []
        for block in data.get("blocks", []):
            if block.get("type", 0) != 0:
                # type 1 is image; skip.
                continue
            lines_text: list[str] = []
            for line in block.get("lines", []):
                spans: list[str] = []
                for span in line.get("spans", []):
                    span_text = span.get("text", "") or ""
                    if span_text.strip():
                        spans.append(span_text)
                if spans:
                    lines_text.append(" ".join(spans))
            if lines_text:
                # Join lines within a block with spaces — block
                # already represents a visually-grouped run of
                # text, so line breaks within it are usually
                # visual-wrap rather than semantic-paragraph.
                paragraph = " ".join(lines_text).strip()
                if paragraph:
                    paragraphs.append(paragraph)
        return paragraphs

    @staticmethod
    def _count_pdf_raster_images(page: Any) -> int:
        """Return the number of raster images on the page."""
        try:
            return len(page.get_images())
        except Exception as exc:
            logger.debug("get_images failed: %s", exc)
            return 0

    @staticmethod
    def _count_significant_drawings(page: Any) -> int:
        """Return the count of "significant" vector drawings.

        Significance rules (from specs4/4-features/doc-convert.md):

        - Any drawing containing Bézier (``c``) or quadratic
          (``qu``) curves → significant
        - A filled path with more than
          :data:`_POLYGON_SIGNIFICANT_SEGMENTS` segments →
          significant
        - Any drawing with more than
          :data:`_PATH_SIGNIFICANT_SEGMENTS` segments →
          significant
        - Simple rectangles and single lines → NOT significant
          (these are just borders and table rules that every
          PDF emits for layout)
        """
        try:
            drawings = page.get_drawings()
        except Exception as exc:
            logger.debug("get_drawings failed: %s", exc)
            return 0

        count = 0
        for drawing in drawings:
            items = drawing.get("items", [])
            if not items:
                continue

            # Check for curves (always significant).
            has_curves = any(
                item and len(item) > 0 and item[0] in ("c", "qu")
                for item in items
            )
            if has_curves:
                count += 1
                continue

            # Check for filled path with multiple segments.
            is_filled = drawing.get("fill") is not None
            if (
                is_filled
                and len(items) > _POLYGON_SIGNIFICANT_SEGMENTS
            ):
                count += 1
                continue

            # Check for complex path (many segments).
            if len(items) > _PATH_SIGNIFICANT_SEGMENTS:
                count += 1
                continue

        return count

    def _export_pdf_page_svg(
        self,
        page: Any,
        svg_name_stem: str,
        assets_dir: Path,
        strip_text: bool = False,
    ) -> tuple[str, list[str]]:
        """Export a page to SVG, externalizing any raster images.

        PyMuPDF emits ``<text>`` elements (text_as_path=0) so the
        output is compact and selectable. What happens next
        depends on ``strip_text``:

        - ``strip_text=False`` (figure-only pages, and every
          page on the LibreOffice → PDF → PyMuPDF route):
          ``<text>`` and ``<tspan>`` elements are preserved.
          For presentations this matters — diagram labels like
          "Runtime Environment" or "Calibration Unit" anchor
          the coloured shapes, and dropping them leaves the
          user staring at nameless rectangles.
        - ``strip_text=True`` (direct-PDF pages with
          extractable text): ``<text>`` and ``<tspan>``
          elements are removed. The same prose already appears
          in the companion markdown file as extracted
          paragraphs, and duplicating it inside the SVG just
          bloats output for real PDFs (papers, reports).

        See specs-reference/4-features/doc-convert.md
        § "SVG text preservation in PDF pipeline" for the
        rationale behind the origin-aware behaviour.

        Parameters
        ----------
        page:
            A PyMuPDF ``Page`` object.
        svg_name_stem:
            Stem used for naming externalized raster image files
            (e.g. ``"02_page"`` → ``02_page_img01.png``).
        assets_dir:
            Directory to save externalized images into.
        strip_text:
            When True, remove ``<text>`` / ``<tspan>`` elements
            from the generated SVG. Default False so callers
            that don't care (tests, figure-only pages) get the
            safe preserve-everything behaviour.

        Returns
        -------
        tuple
            ``(svg_text, externalized_filenames)``.
            ``svg_text`` has any base64 raster images rewritten
            to refer to externalized files and, when
            ``strip_text`` is True, has ``<text>`` / ``<tspan>``
            elements removed.
        """
        try:
            # text_as_path=0 keeps text as <text> elements rather
            # than decomposing into paths. Makes the SVG
            # selectable in a viewer and much smaller.
            svg_text = page.get_svg_image(text_as_path=0)
        except Exception as exc:
            logger.debug("get_svg_image failed: %s", exc)
            # Emit a minimal empty SVG so callers don't crash.
            return (
                '<svg xmlns="http://www.w3.org/2000/svg"/>',
                [],
            )

        # Externalize embedded raster images.
        svg_text, image_files = self._externalize_svg_images(
            svg_text=svg_text,
            stem=svg_name_stem,
            assets_dir=assets_dir,
        )

        # Strip text elements on the direct-PDF path. Runs
        # after image externalization so we can't accidentally
        # strip text from inside an <image> tag's href (not
        # that PyMuPDF would ever produce that, but the
        # ordering keeps each pass independent and easy to
        # reason about).
        if strip_text:
            svg_text = self._strip_svg_text_elements(svg_text)

        return svg_text, image_files

    def _externalize_svg_images(
        self,
        svg_text: str,
        stem: str,
        assets_dir: Path,
    ) -> tuple[str, list[str]]:
        """Extract base64 images from an SVG, save as files.

        PyMuPDF's SVG output embeds raster images as base64
        data URIs inside ``<image>`` elements. For large
        images this bloats the SVG severely; externalising
        them into sibling files keeps the SVG compact and
        matches the approach the DocConvert tab uses
        elsewhere.

        Returns the modified SVG text (with ``data:image/...``
        references replaced by relative filename refs) and
        the list of saved filenames.

        Failures (decode error, write error) leave the original
        data URI in place — broken image ref is better than
        silent content loss.
        """
        saved: list[str] = []

        def _replace(match: re.Match[str]) -> str:
            attr_name = match.group("attr")
            mime_sub = match.group("mime").lower()
            payload = match.group("payload")
            try:
                image_bytes = base64.b64decode(
                    payload, validate=False
                )
            except Exception as exc:
                logger.debug(
                    "SVG image decode failed: %s", exc
                )
                return match.group(0)

            ext = _MIME_TO_EXT.get(mime_sub, ".bin")
            image_index = len(saved) + 1
            image_name = f"{stem}_img{image_index:02d}{ext}"
            image_path = assets_dir / image_name
            try:
                image_path.write_bytes(image_bytes)
            except OSError as exc:
                logger.debug(
                    "SVG image write failed %s: %s",
                    image_path, exc,
                )
                return match.group(0)

            saved.append(image_name)
            # Return the attribute with the relative filename.
            return f'{attr_name}="{image_name}"'

        # Match both `href="data:..."` and
        # `xlink:href="data:..."` attribute forms.
        pattern = re.compile(
            r'(?P<attr>(?:xlink:)?href)='
            r'"data:image/(?P<mime>[^;]+);base64,'
            r'(?P<payload>[^"]+)"',
            re.IGNORECASE,
        )
        modified = pattern.sub(_replace, svg_text)
        return modified, saved

    @staticmethod
    def _strip_svg_text_elements(svg_text: str) -> str:
        """Remove ``<text>...</text>`` and any leftover ``<tspan>``
        elements from an SVG string.

        Used by :meth:`_export_pdf_page_svg` on the direct-PDF
        path when the page has extractable text — the markdown
        already carries the prose, so keeping it in the SVG too
        is just duplication.

        Regex rather than XML parse: PyMuPDF's output is
        consistently structured and regex keeps us dependency-
        free and fast. ``re.DOTALL`` matters for multi-line
        ``<text>`` blocks where tspans span several lines.
        ``re.IGNORECASE`` is defensive — SVG tag names are
        normatively lowercase but case-insensitive matching
        costs nothing and protects against edge cases.

        Two passes:
        1. Strip whole ``<text>...</text>`` blocks (which
           includes any nested tspans).
        2. Strip any stray ``<tspan>...</tspan>`` blocks that
           somehow survived (e.g. tspan outside a text parent —
           invalid but possible in malformed SVG).

        Self-closing variants (``<text ... />``) are rare in
        PyMuPDF output but handled by the first pass too.
        """
        # <text ... />  — self-closing.
        svg_text = re.sub(
            r"<text\b[^>]*/\s*>",
            "",
            svg_text,
            flags=re.IGNORECASE,
        )
        # <text ...>...</text>  — block form.
        svg_text = re.sub(
            r"<text\b[^>]*>.*?</text\s*>",
            "",
            svg_text,
            flags=re.DOTALL | re.IGNORECASE,
        )
        # Stray <tspan>...</tspan> outside a parent text.
        svg_text = re.sub(
            r"<tspan\b[^>]*>.*?</tspan\s*>",
            "",
            svg_text,
            flags=re.DOTALL | re.IGNORECASE,
        )
        # Self-closing tspan.
        svg_text = re.sub(
            r"<tspan\b[^>]*/\s*>",
            "",
            svg_text,
            flags=re.IGNORECASE,
        )
        return svg_text

    def _write_pdf_output(
        self,
        output_abs: Path,
        output_rel: Path,
        source_abs: Path,
        source_hash: str,
        markdown_text: str,
        rel_path: str,
        artefacts: tuple[str, ...],
        display_name: str | None = None,
    ) -> dict[str, Any]:
        """Write PDF pipeline output with provenance header.

        Shared for the normal and empty-PDF cases. The
        ``artefacts`` tuple lists every file produced under the
        assets subdirectory — page SVGs AND externalized images
        — so the orphan-cleanup pass on re-conversion can diff
        against this list.

        ``display_name`` defaults to ``source_abs.name`` — Pass
        A5b overrides it so converted pptx/odp files record the
        original filename in provenance, not the intermediate
        PDF's.
        """
        provenance_line = self._build_provenance_header(
            source_name=display_name or source_abs.name,
            source_hash=source_hash,
            images=artefacts,
        )
        final_content = (
            provenance_line + "\n\n" + markdown_text.lstrip("\n")
        )

        try:
            output_abs.parent.mkdir(parents=True, exist_ok=True)
            output_abs.write_text(final_content, encoding="utf-8")
        except OSError as exc:
            return self._fail(
                rel_path,
                f"Failed to write output: {exc}",
            )

        return {
            "path": rel_path,
            "status": "ok",
            "output_path": str(output_rel).replace("\\", "/"),
            "images": list(artefacts),
        }

    # ------------------------------------------------------------------
    # pptx — python-pptx fallback (renders each slide as an SVG)
    # ------------------------------------------------------------------

    def _convert_via_python_pptx(
        self,
        root: Path,
        source_abs: Path,
        rel_path: str,
    ) -> dict[str, Any]:
        """Convert a .pptx file via the python-pptx fallback.

        Each slide renders as a standalone SVG containing the
        slide's text shapes, embedded images, and tables. The
        output markdown is an index file with per-slide headings
        and image references to the SVGs.

        Output layout:

            docs/presentation.pptx            ← source
            docs/presentation.md               ← index with links
            docs/presentation/
                01_slide.svg
                02_slide.svg
                ...

        Slide filenames use zero-padded numbering (`01_slide.svg`)
        so file listings sort correctly regardless of deck size.
        The assets subdirectory is always created for pptx (unlike
        markitdown docs where it's conditional on image presence)
        because every slide produces an SVG.

        Fails with a per-file error when python-pptx isn't
        installed. Unlike xlsx which falls back to markitdown,
        pptx has no markitdown-equivalent fallback in the current
        release — Pass A5 will add the LibreOffice + PyMuPDF
        primary path.
        """
        # Lazy import — python-pptx is optional in stripped-down
        # releases. A surface-level error with install hint
        # matches the markitdown-missing pattern.
        try:
            from pptx import Presentation
        except ImportError:
            return self._fail(
                rel_path,
                (
                    "python-pptx is not installed. Install with: "
                    "pip install 'ac-dc[docs]'"
                ),
            )

        output_abs = source_abs.with_suffix(".md")
        try:
            output_rel = output_abs.relative_to(root)
        except ValueError:
            return self._fail(
                rel_path,
                "Output path escapes repository root",
            )

        # Hash source for provenance — matches the markitdown and
        # xlsx paths so status classification works uniformly.
        try:
            source_hash = self._hash_file(source_abs)
        except OSError as exc:
            return self._fail(
                rel_path, f"Source hash failed: {exc}"
            )

        # Open the presentation. python-pptx's failure modes on
        # invalid pptx vary — corrupt zip, missing core XML,
        # unsupported schema version. Broad catch to keep the
        # pipeline robust.
        try:
            presentation = Presentation(str(source_abs))
        except Exception as exc:
            return self._fail(
                rel_path,
                f"python-pptx failed to open: {exc}",
            )

        # Slide dimensions — fall back to the 4:3 default if
        # python-pptx reports None (rare but defensive).
        slide_width = presentation.slide_width or _DEFAULT_SLIDE_WIDTH_EMU
        slide_height = presentation.slide_height or _DEFAULT_SLIDE_HEIGHT_EMU
        svg_width_px = slide_width * _EMU_TO_PX
        svg_height_px = slide_height * _EMU_TO_PX

        slides = list(presentation.slides)
        if not slides:
            # Empty deck — produce an informative placeholder
            # output so the scan classifies it as `current`.
            markdown_text = "(empty presentation)\n"
            return self._write_pptx_output(
                output_abs=output_abs,
                output_rel=output_rel,
                source_abs=source_abs,
                source_hash=source_hash,
                markdown_text=markdown_text,
                rel_path=rel_path,
            )

        # Pick zero-padding width — larger of the default minimum
        # and whatever the deck actually needs. A 150-slide deck
        # pads to 3 digits.
        pad_width = max(
            _SLIDE_NUMBER_MIN_WIDTH,
            len(str(len(slides))),
        )

        # Per-source assets subdirectory. Unlike the markitdown
        # path, we always create it here — every slide produces
        # an SVG, so the dir is never empty after a successful
        # conversion.
        assets_dir = source_abs.with_suffix("")
        try:
            assets_dir.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            return self._fail(
                rel_path,
                f"Failed to create assets dir: {exc}",
            )

        # Render each slide. Failures on a single slide are
        # logged and that slide becomes a placeholder in the
        # index — better than aborting the whole deck.
        index_entries: list[str] = []
        saved_slides: list[str] = []
        for idx, slide in enumerate(slides, start=1):
            slide_name = f"{str(idx).zfill(pad_width)}_slide.svg"
            try:
                svg_text = self._render_pptx_slide(
                    slide, svg_width_px, svg_height_px
                )
            except Exception as exc:
                logger.debug(
                    "pptx slide %d of %s render failed: %s",
                    idx, rel_path, exc,
                )
                # Index entry still appears, but without an
                # image link — user sees which slide failed.
                index_entries.append(
                    f"## Slide {idx}\n\n*(rendering failed)*"
                )
                continue

            svg_path = assets_dir / slide_name
            try:
                svg_path.write_text(svg_text, encoding="utf-8")
            except OSError as exc:
                logger.debug(
                    "Failed to write %s: %s", svg_path, exc
                )
                index_entries.append(
                    f"## Slide {idx}\n\n*(write failed)*"
                )
                continue

            saved_slides.append(slide_name)
            # Index entry — heading + image reference. The image
            # ref uses the assets dir name as a relative prefix
            # so the markdown renders correctly from its sibling
            # location.
            rel_ref = f"{assets_dir.name}/{slide_name}"
            index_entries.append(
                f"## Slide {idx}\n\n![Slide {idx}]({rel_ref})"
            )

        markdown_text = "\n\n".join(index_entries) + "\n"

        # Orphan cleanup — re-conversion of a deck with fewer
        # slides than before leaves stale SVGs. Read the prior
        # provenance header (if any) and delete any images
        # listed there that we didn't produce this round.
        prior_images = self._read_prior_images(output_abs)
        if prior_images:
            saved_set = set(saved_slides)
            for orphan in prior_images:
                if orphan in saved_set:
                    continue
                orphan_path = assets_dir / orphan
                try:
                    orphan_path.unlink()
                except OSError as exc:
                    logger.debug(
                        "Failed to remove orphan slide %s: %s",
                        orphan_path, exc,
                    )

        return self._write_pptx_output(
            output_abs=output_abs,
            output_rel=output_rel,
            source_abs=source_abs,
            source_hash=source_hash,
            markdown_text=markdown_text,
            rel_path=rel_path,
            images=tuple(saved_slides),
        )

    def _write_pptx_output(
        self,
        output_abs: Path,
        output_rel: Path,
        source_abs: Path,
        source_hash: str,
        markdown_text: str,
        rel_path: str,
        images: tuple[str, ...] = (),
    ) -> dict[str, Any]:
        """Write the pptx index markdown with provenance header.

        Shared between the normal path and the empty-deck case.
        Builds the provenance header, prepends it to the markdown
        body, writes atomically, returns a per-file result dict.
        """
        provenance_line = self._build_provenance_header(
            source_name=source_abs.name,
            source_hash=source_hash,
            images=images,
        )
        final_content = (
            provenance_line + "\n\n" + markdown_text.lstrip("\n")
        )

        try:
            output_abs.parent.mkdir(parents=True, exist_ok=True)
            output_abs.write_text(final_content, encoding="utf-8")
        except OSError as exc:
            return self._fail(
                rel_path,
                f"Failed to write output: {exc}",
            )

        return {
            "path": rel_path,
            "status": "ok",
            "output_path": str(output_rel).replace("\\", "/"),
            "images": list(images),
        }

    def _read_prior_images(
        self,
        output_abs: Path,
    ) -> tuple[str, ...]:
        """Return the images list from an existing output's header.

        Used by the orphan-cleanup path. Empty tuple when no
        prior output exists or the header is absent / malformed.
        """
        if not output_abs.is_file():
            return ()
        header = self._read_provenance_header(output_abs)
        if header is None:
            return ()
        return header.images

    # ------------------------------------------------------------------
    # pptx slide rendering
    # ------------------------------------------------------------------

    def _render_pptx_slide(
        self,
        slide: Any,
        width_px: float,
        height_px: float,
    ) -> str:
        """Render one slide as a complete SVG document.

        Walks the slide's shapes and emits SVG elements for each.
        Supported shape types:

        - Text frames (including within placeholders like title,
          body, content): `<text>` with font-size, weight, fill,
          and alignment derived from the first run's properties
        - Pictures: `<image>` with the raster content embedded
          as a base64 data URI, preserving position and size
        - Tables: `<rect>` borders + `<text>` cell content,
          rendered as a simple grid (no merged cells)

        Unsupported shape types (charts, SmartArt, OLE objects,
        groups) are skipped with a debug log. A slide with only
        unsupported shapes produces an empty SVG frame — still
        valid output, the user can see the slide exists but the
        content isn't surfaced.
        """
        # SVG header — viewBox lets the rendered image scale to
        # any container size while preserving aspect ratio.
        parts: list[str] = [
            (
                f'<svg xmlns="http://www.w3.org/2000/svg" '
                f'xmlns:xlink="http://www.w3.org/1999/xlink" '
                f'width="{int(width_px)}" height="{int(height_px)}" '
                f'viewBox="0 0 {int(width_px)} {int(height_px)}">'
            ),
            # White background — matches default PowerPoint
            # slide colour. Themes with non-white backgrounds
            # are a future enhancement.
            (
                f'<rect width="{int(width_px)}" '
                f'height="{int(height_px)}" fill="#ffffff"/>'
            ),
        ]

        for shape in slide.shapes:
            try:
                element = self._render_pptx_shape(shape)
            except Exception as exc:
                logger.debug(
                    "Skipping shape %r: %s",
                    getattr(shape, "shape_type", "?"), exc,
                )
                continue
            if element:
                parts.append(element)

        parts.append("</svg>")
        return "\n".join(parts)

    def _render_pptx_shape(self, shape: Any) -> str:
        """Dispatch a single shape to its renderer.

        Returns an SVG element string or empty string for
        skipped shapes. Never raises — all render errors are
        caught at the call site and logged.
        """
        # Dimensions in EMU, converted to pixels.
        left_px = (shape.left or 0) * _EMU_TO_PX
        top_px = (shape.top or 0) * _EMU_TO_PX
        width_px = (shape.width or 0) * _EMU_TO_PX
        height_px = (shape.height or 0) * _EMU_TO_PX

        # Picture — render as <image> with embedded base64 payload.
        if self._is_picture(shape):
            return self._render_picture(
                shape, left_px, top_px, width_px, height_px
            )

        # Table — render as grid of <rect> + <text>.
        if getattr(shape, "has_table", False):
            return self._render_table(
                shape.table, left_px, top_px, width_px, height_px
            )

        # Text frame (including placeholder shapes).
        if getattr(shape, "has_text_frame", False):
            return self._render_text_frame(
                shape.text_frame, left_px, top_px, width_px, height_px
            )

        # Unsupported shape type (chart, SmartArt, group, etc.).
        # Return empty string — caller skips.
        return ""

    @staticmethod
    def _is_picture(shape: Any) -> bool:
        """Return True when the shape is a raster image.

        python-pptx's shape_type constants live under
        `MSO_SHAPE_TYPE`; the PICTURE enum value is 13. Rather
        than importing the enum (which couples us to a specific
        python-pptx version), we probe the attributes the
        Picture class exposes: `.image.blob` returns the raw
        bytes.
        """
        image = getattr(shape, "image", None)
        if image is None:
            return False
        # Image object exposes .blob. A non-picture shape with
        # an "image" attr (rare) wouldn't have this.
        return hasattr(image, "blob")

    def _render_picture(
        self,
        shape: Any,
        left_px: float,
        top_px: float,
        width_px: float,
        height_px: float,
    ) -> str:
        """Render a picture shape as an SVG <image> element.

        The raster payload is base64-encoded and inlined via a
        data URI. Keeps slide layout self-contained — one file
        per slide, no external image refs.

        Returns empty string on any read error — a damaged
        image shouldn't abort the slide.
        """
        try:
            image = shape.image
            blob = image.blob
            content_type = image.content_type or "image/png"
        except Exception as exc:
            logger.debug("Failed to read picture blob: %s", exc)
            return ""

        encoded = base64.b64encode(blob).decode("ascii")
        data_uri = f"data:{content_type};base64,{encoded}"
        return (
            f'<image x="{left_px:.2f}" y="{top_px:.2f}" '
            f'width="{width_px:.2f}" height="{height_px:.2f}" '
            f'xlink:href="{data_uri}"/>'
        )

    def _render_text_frame(
        self,
        text_frame: Any,
        left_px: float,
        top_px: float,
        width_px: float,
        height_px: float,
    ) -> str:
        """Render a text frame as a <g> of <text> elements.

        Each paragraph in the frame becomes a `<text>` line
        positioned vertically by cumulative line height. Font
        properties (size, weight, colour, alignment) come from
        the first run of each paragraph — the spec's
        simplification for Pass A4 is "one style per
        paragraph", not per-run.

        Returns empty string for frames with no rendered text —
        avoids emitting degenerate `<g></g>` wrappers.
        """
        # Collect rendered paragraphs.
        lines: list[str] = []
        cursor_y = top_px
        for paragraph in text_frame.paragraphs:
            line = self._render_paragraph(
                paragraph, left_px, cursor_y, width_px
            )
            if line:
                lines.append(line["svg"])
                cursor_y += line["height"]

        if not lines:
            return ""
        return "<g>\n" + "\n".join(lines) + "\n</g>"

    def _render_paragraph(
        self,
        paragraph: Any,
        x_px: float,
        y_px: float,
        width_px: float,
    ) -> dict[str, Any] | None:
        """Render a paragraph as a single `<text>` element.

        Returns a dict with `svg` (the element string) and
        `height` (vertical space consumed) for cumulative
        layout. None when the paragraph has no text — caller
        skips.

        Font properties are extracted from the first non-empty
        run; subsequent runs inherit these. Richer handling
        (per-run formatting) is deferred to Pass A5's richer
        pipeline.
        """
        text = self._collect_paragraph_text(paragraph)
        if not text:
            return None

        # Extract properties from the first run — python-pptx
        # exposes them under `paragraph.runs[N].font`.
        font_size_pt = _DEFAULT_FONT_SIZE_PT
        font_weight = "normal"
        font_style = "normal"
        font_color = _DEFAULT_FONT_COLOR
        try:
            first_run = next(
                (r for r in paragraph.runs if r.text),
                None,
            )
            if first_run is not None:
                font = first_run.font
                if font.size is not None:
                    # python-pptx returns Emu objects for sizes;
                    # .pt gives points directly.
                    font_size_pt = int(font.size.pt)
                if font.bold:
                    font_weight = "bold"
                if font.italic:
                    font_style = "italic"
                color = self._extract_font_color(font)
                if color:
                    font_color = color
        except Exception as exc:
            logger.debug(
                "Paragraph property extraction failed: %s", exc
            )

        font_size_px = font_size_pt * _PT_TO_PX
        # Text-anchor based on alignment. python-pptx exposes
        # paragraph.alignment as a PP_ALIGN enum member. We
        # probe by name because importing the enum couples to
        # python-pptx internals.
        anchor, anchor_x = self._resolve_text_anchor(
            paragraph, x_px, width_px
        )

        # SVG <text> elements are positioned by their baseline,
        # not their top edge. Shift down by the font size for
        # a reasonable visual match.
        baseline_y = y_px + font_size_px

        escaped = self._escape_svg_text(text)
        svg = (
            f'<text x="{anchor_x:.2f}" y="{baseline_y:.2f}" '
            f'font-family="sans-serif" '
            f'font-size="{font_size_px:.2f}" '
            f'font-weight="{font_weight}" '
            f'font-style="{font_style}" '
            f'fill="{font_color}" '
            f'text-anchor="{anchor}">{escaped}</text>'
        )
        # Line height — 1.2x font size is a common typographic
        # default and matches PowerPoint's single-spacing
        # behaviour closely enough.
        return {"svg": svg, "height": font_size_px * 1.2}

    @staticmethod
    def _collect_paragraph_text(paragraph: Any) -> str:
        """Concatenate all run text in a paragraph.

        python-pptx's `paragraph.text` is equivalent, but using
        it would prevent us from choosing a specific run for
        property extraction. Collecting here keeps both the
        text and the run list available.
        """
        try:
            return "".join(run.text or "" for run in paragraph.runs)
        except Exception:
            return ""

    @staticmethod
    def _extract_font_color(font: Any) -> str | None:
        """Return a `#rrggbb` colour string from a font, or None.

        python-pptx's `font.color.rgb` is an `RGBColor` when an
        explicit colour is set. When the font uses a theme
        colour, `.rgb` raises `AttributeError`. We swallow
        that case and return None — callers use the default
        black.
        """
        try:
            color = font.color
            if color is None:
                return None
            rgb = color.rgb
            if rgb is None:
                return None
            # RGBColor's str is the hex without prefix.
            return "#" + str(rgb)
        except (AttributeError, ValueError):
            return None
        except Exception:
            return None

    @staticmethod
    def _resolve_text_anchor(
        paragraph: Any,
        left_px: float,
        width_px: float,
    ) -> tuple[str, float]:
        """Resolve paragraph alignment to an SVG text-anchor.

        Returns `(anchor_value, x_coordinate)`. SVG
        `text-anchor` is one of `start`, `middle`, `end`, with
        `start` being the default. The `x` coordinate shifts
        depending on the anchor so the text box fills from the
        correct side.

        Alignment detected by the name attribute on
        `paragraph.alignment` — avoids importing
        `PP_ALIGN` from python-pptx (which would couple us to
        library internals).
        """
        alignment = getattr(paragraph, "alignment", None)
        name = str(alignment).upper() if alignment is not None else ""
        if "CENTER" in name:
            return "middle", left_px + (width_px / 2)
        if "RIGHT" in name:
            return "end", left_px + width_px
        return "start", left_px

    def _render_table(
        self,
        table: Any,
        left_px: float,
        top_px: float,
        width_px: float,
        height_px: float,
    ) -> str:
        """Render a table as a grid of cells.

        Each cell is a `<rect>` border plus a `<text>` with the
        cell's string content centred. No merged-cell handling
        — python-pptx tracks merges via `.cell.is_merge_origin`
        but supporting them would require per-cell span
        tracking that doesn't fit Pass A4's scope.

        Rows and columns use uniform heights / widths derived
        from `table.rows` / `table.columns` where available,
        falling back to equal division of the shape's bounds.
        """
        try:
            rows = list(table.rows)
            columns = list(table.columns)
        except Exception as exc:
            logger.debug("Table enumeration failed: %s", exc)
            return ""

        if not rows or not columns:
            return ""

        # Per-row height and per-column width. python-pptx
        # exposes these via `.height` / `.width` in EMU.
        row_heights_px = [
            (row.height or 0) * _EMU_TO_PX for row in rows
        ]
        col_widths_px = [
            (col.width or 0) * _EMU_TO_PX for col in columns
        ]

        # Fall back to equal distribution if any dimension is
        # zero (happens with some templates).
        if not any(row_heights_px):
            row_heights_px = [
                height_px / len(rows) for _ in rows
            ]
        if not any(col_widths_px):
            col_widths_px = [
                width_px / len(columns) for _ in columns
            ]

        parts: list[str] = ["<g>"]
        y_cursor = top_px
        for row_idx, row in enumerate(rows):
            row_height = row_heights_px[row_idx] or (
                height_px / len(rows)
            )
            x_cursor = left_px
            for col_idx, cell in enumerate(row.cells):
                col_width = col_widths_px[col_idx] or (
                    width_px / len(columns)
                )
                # Cell border.
                parts.append(
                    f'<rect x="{x_cursor:.2f}" '
                    f'y="{y_cursor:.2f}" '
                    f'width="{col_width:.2f}" '
                    f'height="{row_height:.2f}" '
                    f'fill="none" stroke="#808080" '
                    f'stroke-width="1"/>'
                )
                # Cell text — centred vertically and
                # horizontally within the cell. Font uses
                # defaults; per-cell formatting would require
                # walking cell.text_frame.paragraphs, which
                # doesn't fit A4's scope.
                text = self._escape_svg_text(cell.text or "")
                if text:
                    text_x = x_cursor + (col_width / 2)
                    text_y = y_cursor + (row_height / 2) + (
                        _DEFAULT_FONT_SIZE_PT * _PT_TO_PX / 2
                    )
                    font_size_px = (
                        _DEFAULT_FONT_SIZE_PT * _PT_TO_PX
                    )
                    parts.append(
                        f'<text x="{text_x:.2f}" y="{text_y:.2f}" '
                        f'font-family="sans-serif" '
                        f'font-size="{font_size_px:.2f}" '
                        f'fill="{_DEFAULT_FONT_COLOR}" '
                        f'text-anchor="middle">{text}</text>'
                    )
                x_cursor += col_width
            y_cursor += row_height

        parts.append("</g>")
        return "\n".join(parts)

    @staticmethod
    def _escape_svg_text(text: str) -> str:
        """Escape characters that break an SVG text node.

        SVG is XML, so `<`, `>`, `&` need escaping. Quotes don't
        need escaping inside a text node (only inside attribute
        values), but we encode them anyway to be robust against
        accidental attribute-context use. Leading/trailing
        whitespace is stripped — PowerPoint often pads bullet
        text.
        """
        return (
            text.strip()
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&apos;")
        )

    # ------------------------------------------------------------------
    # Provenance header writing
    # ------------------------------------------------------------------

    @staticmethod
    def _build_provenance_header(
        source_name: str,
        source_hash: str,
        images: tuple[str, ...],
    ) -> str:
        """Build the `<!-- docuvert: ... -->` header string.

        Fields rendered in a stable order (source, sha256,
        images) so diff noise is minimal when a file is
        re-converted with only one field changing.

        The images field is omitted entirely when no images
        were extracted — keeps the header compact for text-only
        documents. The scan's parser tolerates both shapes
        (pinned by Pass A's `test_empty_images_list_is_empty_tuple`).
        """
        parts = [
            f"source={source_name}",
            f"sha256={source_hash}",
        ]
        if images:
            parts.append(f"images={','.join(images)}")
        body = " ".join(parts)
        return f"<!-- docuvert: {body} -->"

    # ------------------------------------------------------------------
    # Result-dict builders
    # ------------------------------------------------------------------

    @staticmethod
    def _fail(rel_path: str, message: str) -> dict[str, Any]:
        """Per-file error result."""
        return {
            "path": rel_path,
            "status": "error",
            "message": message,
        }

    @staticmethod
    def _skip(rel_path: str, message: str) -> dict[str, Any]:
        """Per-file skip result (not a failure, just not converted).

        Distinct from `error` so the frontend can render a
        different icon — skipped files may be retried later
        (over-size, extension deferred to a later pass) while
        errors typically indicate a real problem.
        """
        return {
            "path": rel_path,
            "status": "skipped",
            "message": message,
        }