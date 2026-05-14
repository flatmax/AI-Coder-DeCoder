"""python-pptx fallback pipeline rendering each slide as an SVG. Extracted from the original monolithic `doc_convert.py` during the package split."""

from __future__ import annotations

import base64
import logging
from pathlib import Path
from typing import Any

from .constants import (
    _DEFAULT_FONT_COLOR,
    _DEFAULT_FONT_SIZE_PT,
    _DEFAULT_SLIDE_HEIGHT_EMU,
    _DEFAULT_SLIDE_WIDTH_EMU,
    _EMU_TO_PX,
    _PT_TO_PX,
    _SLIDE_NUMBER_MIN_WIDTH,
)
from .provenance import build_provenance_header, hash_file, read_prior_images

logger = logging.getLogger(__name__)


class PptxPipeline:
    def __init__(self, fail, skip) -> None:
        self._fail = fail
        self._skip = skip

    def convert(
        self,
        root: Path,
        source_abs: Path,
        rel_path: str,
    ) -> dict[str, Any]:
        """Convert a .pptx file via the python-pptx fallback.

        Each slide renders as a standalone SVG containing the
        slide's text shapes, embedded images, and tables. The
        output markdown is an index file with per-slide headings
        and image references to the SVGs.

        Output layout:

            docs/presentation.pptx            ← source
            docs/presentation.md               ← index with links
            docs/presentation/
                01_slide.svg
                02_slide.svg
                ...

        Slide filenames use zero-padded numbering (`01_slide.svg`)
        so file listings sort correctly regardless of deck size.
        The assets subdirectory is always created for pptx (unlike
        markitdown docs where it's conditional on image presence)
        because every slide produces an SVG.

        Fails with a per-file error when python-pptx isn't
        installed. Unlike xlsx which falls back to markitdown,
        pptx has no markitdown-equivalent fallback in the current
        release — Pass A5 will add the LibreOffice + PyMuPDF
        primary path.
        """
        # Lazy import — python-pptx is optional in stripped-down
        # releases. A surface-level error with install hint
        # matches the markitdown-missing pattern.
        try:
            from pptx import Presentation
        except ImportError:
            return self._fail(
                rel_path,
                (
                    "python-pptx is not installed. Install with: "
                    "pip install 'ac-dc[docs]'"
                ),
            )

        output_abs = source_abs.with_suffix(".md")
        try:
            output_rel = output_abs.relative_to(root)
        except ValueError:
            return self._fail(
                rel_path,
                "Output path escapes repository root",
            )

        # Hash source for provenance — matches the markitdown and
        # xlsx paths so status classification works uniformly.
        try:
            source_hash = hash_file(source_abs)
        except OSError as exc:
            return self._fail(
                rel_path, f"Source hash failed: {exc}"
            )

        # Detect password-protected files before python-pptx
        # raises its opaque "Package not found" error. Encrypted
        # Office documents are wrapped in a CDFV2 (OLE compound)
        # container with the signature D0 CF 11 E0 A1 B1 1A E1,
        # whereas valid .pptx files are ZIP archives starting
        # with PK\x03\x04. The CDFV2 wrapper holds an encrypted
        # payload that neither LibreOffice (without --password)
        # nor python-pptx can read, so a clear early error
        # beats a misleading "Package not found".
        try:
            with open(source_abs, "rb") as fh:
                header = fh.read(8)
        except OSError as exc:
            return self._fail(
                rel_path, f"Failed to read source: {exc}"
            )
        if header.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"):
            return self._fail(
                rel_path,
                (
                    "File is password-protected (encrypted "
                    "CDFV2 container). Remove the password in "
                    "PowerPoint (File → Info → Protect "
                    "Presentation → Encrypt with Password → "
                    "clear) or decrypt with msoffcrypto-tool, "
                    "then retry."
                ),
            )

        # Open the presentation. python-pptx's failure modes on
        # invalid pptx vary — corrupt zip, missing core XML,
        # unsupported schema version. Broad catch to keep the
        # pipeline robust.
        try:
            presentation = Presentation(str(source_abs))
        except Exception as exc:
            return self._fail(
                rel_path,
                f"python-pptx failed to open: {exc}",
            )

        # Slide dimensions — fall back to the 4:3 default if
        # python-pptx reports None (rare but defensive).
        slide_width = presentation.slide_width or _DEFAULT_SLIDE_WIDTH_EMU
        slide_height = presentation.slide_height or _DEFAULT_SLIDE_HEIGHT_EMU
        svg_width_px = slide_width * _EMU_TO_PX
        svg_height_px = slide_height * _EMU_TO_PX

        slides = list(presentation.slides)
        if not slides:
            # Empty deck — produce an informative placeholder
            # output so the scan classifies it as `current`.
            markdown_text = "(empty presentation)\n"
            return self._write_pptx_output(
                output_abs=output_abs,
                output_rel=output_rel,
                source_abs=source_abs,
                source_hash=source_hash,
                markdown_text=markdown_text,
                rel_path=rel_path,
            )

        # Pick zero-padding width — larger of the default minimum
        # and whatever the deck actually needs. A 150-slide deck
        # pads to 3 digits.
        pad_width = max(
            _SLIDE_NUMBER_MIN_WIDTH,
            len(str(len(slides))),
        )

        # Per-source assets subdirectory. Unlike the markitdown
        # path, we always create it here — every slide produces
        # an SVG, so the dir is never empty after a successful
        # conversion.
        assets_dir = source_abs.with_suffix("")
        try:
            assets_dir.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            return self._fail(
                rel_path,
                f"Failed to create assets dir: {exc}",
            )

        # Render each slide. Failures on a single slide are
        # logged and that slide becomes a placeholder in the
        # index — better than aborting the whole deck.
        index_entries: list[str] = []
        saved_slides: list[str] = []
        for idx, slide in enumerate(slides, start=1):
            slide_name = f"{str(idx).zfill(pad_width)}_slide.svg"
            try:
                svg_text = self._render_pptx_slide(
                    slide, svg_width_px, svg_height_px
                )
            except Exception as exc:
                logger.debug(
                    "pptx slide %d of %s render failed: %s",
                    idx, rel_path, exc,
                )
                # Index entry still appears, but without an
                # image link — user sees which slide failed.
                index_entries.append(
                    f"## Slide {idx}\n\n*(rendering failed)*"
                )
                continue

            svg_path = assets_dir / slide_name
            try:
                svg_path.write_text(svg_text, encoding="utf-8")
            except OSError as exc:
                logger.debug(
                    "Failed to write %s: %s", svg_path, exc
                )
                index_entries.append(
                    f"## Slide {idx}\n\n*(write failed)*"
                )
                continue

            saved_slides.append(slide_name)
            # Index entry — heading + image reference. The image
            # ref uses the assets dir name as a relative prefix
            # so the markdown renders correctly from its sibling
            # location.
            rel_ref = f"{assets_dir.name}/{slide_name}"
            index_entries.append(
                f"## Slide {idx}\n\n![Slide {idx}]({rel_ref})"
            )

        markdown_text = "\n\n".join(index_entries) + "\n"

        # Orphan cleanup — re-conversion of a deck with fewer
        # slides than before leaves stale SVGs. Read the prior
        # provenance header (if any) and delete any images
        # listed there that we didn't produce this round.
        prior_images = read_prior_images(output_abs)
        if prior_images:
            saved_set = set(saved_slides)
            for orphan in prior_images:
                if orphan in saved_set:
                    continue
                orphan_path = assets_dir / orphan
                try:
                    orphan_path.unlink()
                except OSError as exc:
                    logger.debug(
                        "Failed to remove orphan slide %s: %s",
                        orphan_path, exc,
                    )

        return self._write_pptx_output(
            output_abs=output_abs,
            output_rel=output_rel,
            source_abs=source_abs,
            source_hash=source_hash,
            markdown_text=markdown_text,
            rel_path=rel_path,
            images=tuple(saved_slides),
        )

    def _write_pptx_output(
        self,
        output_abs: Path,
        output_rel: Path,
        source_abs: Path,
        source_hash: str,
        markdown_text: str,
        rel_path: str,
        images: tuple[str, ...] = (),
    ) -> dict[str, Any]:
        """Write the pptx index markdown with provenance header.

        Shared between the normal path and the empty-deck case.
        Builds the provenance header, prepends it to the markdown
        body, writes atomically, returns a per-file result dict.
        """
        provenance_line = build_provenance_header(
            source_name=source_abs.name,
            source_hash=source_hash,
            images=images,
        )
        final_content = (
            provenance_line + "\n\n" + markdown_text.lstrip("\n")
        )

        try:
            output_abs.parent.mkdir(parents=True, exist_ok=True)
            output_abs.write_text(final_content, encoding="utf-8")
        except OSError as exc:
            return self._fail(
                rel_path,
                f"Failed to write output: {exc}",
            )

        return {
            "path": rel_path,
            "status": "ok",
            "output_path": str(output_rel).replace("\\", "/"),
            "images": list(images),
        }

    def _render_pptx_slide(
        self,
        slide: Any,
        width_px: float,
        height_px: float,
    ) -> str:
        """Render one slide as a complete SVG document.

        Walks the slide's shapes and emits SVG elements for each.
        Supported shape types:

        - Text frames (including within placeholders like title,
          body, content): `<text>` with font-size, weight, fill,
          and alignment derived from the first run's properties
        - Pictures: `<image>` with the raster content embedded
          as a base64 data URI, preserving position and size
        - Tables: `<rect>` borders + `<text>` cell content,
          rendered as a simple grid (no merged cells)

        Unsupported shape types (charts, SmartArt, OLE objects,
        groups) are skipped with a debug log. A slide with only
        unsupported shapes produces an empty SVG frame — still
        valid output, the user can see the slide exists but the
        content isn't surfaced.
        """
        # SVG header — viewBox lets the rendered image scale to
        # any container size while preserving aspect ratio.
        parts: list[str] = [
            (
                f'<svg xmlns="http://www.w3.org/2000/svg" '
                f'xmlns:xlink="http://www.w3.org/1999/xlink" '
                f'width="{int(width_px)}" height="{int(height_px)}" '
                f'viewBox="0 0 {int(width_px)} {int(height_px)}">'
            ),
            # White background — matches default PowerPoint
            # slide colour. Themes with non-white backgrounds
            # are a future enhancement.
            (
                f'<rect width="{int(width_px)}" '
                f'height="{int(height_px)}" fill="#ffffff"/>'
            ),
        ]

        for shape in slide.shapes:
            try:
                element = self._render_pptx_shape(shape)
            except Exception as exc:
                logger.debug(
                    "Skipping shape %r: %s",
                    getattr(shape, "shape_type", "?"), exc,
                )
                continue
            if element:
                parts.append(element)

        parts.append("</svg>")
        return "\n".join(parts)

    def _render_pptx_shape(self, shape: Any) -> str:
        """Dispatch a single shape to its renderer.

        Returns an SVG element string or empty string for
        skipped shapes. Never raises — all render errors are
        caught at the call site and logged.
        """
        # Dimensions in EMU, converted to pixels.
        left_px = (shape.left or 0) * _EMU_TO_PX
        top_px = (shape.top or 0) * _EMU_TO_PX
        width_px = (shape.width or 0) * _EMU_TO_PX
        height_px = (shape.height or 0) * _EMU_TO_PX

        # Picture — render as <image> with embedded base64 payload.
        if self._is_picture(shape):
            return self._render_picture(
                shape, left_px, top_px, width_px, height_px
            )

        # Table — render as grid of <rect> + <text>.
        if getattr(shape, "has_table", False):
            return self._render_table(
                shape.table, left_px, top_px, width_px, height_px
            )

        # Text frame (including placeholder shapes).
        if getattr(shape, "has_text_frame", False):
            return self._render_text_frame(
                shape.text_frame, left_px, top_px, width_px, height_px
            )

        # Unsupported shape type (chart, SmartArt, group, etc.).
        # Return empty string — caller skips.
        return ""

    @staticmethod
    def _is_picture(shape: Any) -> bool:
        """Return True when the shape is a raster image.

        python-pptx's shape_type constants live under
        `MSO_SHAPE_TYPE`; the PICTURE enum value is 13. Rather
        than importing the enum (which couples us to a specific
        python-pptx version), we probe the attributes the
        Picture class exposes: `.image.blob` returns the raw
        bytes.
        """
        image = getattr(shape, "image", None)
        if image is None:
            return False
        # Image object exposes .blob. A non-picture shape with
        # an "image" attr (rare) wouldn't have this.
        return hasattr(image, "blob")

    def _render_picture(
        self,
        shape: Any,
        left_px: float,
        top_px: float,
        width_px: float,
        height_px: float,
    ) -> str:
        """Render a picture shape as an SVG <image> element.

        The raster payload is base64-encoded and inlined via a
        data URI. Keeps slide layout self-contained — one file
        per slide, no external image refs.

        Returns empty string on any read error — a damaged
        image shouldn't abort the slide.
        """
        try:
            image = shape.image
            blob = image.blob
            content_type = image.content_type or "image/png"
        except Exception as exc:
            logger.debug("Failed to read picture blob: %s", exc)
            return ""

        encoded = base64.b64encode(blob).decode("ascii")
        data_uri = f"data:{content_type};base64,{encoded}"
        return (
            f'<image x="{left_px:.2f}" y="{top_px:.2f}" '
            f'width="{width_px:.2f}" height="{height_px:.2f}" '
            f'xlink:href="{data_uri}"/>'
        )

    def _render_text_frame(
        self,
        text_frame: Any,
        left_px: float,
        top_px: float,
        width_px: float,
        height_px: float,
    ) -> str:
        """Render a text frame as a <g> of <text> elements.

        Each paragraph in the frame becomes a `<text>` line
        positioned vertically by cumulative line height. Font
        properties (size, weight, colour, alignment) come from
        the first run of each paragraph — the spec's
        simplification for Pass A4 is "one style per
        paragraph", not per-run.

        Returns empty string for frames with no rendered text —
        avoids emitting degenerate `<g></g>` wrappers.
        """
        # Collect rendered paragraphs.
        lines: list[str] = []
        cursor_y = top_px
        for paragraph in text_frame.paragraphs:
            line = self._render_paragraph(
                paragraph, left_px, cursor_y, width_px
            )
            if line:
                lines.append(line["svg"])
                cursor_y += line["height"]

        if not lines:
            return ""
        return "<g>\n" + "\n".join(lines) + "\n</g>"

    def _render_paragraph(
        self,
        paragraph: Any,
        x_px: float,
        y_px: float,
        width_px: float,
    ) -> dict[str, Any] | None:
        """Render a paragraph as a single `<text>` element.

        Returns a dict with `svg` (the element string) and
        `height` (vertical space consumed) for cumulative
        layout. None when the paragraph has no text — caller
        skips.

        Font properties are extracted from the first non-empty
        run; subsequent runs inherit these. Richer handling
        (per-run formatting) is deferred to Pass A5's richer
        pipeline.
        """
        text = self._collect_paragraph_text(paragraph)
        if not text:
            return None

        # Extract properties from the first run — python-pptx
        # exposes them under `paragraph.runs[N].font`.
        font_size_pt = _DEFAULT_FONT_SIZE_PT
        font_weight = "normal"
        font_style = "normal"
        font_color = _DEFAULT_FONT_COLOR
        try:
            first_run = next(
                (r for r in paragraph.runs if r.text),
                None,
            )
            if first_run is not None:
                font = first_run.font
                if font.size is not None:
                    # python-pptx returns Emu objects for sizes;
                    # .pt gives points directly.
                    font_size_pt = int(font.size.pt)
                if font.bold:
                    font_weight = "bold"
                if font.italic:
                    font_style = "italic"
                color = self._extract_font_color(font)
                if color:
                    font_color = color
        except Exception as exc:
            logger.debug(
                "Paragraph property extraction failed: %s", exc
            )

        font_size_px = font_size_pt * _PT_TO_PX
        # Text-anchor based on alignment. python-pptx exposes
        # paragraph.alignment as a PP_ALIGN enum member. We
        # probe by name because importing the enum couples to
        # python-pptx internals.
        anchor, anchor_x = self._resolve_text_anchor(
            paragraph, x_px, width_px
        )

        # SVG <text> elements are positioned by their baseline,
        # not their top edge. Shift down by the font size for
        # a reasonable visual match.
        baseline_y = y_px + font_size_px

        escaped = self._escape_svg_text(text)
        svg = (
            f'<text x="{anchor_x:.2f}" y="{baseline_y:.2f}" '
            f'font-family="sans-serif" '
            f'font-size="{font_size_px:.2f}" '
            f'font-weight="{font_weight}" '
            f'font-style="{font_style}" '
            f'fill="{font_color}" '
            f'text-anchor="{anchor}">{escaped}</text>'
        )
        # Line height — 1.2x font size is a common typographic
        # default and matches PowerPoint's single-spacing
        # behaviour closely enough.
        return {"svg": svg, "height": font_size_px * 1.2}

    @staticmethod
    def _collect_paragraph_text(paragraph: Any) -> str:
        """Concatenate all run text in a paragraph.

        python-pptx's `paragraph.text` is equivalent, but using
        it would prevent us from choosing a specific run for
        property extraction. Collecting here keeps both the
        text and the run list available.
        """
        try:
            return "".join(run.text or "" for run in paragraph.runs)
        except Exception:
            return ""

    @staticmethod
    def _extract_font_color(font: Any) -> str | None:
        """Return a `#rrggbb` colour string from a font, or None.

        python-pptx's `font.color.rgb` is an `RGBColor` when an
        explicit colour is set. When the font uses a theme
        colour, `.rgb` raises `AttributeError`. We swallow
        that case and return None — callers use the default
        black.
        """
        try:
            color = font.color
            if color is None:
                return None
            rgb = color.rgb
            if rgb is None:
                return None
            # RGBColor's str is the hex without prefix.
            return "#" + str(rgb)
        except (AttributeError, ValueError):
            return None
        except Exception:
            return None

    @staticmethod
    def _resolve_text_anchor(
        paragraph: Any,
        left_px: float,
        width_px: float,
    ) -> tuple[str, float]:
        """Resolve paragraph alignment to an SVG text-anchor.

        Returns `(anchor_value, x_coordinate)`. SVG
        `text-anchor` is one of `start`, `middle`, `end`, with
        `start` being the default. The `x` coordinate shifts
        depending on the anchor so the text box fills from the
        correct side.

        Alignment detected by the name attribute on
        `paragraph.alignment` — avoids importing
        `PP_ALIGN` from python-pptx (which would couple us to
        library internals).
        """
        alignment = getattr(paragraph, "alignment", None)
        name = str(alignment).upper() if alignment is not None else ""
        if "CENTER" in name:
            return "middle", left_px + (width_px / 2)
        if "RIGHT" in name:
            return "end", left_px + width_px
        return "start", left_px

    def _render_table(
        self,
        table: Any,
        left_px: float,
        top_px: float,
        width_px: float,
        height_px: float,
    ) -> str:
        """Render a table as a grid of cells.

        Each cell is a `<rect>` border plus a `<text>` with the
        cell's string content centred. No merged-cell handling
        — python-pptx tracks merges via `.cell.is_merge_origin`
        but supporting them would require per-cell span
        tracking that doesn't fit Pass A4's scope.

        Rows and columns use uniform heights / widths derived
        from `table.rows` / `table.columns` where available,
        falling back to equal division of the shape's bounds.
        """
        try:
            rows = list(table.rows)
            columns = list(table.columns)
        except Exception as exc:
            logger.debug("Table enumeration failed: %s", exc)
            return ""

        if not rows or not columns:
            return ""

        # Per-row height and per-column width. python-pptx
        # exposes these via `.height` / `.width` in EMU.
        row_heights_px = [
            (row.height or 0) * _EMU_TO_PX for row in rows
        ]
        col_widths_px = [
            (col.width or 0) * _EMU_TO_PX for col in columns
        ]

        # Fall back to equal distribution if any dimension is
        # zero (happens with some templates).
        if not any(row_heights_px):
            row_heights_px = [
                height_px / len(rows) for _ in rows
            ]
        if not any(col_widths_px):
            col_widths_px = [
                width_px / len(columns) for _ in columns
            ]

        parts: list[str] = ["<g>"]
        y_cursor = top_px
        for row_idx, row in enumerate(rows):
            row_height = row_heights_px[row_idx] or (
                height_px / len(rows)
            )
            x_cursor = left_px
            for col_idx, cell in enumerate(row.cells):
                col_width = col_widths_px[col_idx] or (
                    width_px / len(columns)
                )
                # Cell border.
                parts.append(
                    f'<rect x="{x_cursor:.2f}" '
                    f'y="{y_cursor:.2f}" '
                    f'width="{col_width:.2f}" '
                    f'height="{row_height:.2f}" '
                    f'fill="none" stroke="#808080" '
                    f'stroke-width="1"/>'
                )
                # Cell text — centred vertically and
                # horizontally within the cell. Font uses
                # defaults; per-cell formatting would require
                # walking cell.text_frame.paragraphs, which
                # doesn't fit A4's scope.
                text = self._escape_svg_text(cell.text or "")
                if text:
                    text_x = x_cursor + (col_width / 2)
                    text_y = y_cursor + (row_height / 2) + (
                        _DEFAULT_FONT_SIZE_PT * _PT_TO_PX / 2
                    )
                    font_size_px = (
                        _DEFAULT_FONT_SIZE_PT * _PT_TO_PX
                    )
                    parts.append(
                        f'<text x="{text_x:.2f}" y="{text_y:.2f}" '
                        f'font-family="sans-serif" '
                        f'font-size="{font_size_px:.2f}" '
                        f'fill="{_DEFAULT_FONT_COLOR}" '
                        f'text-anchor="middle">{text}</text>'
                    )
                x_cursor += col_width
            y_cursor += row_height

        parts.append("</g>")
        return "\n".join(parts)

    @staticmethod
    def _escape_svg_text(text: str) -> str:
        """Escape characters that break an SVG text node.

        SVG is XML, so `<`, `>`, `&` need escaping. Quotes don't
        need escaping inside a text node (only inside attribute
        values), but we encode them anyway to be robust against
        accidental attribute-context use. Leading/trailing
        whitespace is stripped — PowerPoint often pads bullet
        text.
        """
        return (
            text.strip()
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&apos;")
        )