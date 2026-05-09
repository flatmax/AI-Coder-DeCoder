"""Shared helpers and stubs for doc_convert test suite.

Pure-function helpers and stub classes live here so they can be
imported from individual test modules without going through the
pytest fixture system. Fixtures live in ``conftest.py``.
"""

from __future__ import annotations

import base64
import hashlib
from pathlib import Path
from typing import Any

import pytest


# ---------------------------------------------------------------------------
# Stub collab (matches test_collab_restrictions.py + test_settings.py)
# ---------------------------------------------------------------------------


class _StubCollab:
    def __init__(self, is_localhost: bool = True) -> None:
        self._is_localhost = is_localhost
        self.call_count = 0

    def is_caller_localhost(self) -> bool:
        self.call_count += 1
        return self._is_localhost


class _RaisingCollab:
    def is_caller_localhost(self) -> bool:
        raise RuntimeError("collab check failed")


# ---------------------------------------------------------------------------
# Generic helpers
# ---------------------------------------------------------------------------


def _assert_restricted(result: Any) -> None:
    assert isinstance(result, dict)
    assert result.get("error") == "restricted"
    reason = result.get("reason")
    assert isinstance(reason, str) and reason


def _write_source(scan_root: Path, rel_path: str, content: bytes) -> Path:
    """Helper — write a binary source file at a given relative path."""
    path = scan_root / rel_path
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(content)
    return path


def _write_output(
    scan_root: Path,
    rel_path: str,
    provenance: str | None,
    body: str = "converted content\n",
) -> Path:
    """Helper — write a markdown output with an optional provenance header."""
    path = scan_root / rel_path
    path.parent.mkdir(parents=True, exist_ok=True)
    lines: list[str] = []
    if provenance is not None:
        lines.append(f"<!-- docuvert: {provenance} -->")
        lines.append("")
    lines.append(body)
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def _sha256_of(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


# ---------------------------------------------------------------------------
# Markitdown stub (used by the markitdown-path tests)
# ---------------------------------------------------------------------------


class _FakeMarkItDownResult:
    """Stand-in for ``markitdown.MarkItDown().convert().text_content``."""

    def __init__(self, text: str) -> None:
        self.text_content = text


class _FakeMarkItDown:
    """Minimal MarkItDown stub.

    Instances return `_FakeMarkItDownResult` from convert(). Output
    text comes from a module-level mapping keyed by source path, so
    each test can set its own expected conversion output without
    fighting global state.
    """

    # Map of absolute source path (as string) → text to return.
    # Tests populate this before calling convert_files.
    outputs: dict[str, str] = {}

    # If set, convert() raises this exception instead of returning.
    raise_on_convert: Exception | None = None

    def convert(self, source: str) -> _FakeMarkItDownResult:
        if _FakeMarkItDown.raise_on_convert is not None:
            raise _FakeMarkItDown.raise_on_convert
        text = _FakeMarkItDown.outputs.get(source, "")
        return _FakeMarkItDownResult(text)


# ---------------------------------------------------------------------------
# Image / docx helpers
# ---------------------------------------------------------------------------


def _make_png_bytes() -> bytes:
    """Return a minimal valid 1x1 PNG payload.

    Constructed at runtime with real CRC checksums via `zlib.crc32`.
    A hand-encoded hex string is error-prone — one wrong nibble in
    a length or CRC field makes the file fail strict decoders like
    PyMuPDF with "premature end of data". Building fresh means
    every chunk is structurally valid by construction.

    The image is a 1x1 red pixel (grayscale-simplest-form would
    save ~10 bytes but isn't worth the added obscurity).
    """
    import struct
    import zlib

    def _chunk(chunk_type: bytes, data: bytes) -> bytes:
        """Build a PNG chunk: length + type + data + CRC."""
        return (
            struct.pack(">I", len(data))
            + chunk_type
            + data
            + struct.pack(">I", zlib.crc32(chunk_type + data))
        )

    signature = b"\x89PNG\r\n\x1a\n"
    # IHDR — 13 bytes: width, height, bit depth, colour type,
    # compression, filter, interlace.
    ihdr = struct.pack(
        ">IIBBBBB",
        1,    # width
        1,    # height
        8,    # bit depth
        2,    # colour type: truecolour (RGB)
        0,    # compression method
        0,    # filter method
        0,    # interlace method
    )
    # IDAT — zlib-compressed raw scanlines. For a 1x1 RGB image,
    # the raw data is 1 filter byte + 3 RGB bytes = 4 bytes.
    raw = b"\x00\xff\x00\x00"  # filter=None, red=255, green=0, blue=0
    idat = zlib.compress(raw, level=9)
    # IEND — no data.
    return (
        signature
        + _chunk(b"IHDR", ihdr)
        + _chunk(b"IDAT", idat)
        + _chunk(b"IEND", b"")
    )


def _make_docx_zip(path: Path, media: dict[str, bytes] | None = None) -> None:
    """Create a minimal .docx (zip) at path with optional media files.

    Only populates the `word/media/` subdirectory — we don't need
    a fully-valid document.xml for image-extraction tests. The zip
    reader only looks under `word/media/`.
    """
    import zipfile
    with zipfile.ZipFile(path, "w") as zf:
        # A minimal [Content_Types].xml so the zip is valid.
        zf.writestr(
            "[Content_Types].xml",
            b"<?xml version='1.0'?><Types/>",
        )
        for name, data in (media or {}).items():
            zf.writestr(f"word/media/{name}", data)


def _make_data_uri(payload: bytes, mime: str = "image/png") -> str:
    """Build a data:image URI from raw bytes."""
    encoded = base64.b64encode(payload).decode("ascii")
    return f"data:{mime};base64,{encoded}"


# ---------------------------------------------------------------------------
# Optional-dependency skip helpers
# ---------------------------------------------------------------------------


def _require_openpyxl():
    """Skip the test if openpyxl isn't installed."""
    try:
        import openpyxl  # noqa: F401
    except ImportError:
        pytest.skip("openpyxl not installed")


def _require_pptx():
    """Skip the test if python-pptx isn't installed."""
    try:
        import pptx  # noqa: F401
    except ImportError:
        pytest.skip("python-pptx not installed")


def _require_pymupdf():
    """Skip the test if PyMuPDF isn't installed."""
    try:
        import fitz  # noqa: F401
    except ImportError:
        pytest.skip("PyMuPDF not installed")


# ---------------------------------------------------------------------------
# xlsx builder
# ---------------------------------------------------------------------------


def _write_xlsx(
    path: Path,
    sheets: dict[str, list[list[tuple[str, str | None]]]],
) -> None:
    """Write an xlsx file with the given sheets.

    Each sheet is a list of rows; each row is a list of
    `(value, fill_hex)` tuples where `fill_hex` is None for no
    fill or a 6-char hex string (without leading #).
    """
    _require_openpyxl()
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill

    wb = Workbook()
    # Remove the default sheet — we'll add named sheets below.
    default = wb.active
    wb.remove(default)

    for sheet_name, rows in sheets.items():
        ws = wb.create_sheet(title=sheet_name)
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, (value, fill_hex) in enumerate(row, start=1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                if fill_hex:
                    cell.fill = PatternFill(
                        start_color=fill_hex,
                        end_color=fill_hex,
                        fill_type="solid",
                    )
    wb.save(str(path))


# ---------------------------------------------------------------------------
# pptx builders
# ---------------------------------------------------------------------------


def _make_pptx_with_title(
    path: Path,
    title: str,
    body: str = "",
) -> None:
    """Create a minimal pptx with a title slide.

    Uses python-pptx's default layout which provides title +
    body placeholders. `body` defaults to empty.
    """
    _require_pptx()
    from pptx import Presentation

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if body and len(slide.placeholders) > 1:
        slide.placeholders[1].text = body
    prs.save(str(path))


def _make_pptx_with_n_slides(path: Path, n: int) -> None:
    """Create a pptx with n title slides, each numbered."""
    _require_pptx()
    from pptx import Presentation

    prs = Presentation()
    for i in range(n):
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {i + 1}"
    prs.save(str(path))


def _make_pptx_with_image(
    path: Path,
    image_bytes: bytes,
) -> None:
    """Create a pptx with one image-containing slide."""
    _require_pptx()
    import io
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Blank layout — index 6 in the default template.
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(
        io.BytesIO(image_bytes),
        Inches(1), Inches(1),
        width=Inches(3), height=Inches(2),
    )
    prs.save(str(path))


def _make_pptx_with_table(path: Path) -> None:
    """Create a pptx with one table-containing slide."""
    _require_pptx()
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    rows, cols = 2, 3
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(1), Inches(1),
        Inches(6), Inches(2),
    )
    table = table_shape.table
    table.cell(0, 0).text = "header1"
    table.cell(0, 1).text = "header2"
    table.cell(0, 2).text = "header3"
    table.cell(1, 0).text = "a"
    table.cell(1, 1).text = "b"
    table.cell(1, 2).text = "c"
    prs.save(str(path))


# ---------------------------------------------------------------------------
# pdf builders
# ---------------------------------------------------------------------------


def _make_pdf_with_text(path: Path, pages: list[str]) -> None:
    """Create a minimal PDF with one text block per page.

    Each page gets a single text block positioned near the
    top-left. No images or drawings — pure text.
    """
    _require_pymupdf()
    import fitz

    doc = fitz.open()
    for text in pages:
        page = doc.new_page(width=612, height=792)  # US Letter
        # Insert text at a visible position.
        page.insert_text((72, 72), text, fontsize=12)
    doc.save(str(path))
    doc.close()


def _make_pdf_with_image(path: Path, image_bytes: bytes) -> None:
    """Create a PDF with one image-containing page."""
    _require_pymupdf()
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=612, height=792)
    # Insert a raster image at a fixed position.
    rect = fitz.Rect(72, 72, 372, 372)
    page.insert_image(rect, stream=image_bytes)
    doc.save(str(path))
    doc.close()


def _make_pdf_with_text_and_image(
    path: Path, text: str, image_bytes: bytes
) -> None:
    """Create a PDF with both text and an image on one page."""
    _require_pymupdf()
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=612, height=792)
    page.insert_text((72, 72), text, fontsize=12)
    rect = fitz.Rect(72, 400, 372, 700)
    page.insert_image(rect, stream=image_bytes)
    doc.save(str(path))
    doc.close()


def _make_empty_pdf(path: Path) -> None:
    """Create a PDF with zero pages.

    PyMuPDF allows empty documents but the save invokes the
    writer which requires at least one page — so we use a
    workaround by writing a minimal valid PDF directly.
    """
    # Minimal valid PDF with 0 pages. The xref is the trick —
    # PDF spec allows empty page trees.
    path.write_bytes(
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[]/Count 0>>endobj\n"
        b"xref\n0 3\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000052 00000 n \n"
        b"trailer<</Size 3/Root 1 0 R>>\n"
        b"startxref\n94\n%%EOF\n"
    )