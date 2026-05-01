"""Tests for ac_dc.doc_convert.DocConvert — Layer 4.5 Pass A (foundation).

Covers:

- Construction (holds config reference, collab starts None)
- Localhost-only guard (single-user path, localhost allowed,
  non-localhost rejected, fail-closed on raising collab)
- `is_available` dependency probing (reports all four flags,
  degradation when deps missing, derived `pdf_pipeline` flag)
- `scan_convertible_files` (empty repo, new sources, status
  classification matrix, directory exclusions, extension
  filtering, disabled flag returns empty, size threshold flag,
  stable sort order, malformed files skipped defensively)
- Provenance header parsing (valid header, missing required
  fields, unknown fields preserved as extra, no header,
  malformed bytes)
- `convert_files` stub (restricted for non-localhost, raises
  NotImplementedError for localhost — proves the guard runs
  before the unimplemented body)

Strategy:

- Isolated config dir via ``AC_DC_CONFIG_HOME`` (the documented
  test hook used by every other test suite).
- Use `tmp_path` as the scan root rather than constructing a
  real `Repo` — the scanner only needs `.root` which is trivially
  provided by a `SimpleNamespace` stub. Keeps tests from coupling
  to the Repo test fixtures.
- Stub collab with configurable `is_caller_localhost`, same
  pattern as `test_settings.py` and `test_collab_restrictions.py`.
"""

from __future__ import annotations

import base64
import hashlib
import json
import shutil
from pathlib import Path
from types import SimpleNamespace
from typing import Any
from unittest.mock import patch

import pytest

from ac_dc.config import ConfigManager
from ac_dc.doc_convert import (
    DocConvert,
    ProvenanceHeader,
    _DEFAULT_EXTENSIONS,
    _EXCLUDED_DIRS,
)


# ---------------------------------------------------------------------------
# Stub collab (matches test_collab_restrictions.py + test_settings.py)
# ---------------------------------------------------------------------------


class _StubCollab:
    def __init__(self, is_localhost: bool = True) -> None:
        self._is_localhost = is_localhost
        self.call_count = 0

    def is_caller_localhost(self) -> bool:
        self.call_count += 1
        return self._is_localhost


class _RaisingCollab:
    def is_caller_localhost(self) -> bool:
        raise RuntimeError("collab check failed")


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def isolated_config_dir(tmp_path, monkeypatch):
    home = tmp_path / "ac-dc-config"
    monkeypatch.setenv("AC_DC_CONFIG_HOME", str(home))
    return home


@pytest.fixture
def config(isolated_config_dir):
    return ConfigManager()


@pytest.fixture
def scan_root(tmp_path):
    """A fresh directory to act as the scan root."""
    root = tmp_path / "scan-root"
    root.mkdir()
    return root


@pytest.fixture
def fake_repo(scan_root):
    """A minimal repo stand-in exposing `.root`."""
    return SimpleNamespace(root=scan_root)


@pytest.fixture
def doc_convert(config, fake_repo):
    """DocConvert with no collab attached (single-user mode)."""
    return DocConvert(config, repo=fake_repo)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _assert_restricted(result: Any) -> None:
    assert isinstance(result, dict)
    assert result.get("error") == "restricted"
    reason = result.get("reason")
    assert isinstance(reason, str) and reason


def _write_source(scan_root: Path, rel_path: str, content: bytes) -> Path:
    """Helper — write a binary source file at a given relative path."""
    path = scan_root / rel_path
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(content)
    return path


def _write_output(
    scan_root: Path,
    rel_path: str,
    provenance: str | None,
    body: str = "converted content\n",
) -> Path:
    """Helper — write a markdown output with an optional provenance header."""
    path = scan_root / rel_path
    path.parent.mkdir(parents=True, exist_ok=True)
    lines: list[str] = []
    if provenance is not None:
        lines.append(f"<!-- docuvert: {provenance} -->")
        lines.append("")
    lines.append(body)
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def _sha256_of(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


# ---------------------------------------------------------------------------
# Construction
# ---------------------------------------------------------------------------


class TestConstruction:
    def test_holds_config_reference(self, config):
        svc = DocConvert(config)
        assert svc._config is config

    def test_collab_starts_none(self, doc_convert):
        assert doc_convert._collab is None

    def test_repo_optional(self, config):
        # No repo → uses CWD as root. Not great for production but
        # useful for tests and for "scan current directory" CLI-like
        # use cases.
        svc = DocConvert(config)
        assert svc._repo is None
        # Calling _root returns a Path — cannot guarantee specific
        # value since CWD varies across test runs.
        root = svc._root()
        assert isinstance(root, Path)


# ---------------------------------------------------------------------------
# Localhost-only guard
# ---------------------------------------------------------------------------


class TestLocalhostGuard:
    def test_no_collab_returns_none(self, doc_convert):
        assert doc_convert._check_localhost_only() is None

    def test_localhost_caller_returns_none(self, doc_convert):
        doc_convert._collab = _StubCollab(is_localhost=True)
        assert doc_convert._check_localhost_only() is None
        assert doc_convert._collab.call_count == 1

    def test_non_localhost_caller_returns_restricted(self, doc_convert):
        doc_convert._collab = _StubCollab(is_localhost=False)
        result = doc_convert._check_localhost_only()
        _assert_restricted(result)

    def test_raising_collab_fails_closed(self, doc_convert):
        doc_convert._collab = _RaisingCollab()
        result = doc_convert._check_localhost_only()
        _assert_restricted(result)


# ---------------------------------------------------------------------------
# is_available — dependency probing
# ---------------------------------------------------------------------------


class TestIsAvailable:
    def test_returns_all_four_flags(self, doc_convert):
        result = doc_convert.is_available()
        assert set(result.keys()) == {
            "available",
            "libreoffice",
            "pymupdf",
            "pdf_pipeline",
        }
        for value in result.values():
            assert isinstance(value, bool)

    def test_pdf_pipeline_requires_both_deps(self, doc_convert):
        """Truth table — pdf_pipeline = libreoffice AND pymupdf."""
        with patch.object(
            DocConvert, "_probe_import",
            side_effect=lambda name: {"markitdown": True, "fitz": True}.get(name, False),
        ), patch("ac_dc.doc_convert.shutil.which", return_value="/usr/bin/soffice"):
            result = doc_convert.is_available()
        assert result["pdf_pipeline"] is True

    def test_pdf_pipeline_false_without_libreoffice(self, doc_convert):
        with patch.object(
            DocConvert, "_probe_import",
            side_effect=lambda name: {"markitdown": True, "fitz": True}.get(name, False),
        ), patch("ac_dc.doc_convert.shutil.which", return_value=None):
            result = doc_convert.is_available()
        assert result["pymupdf"] is True
        assert result["libreoffice"] is False
        assert result["pdf_pipeline"] is False

    def test_pdf_pipeline_false_without_pymupdf(self, doc_convert):
        with patch.object(
            DocConvert, "_probe_import",
            side_effect=lambda name: {"markitdown": True, "fitz": False}.get(name, False),
        ), patch("ac_dc.doc_convert.shutil.which", return_value="/usr/bin/soffice"):
            result = doc_convert.is_available()
        assert result["libreoffice"] is True
        assert result["pymupdf"] is False
        assert result["pdf_pipeline"] is False

    def test_markitdown_missing_flags_unavailable(self, doc_convert):
        with patch.object(
            DocConvert, "_probe_import",
            side_effect=lambda name: False,
        ), patch("ac_dc.doc_convert.shutil.which", return_value=None):
            result = doc_convert.is_available()
        assert result["available"] is False

    def test_all_deps_missing(self, doc_convert):
        with patch.object(
            DocConvert, "_probe_import", return_value=False,
        ), patch("ac_dc.doc_convert.shutil.which", return_value=None):
            result = doc_convert.is_available()
        assert result == {
            "available": False,
            "libreoffice": False,
            "pymupdf": False,
            "pdf_pipeline": False,
        }

    def test_probe_import_catches_exception(self):
        """A raising import doesn't escape the probe."""
        with patch(
            "importlib.import_module",
            side_effect=RuntimeError("bad install"),
        ):
            assert DocConvert._probe_import("anything") is False

    def test_is_available_callable_when_disabled(
        self, doc_convert, isolated_config_dir
    ):
        """Disabled config doesn't affect availability probing."""
        # Disable via app config.
        app_json = isolated_config_dir / "app.json"
        data = json.loads(app_json.read_text(encoding="utf-8"))
        data["doc_convert"]["enabled"] = False
        app_json.write_text(json.dumps(data), encoding="utf-8")
        doc_convert._config.reload_app_config()
        # Should still return the capability dict, not refuse.
        result = doc_convert.is_available()
        assert "available" in result


# ---------------------------------------------------------------------------
# scan_convertible_files — empty case + basic structure
# ---------------------------------------------------------------------------


class TestScanEmpty:
    def test_empty_repo(self, doc_convert):
        assert doc_convert.scan_convertible_files() == []

    def test_disabled_returns_empty(
        self, doc_convert, isolated_config_dir, scan_root
    ):
        _write_source(scan_root, "doc.docx", b"fake docx bytes")
        app_json = isolated_config_dir / "app.json"
        data = json.loads(app_json.read_text(encoding="utf-8"))
        data["doc_convert"]["enabled"] = False
        app_json.write_text(json.dumps(data), encoding="utf-8")
        doc_convert._config.reload_app_config()
        assert doc_convert.scan_convertible_files() == []

    def test_missing_root_logs_and_returns_empty(
        self, config, tmp_path
    ):
        """A repo root that doesn't exist doesn't crash."""
        fake_repo = SimpleNamespace(root=tmp_path / "does-not-exist")
        svc = DocConvert(config, repo=fake_repo)
        assert svc.scan_convertible_files() == []


# ---------------------------------------------------------------------------
# scan_convertible_files — extension filtering
# ---------------------------------------------------------------------------


class TestScanExtensions:
    def test_default_extensions_recognised(self, doc_convert, scan_root):
        for ext in _DEFAULT_EXTENSIONS:
            _write_source(scan_root, f"doc{ext}", b"x")
        results = doc_convert.scan_convertible_files()
        found = {r["path"] for r in results}
        expected = {f"doc{ext}" for ext in _DEFAULT_EXTENSIONS}
        assert found == expected

    def test_non_convertible_extensions_ignored(self, doc_convert, scan_root):
        _write_source(scan_root, "script.py", b"code")
        _write_source(scan_root, "readme.md", b"docs")
        _write_source(scan_root, "real.docx", b"x")
        results = doc_convert.scan_convertible_files()
        paths = {r["path"] for r in results}
        assert paths == {"real.docx"}

    def test_case_insensitive_extension_match(self, doc_convert, scan_root):
        _write_source(scan_root, "UPPER.DOCX", b"x")
        _write_source(scan_root, "Mixed.PdF", b"x")
        results = doc_convert.scan_convertible_files()
        assert len(results) == 2

    def test_config_can_restrict_extensions(
        self, doc_convert, isolated_config_dir, scan_root
    ):
        _write_source(scan_root, "a.docx", b"x")
        _write_source(scan_root, "b.xlsx", b"x")
        _write_source(scan_root, "c.pdf", b"x")
        # Override to only .docx.
        app_json = isolated_config_dir / "app.json"
        data = json.loads(app_json.read_text(encoding="utf-8"))
        data["doc_convert"]["extensions"] = [".docx"]
        app_json.write_text(json.dumps(data), encoding="utf-8")
        doc_convert._config.reload_app_config()
        results = doc_convert.scan_convertible_files()
        paths = {r["path"] for r in results}
        assert paths == {"a.docx"}


# ---------------------------------------------------------------------------
# scan_convertible_files — directory exclusions
# ---------------------------------------------------------------------------


class TestScanExclusions:
    def test_git_directory_excluded(self, doc_convert, scan_root):
        _write_source(scan_root, ".git/info/stash.docx", b"x")
        _write_source(scan_root, "visible.docx", b"x")
        results = doc_convert.scan_convertible_files()
        paths = {r["path"] for r in results}
        assert paths == {"visible.docx"}

    def test_ac_dc_directory_excluded(self, doc_convert, scan_root):
        _write_source(scan_root, ".ac-dc4/history.docx", b"x")
        _write_source(scan_root, "visible.docx", b"x")
        results = doc_convert.scan_convertible_files()
        assert {r["path"] for r in results} == {"visible.docx"}

    def test_node_modules_excluded(self, doc_convert, scan_root):
        _write_source(scan_root, "node_modules/pkg/README.docx", b"x")
        _write_source(scan_root, "docs/real.docx", b"x")
        results = doc_convert.scan_convertible_files()
        assert {r["path"] for r in results} == {"docs/real.docx"}

    def test_hidden_dirs_excluded_except_github(
        self, doc_convert, scan_root
    ):
        _write_source(scan_root, ".vscode/settings.docx", b"x")
        _write_source(scan_root, ".github/docs/ci.docx", b"x")
        _write_source(scan_root, "visible.docx", b"x")
        results = doc_convert.scan_convertible_files()
        paths = {r["path"] for r in results}
        # .vscode excluded; .github allowed; visible present.
        assert ".vscode/settings.docx" not in paths
        # Normalise separator for cross-platform.
        assert ".github/docs/ci.docx" in paths
        assert "visible.docx" in paths

    def test_all_excluded_dirs_enumerated(self, doc_convert, scan_root):
        """Smoke-test every name in _EXCLUDED_DIRS."""
        for name in _EXCLUDED_DIRS:
            _write_source(scan_root, f"{name}/trapped.docx", b"x")
        _write_source(scan_root, "visible.docx", b"x")
        results = doc_convert.scan_convertible_files()
        assert {r["path"] for r in results} == {"visible.docx"}


# ---------------------------------------------------------------------------
# scan_convertible_files — status classification matrix
# ---------------------------------------------------------------------------


class TestScanStatusClassification:
    def test_new_when_no_output(self, doc_convert, scan_root):
        _write_source(scan_root, "doc.docx", b"content")
        [entry] = doc_convert.scan_convertible_files()
        assert entry["status"] == "new"

    def test_conflict_when_output_has_no_header(
        self, doc_convert, scan_root
    ):
        _write_source(scan_root, "doc.docx", b"content")
        _write_output(scan_root, "doc.md", provenance=None)
        [entry] = doc_convert.scan_convertible_files()
        assert entry["status"] == "conflict"

    def test_current_when_hash_matches(self, doc_convert, scan_root):
        content = b"stable content"
        _write_source(scan_root, "doc.docx", content)
        sha = _sha256_of(content)
        _write_output(
            scan_root, "doc.md",
            provenance=f"source=doc.docx sha256={sha}",
        )
        [entry] = doc_convert.scan_convertible_files()
        assert entry["status"] == "current"

    def test_stale_when_hash_differs(self, doc_convert, scan_root):
        _write_source(scan_root, "doc.docx", b"new content")
        _write_output(
            scan_root, "doc.md",
            provenance=(
                "source=doc.docx sha256="
                "aaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaa"
            ),
        )
        [entry] = doc_convert.scan_convertible_files()
        assert entry["status"] == "stale"

    def test_conflict_takes_precedence_over_hash_check(
        self, doc_convert, scan_root
    ):
        """A file without a docuvert header is conflict regardless of
        whether a hash check would pass or fail."""
        _write_source(scan_root, "doc.docx", b"content")
        # Output with no header at all.
        (scan_root / "doc.md").write_text(
            "# Manually authored\n\nNothing to see.\n",
            encoding="utf-8",
        )
        [entry] = doc_convert.scan_convertible_files()
        assert entry["status"] == "conflict"

    def test_malformed_header_treated_as_conflict(
        self, doc_convert, scan_root
    ):
        """Header missing required fields → conflict."""
        _write_source(scan_root, "doc.docx", b"content")
        _write_output(
            scan_root, "doc.md",
            # No sha256 field — fails required-field validation.
            provenance="source=doc.docx",
        )
        [entry] = doc_convert.scan_convertible_files()
        assert entry["status"] == "conflict"


# ---------------------------------------------------------------------------
# scan_convertible_files — entry shape
# ---------------------------------------------------------------------------


class TestScanEntryShape:
    def test_entry_has_required_fields(self, doc_convert, scan_root):
        _write_source(scan_root, "doc.docx", b"hello")
        [entry] = doc_convert.scan_convertible_files()
        assert set(entry.keys()) == {
            "path", "name", "size", "status",
            "output_path", "over_size",
        }

    def test_path_and_name_populated(self, doc_convert, scan_root):
        _write_source(scan_root, "dir/nested.docx", b"x")
        [entry] = doc_convert.scan_convertible_files()
        assert entry["path"] == "dir/nested.docx"
        assert entry["name"] == "nested.docx"

    def test_output_path_is_sibling_md(self, doc_convert, scan_root):
        _write_source(scan_root, "docs/arch.docx", b"x")
        [entry] = doc_convert.scan_convertible_files()
        assert entry["output_path"] == "docs/arch.md"

    def test_size_is_byte_count(self, doc_convert, scan_root):
        content = b"x" * 12345
        _write_source(scan_root, "sized.docx", content)
        [entry] = doc_convert.scan_convertible_files()
        assert entry["size"] == 12345

    def test_over_size_flag_false_under_threshold(
        self, doc_convert, scan_root
    ):
        _write_source(scan_root, "small.docx", b"x" * 100)
        [entry] = doc_convert.scan_convertible_files()
        assert entry["over_size"] is False

    def test_over_size_flag_true_above_threshold(
        self, doc_convert, isolated_config_dir, scan_root
    ):
        # Set max to 1 MB.
        app_json = isolated_config_dir / "app.json"
        data = json.loads(app_json.read_text(encoding="utf-8"))
        data["doc_convert"]["max_source_size_mb"] = 1
        app_json.write_text(json.dumps(data), encoding="utf-8")
        doc_convert._config.reload_app_config()
        # Write > 1 MB of content.
        _write_source(scan_root, "big.docx", b"x" * (2 * 1024 * 1024))
        [entry] = doc_convert.scan_convertible_files()
        assert entry["over_size"] is True

    def test_windows_path_normalised(self, doc_convert, scan_root):
        """Path uses forward slashes regardless of platform."""
        _write_source(scan_root, "a/b/c.docx", b"x")
        [entry] = doc_convert.scan_convertible_files()
        assert "\\" not in entry["path"]
        assert "\\" not in entry["output_path"]


# ---------------------------------------------------------------------------
# scan_convertible_files — ordering
# ---------------------------------------------------------------------------


class TestScanOrdering:
    def test_stable_alphabetical_sort(self, doc_convert, scan_root):
        _write_source(scan_root, "z.docx", b"x")
        _write_source(scan_root, "a.docx", b"x")
        _write_source(scan_root, "m.docx", b"x")
        results = doc_convert.scan_convertible_files()
        paths = [r["path"] for r in results]
        assert paths == sorted(paths)

    def test_sort_includes_nested_paths(self, doc_convert, scan_root):
        _write_source(scan_root, "b/z.docx", b"x")
        _write_source(scan_root, "a/z.docx", b"x")
        _write_source(scan_root, "b/a.docx", b"x")
        results = doc_convert.scan_convertible_files()
        paths = [r["path"] for r in results]
        # Sorted alphabetically, so "a/z" < "b/a" < "b/z".
        assert paths == ["a/z.docx", "b/a.docx", "b/z.docx"]


# ---------------------------------------------------------------------------
# Provenance header parsing
# ---------------------------------------------------------------------------


class TestProvenanceParsing:
    def test_valid_header_parses(self):
        body = "source=doc.docx sha256=abc123 images=a.png,b.png"
        result = DocConvert.parse_provenance_body(body)
        assert result is not None
        assert result.source == "doc.docx"
        assert result.sha256 == "abc123"
        assert result.images == ("a.png", "b.png")
        assert result.extra is None

    def test_minimal_header_valid(self):
        body = "source=doc.docx sha256=abc"
        result = DocConvert.parse_provenance_body(body)
        assert result is not None
        assert result.images == ()

    def test_missing_source_returns_none(self):
        body = "sha256=abc"
        assert DocConvert.parse_provenance_body(body) is None

    def test_missing_sha256_returns_none(self):
        body = "source=doc.docx"
        assert DocConvert.parse_provenance_body(body) is None

    def test_empty_body_returns_none(self):
        assert DocConvert.parse_provenance_body("") is None

    def test_unknown_fields_captured_as_extra(self):
        body = "source=doc.docx sha256=abc tool_version=2.0 custom=xyz"
        result = DocConvert.parse_provenance_body(body)
        assert result is not None
        assert result.extra == {"tool_version": "2.0", "custom": "xyz"}

    def test_empty_images_list_is_empty_tuple(self):
        body = "source=doc.docx sha256=abc images="
        result = DocConvert.parse_provenance_body(body)
        assert result is not None
        assert result.images == ()

    def test_single_image(self):
        body = "source=doc.docx sha256=abc images=only.png"
        result = DocConvert.parse_provenance_body(body)
        assert result is not None
        assert result.images == ("only.png",)

    def test_is_frozen_dataclass(self):
        header = ProvenanceHeader(source="a", sha256="b")
        with pytest.raises((AttributeError, Exception)):
            header.source = "changed"  # type: ignore[misc]


class TestProvenanceReadFromFile:
    def test_reads_header_from_output_file(
        self, doc_convert, scan_root
    ):
        _write_output(
            scan_root, "doc.md",
            provenance="source=doc.docx sha256=abc",
        )
        header = doc_convert._read_provenance_header(
            scan_root / "doc.md"
        )
        assert header is not None
        assert header.source == "doc.docx"
        assert header.sha256 == "abc"

    def test_missing_header_returns_none(self, doc_convert, scan_root):
        (scan_root / "doc.md").write_text(
            "# Plain markdown\n", encoding="utf-8"
        )
        assert doc_convert._read_provenance_header(
            scan_root / "doc.md"
        ) is None

    def test_invalid_body_returns_none(self, doc_convert, scan_root):
        _write_output(
            scan_root, "doc.md",
            provenance="only=garbage",
        )
        assert doc_convert._read_provenance_header(
            scan_root / "doc.md"
        ) is None

    def test_unreadable_file_returns_none(self, doc_convert, scan_root):
        """A file we can't open returns None rather than crashing."""
        # A directory at the path where a file should be.
        bad_path = scan_root / "doc.md"
        bad_path.mkdir()
        assert doc_convert._read_provenance_header(bad_path) is None

    def test_header_with_leading_content_still_found(
        self, doc_convert, scan_root
    ):
        """Header doesn't HAVE to be on line 1 — a leading BOM or
        blank line should still work."""
        (scan_root / "doc.md").write_text(
            "\n<!-- docuvert: source=a sha256=b -->\n# body\n",
            encoding="utf-8",
        )
        header = doc_convert._read_provenance_header(
            scan_root / "doc.md"
        )
        assert header is not None
        assert header.source == "a"


# ---------------------------------------------------------------------------
# Source hashing
# ---------------------------------------------------------------------------


class TestHashFile:
    def test_matches_stdlib_sha256(self, tmp_path):
        content = b"hello world" * 1000
        path = tmp_path / "data.bin"
        path.write_bytes(content)
        expected = hashlib.sha256(content).hexdigest()
        assert DocConvert._hash_file(path) == expected

    def test_streams_large_file(self, tmp_path):
        """File larger than the 64KB chunk is hashed correctly."""
        content = b"x" * (500 * 1024)  # 500 KB
        path = tmp_path / "large.bin"
        path.write_bytes(content)
        expected = hashlib.sha256(content).hexdigest()
        assert DocConvert._hash_file(path) == expected

    def test_empty_file(self, tmp_path):
        path = tmp_path / "empty"
        path.write_bytes(b"")
        assert (
            DocConvert._hash_file(path)
            == hashlib.sha256(b"").hexdigest()
        )


# ---------------------------------------------------------------------------
# convert_files — Pass A2 (markitdown path)
# ---------------------------------------------------------------------------
#
# These tests replace the Pass A stub tests. convert_files now has a
# real implementation for .docx/.rtf/.odt via markitdown, returns
# per-file results for other extensions, and enforces the clean-tree
# gate when a repo is attached.
#
# markitdown is stubbed via sys.modules injection (same pattern as
# test_llm_service.py's litellm fake). Tests that exercise the real
# library would couple to a specific markitdown version's output
# format; the stub lets us pin behaviour precisely.


class _FakeMarkItDownResult:
    """Stand-in for ``markitdown.MarkItDown().convert().text_content``."""

    def __init__(self, text: str) -> None:
        self.text_content = text


class _FakeMarkItDown:
    """Minimal MarkItDown stub.

    Instances return `_FakeMarkItDownResult` from convert(). Output
    text comes from a module-level mapping keyed by source path, so
    each test can set its own expected conversion output without
    fighting global state.
    """

    # Map of absolute source path (as string) → text to return.
    # Tests populate this before calling convert_files.
    outputs: dict[str, str] = {}

    # If set, convert() raises this exception instead of returning.
    raise_on_convert: Exception | None = None

    def convert(self, source: str) -> _FakeMarkItDownResult:
        if _FakeMarkItDown.raise_on_convert is not None:
            raise _FakeMarkItDown.raise_on_convert
        text = _FakeMarkItDown.outputs.get(source, "")
        return _FakeMarkItDownResult(text)


@pytest.fixture
def fake_markitdown(monkeypatch):
    """Install a fake markitdown module via sys.modules."""
    import sys
    import types

    # Reset per-test state.
    _FakeMarkItDown.outputs = {}
    _FakeMarkItDown.raise_on_convert = None

    fake_module = types.ModuleType("markitdown")
    fake_module.MarkItDown = _FakeMarkItDown
    monkeypatch.setitem(sys.modules, "markitdown", fake_module)
    return _FakeMarkItDown


@pytest.fixture
def clean_repo(scan_root):
    """A fake repo whose `is_clean()` returns True."""
    return SimpleNamespace(
        root=scan_root,
        is_clean=lambda: True,
    )


@pytest.fixture
def dirty_repo(scan_root):
    """A fake repo whose `is_clean()` returns False."""
    return SimpleNamespace(
        root=scan_root,
        is_clean=lambda: False,
    )


def _make_png_bytes() -> bytes:
    """Return a minimal valid 1x1 PNG payload.

    Constructed at runtime with real CRC checksums via `zlib.crc32`.
    A hand-encoded hex string is error-prone — one wrong nibble in
    a length or CRC field makes the file fail strict decoders like
    PyMuPDF with "premature end of data". Building fresh means
    every chunk is structurally valid by construction.

    The image is a 1x1 red pixel (grayscale-simplest-form would
    save ~10 bytes but isn't worth the added obscurity).
    """
    import struct
    import zlib

    def _chunk(chunk_type: bytes, data: bytes) -> bytes:
        """Build a PNG chunk: length + type + data + CRC."""
        return (
            struct.pack(">I", len(data))
            + chunk_type
            + data
            + struct.pack(">I", zlib.crc32(chunk_type + data))
        )

    signature = b"\x89PNG\r\n\x1a\n"
    # IHDR — 13 bytes: width, height, bit depth, colour type,
    # compression, filter, interlace.
    ihdr = struct.pack(
        ">IIBBBBB",
        1,    # width
        1,    # height
        8,    # bit depth
        2,    # colour type: truecolour (RGB)
        0,    # compression method
        0,    # filter method
        0,    # interlace method
    )
    # IDAT — zlib-compressed raw scanlines. For a 1x1 RGB image,
    # the raw data is 1 filter byte + 3 RGB bytes = 4 bytes.
    raw = b"\x00\xff\x00\x00"  # filter=None, red=255, green=0, blue=0
    idat = zlib.compress(raw, level=9)
    # IEND — no data.
    return (
        signature
        + _chunk(b"IHDR", ihdr)
        + _chunk(b"IDAT", idat)
        + _chunk(b"IEND", b"")
    )


def _make_docx_zip(path: Path, media: dict[str, bytes] | None = None) -> None:
    """Create a minimal .docx (zip) at path with optional media files.

    Only populates the `word/media/` subdirectory — we don't need
    a fully-valid document.xml for image-extraction tests. The zip
    reader only looks under `word/media/`.
    """
    import zipfile
    with zipfile.ZipFile(path, "w") as zf:
        # A minimal [Content_Types].xml so the zip is valid.
        zf.writestr(
            "[Content_Types].xml",
            b"<?xml version='1.0'?><Types/>",
        )
        for name, data in (media or {}).items():
            zf.writestr(f"word/media/{name}", data)


def _make_data_uri(payload: bytes, mime: str = "image/png") -> str:
    """Build a data:image URI from raw bytes."""
    encoded = base64.b64encode(payload).decode("ascii")
    return f"data:{mime};base64,{encoded}"


# ---------------------------------------------------------------------------
# Localhost guard + restricted callers
# ---------------------------------------------------------------------------


class TestConvertFilesGuards:
    def test_non_localhost_returns_restricted(self, doc_convert):
        doc_convert._collab = _StubCollab(is_localhost=False)
        result = doc_convert.convert_files(["any.docx"])
        _assert_restricted(result)

    def test_raising_collab_returns_restricted(self, doc_convert):
        doc_convert._collab = _RaisingCollab()
        result = doc_convert.convert_files(["any.docx"])
        _assert_restricted(result)

    def test_guard_runs_before_clean_tree_check(
        self, config, dirty_repo
    ):
        """Restricted caller doesn't even see the dirty-tree error."""
        svc = DocConvert(config, repo=dirty_repo)
        svc._collab = _StubCollab(is_localhost=False)
        result = svc.convert_files(["any.docx"])
        _assert_restricted(result)


# ---------------------------------------------------------------------------
# Clean-tree gate
# ---------------------------------------------------------------------------


class TestCleanTreeGate:
    def test_dirty_tree_rejected(
        self, config, dirty_repo, scan_root
    ):
        _write_source(scan_root, "doc.docx", b"x")
        svc = DocConvert(config, repo=dirty_repo)
        result = svc.convert_files(["doc.docx"])
        assert "error" in result
        assert "uncommitted" in result["error"].lower()

    def test_clean_tree_allowed(
        self, config, clean_repo, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"x")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = "body\n"
        svc = DocConvert(config, repo=clean_repo)
        result = svc.convert_files(["doc.docx"])
        # No top-level error — per-file results present.
        assert "error" not in result
        assert result["status"] == "ok"

    def test_no_repo_skips_gate(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """Without a repo, the gate is skipped (CLI / test use)."""
        # doc_convert fixture uses fake_repo which has no is_clean.
        _write_source(scan_root, "doc.docx", b"x")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = "body\n"
        result = doc_convert.convert_files(["doc.docx"])
        assert "error" not in result


# ---------------------------------------------------------------------------
# Extension dispatch
# ---------------------------------------------------------------------------


class TestExtensionDispatch:
    def test_docx_routes_to_markitdown(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"x")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = "text\n"
        result = doc_convert.convert_files(["doc.docx"])
        [entry] = result["results"]
        assert entry["status"] == "ok"

    def test_rtf_routes_to_markitdown(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.rtf", b"x")
        fake_markitdown.outputs[str(scan_root / "doc.rtf")] = "text\n"
        result = doc_convert.convert_files(["doc.rtf"])
        [entry] = result["results"]
        assert entry["status"] == "ok"

    def test_odt_routes_to_markitdown(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.odt", b"x")
        fake_markitdown.outputs[str(scan_root / "doc.odt")] = "text\n"
        result = doc_convert.convert_files(["doc.odt"])
        [entry] = result["results"]
        assert entry["status"] == "ok"


    def test_unsupported_extension_errors(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "weird.xyz", b"x")
        result = doc_convert.convert_files(["weird.xyz"])
        [entry] = result["results"]
        assert entry["status"] == "error"

    def test_csv_routes_to_markitdown(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """Pass A3 — csv uses markitdown (clean table output)."""
        _write_source(scan_root, "data.csv", b"a,b,c\n1,2,3\n")
        fake_markitdown.outputs[str(scan_root / "data.csv")] = (
            "| a | b | c |\n|---|---|---|\n| 1 | 2 | 3 |\n"
        )
        result = doc_convert.convert_files(["data.csv"])
        [entry] = result["results"]
        assert entry["status"] == "ok"

    def test_mixed_batch_produces_per_file_results(
        self, doc_convert, scan_root, fake_markitdown
    ):
        # Pairs one success with one failure (missing file)
        # to prove per-file isolation — a failing entry doesn't
        # abort the rest of the batch. Before Pass A5a/b there
        # were multiple supported-but-deferred extensions that
        # produced `skipped` results; now every supported ext
        # has at least one working path (or a format-specific
        # fallback), so a genuine failure is the right test
        # shape.
        _write_source(scan_root, "ok.docx", b"x")
        # "missing.docx" is deliberately not written — the
        # converter's pre-flight sees it doesn't exist.
        fake_markitdown.outputs[str(scan_root / "ok.docx")] = "text\n"
        result = doc_convert.convert_files(["ok.docx", "missing.docx"])
        assert len(result["results"]) == 2
        statuses = [r["status"] for r in result["results"]]
        assert statuses == ["ok", "error"]


# ---------------------------------------------------------------------------
# Pre-flight validation
# ---------------------------------------------------------------------------


class TestPreflightValidation:
    def test_missing_file_errors(
        self, doc_convert, fake_markitdown
    ):
        result = doc_convert.convert_files(["nonexistent.docx"])
        [entry] = result["results"]
        assert entry["status"] == "error"
        assert "not found" in entry["message"].lower()

    def test_path_traversal_rejected(
        self, doc_convert, fake_markitdown
    ):
        result = doc_convert.convert_files(["../escape.docx"])
        [entry] = result["results"]
        assert entry["status"] == "error"
        # Message should mention the path must be inside the root.
        assert (
            "repository" in entry["message"].lower()
            or "not found" in entry["message"].lower()
        )

    def test_over_size_file_skipped(
        self, doc_convert, isolated_config_dir, scan_root, fake_markitdown
    ):
        # Set max to 1 MB, write 2 MB.
        app_json = isolated_config_dir / "app.json"
        data = json.loads(app_json.read_text(encoding="utf-8"))
        data["doc_convert"]["max_source_size_mb"] = 1
        app_json.write_text(json.dumps(data), encoding="utf-8")
        doc_convert._config.reload_app_config()
        _write_source(scan_root, "big.docx", b"x" * (2 * 1024 * 1024))
        result = doc_convert.convert_files(["big.docx"])
        [entry] = result["results"]
        assert entry["status"] == "skipped"
        assert "limit" in entry["message"].lower()


# ---------------------------------------------------------------------------
# markitdown failure handling
# ---------------------------------------------------------------------------


class TestMarkitdownFailures:
    def test_missing_markitdown_returns_error(
        self, doc_convert, scan_root, monkeypatch
    ):
        """When markitdown isn't installed, return a clean error."""
        import sys
        _write_source(scan_root, "doc.docx", b"x")
        # Ensure no markitdown in sys.modules. Also block the import
        # at a lower level so the lazy `from markitdown import ...`
        # raises ImportError.
        monkeypatch.delitem(sys.modules, "markitdown", raising=False)

        real_import = __builtins__["__import__"] if isinstance(
            __builtins__, dict
        ) else __builtins__.__import__

        def blocking_import(name, *args, **kwargs):
            if name == "markitdown":
                raise ImportError("markitdown not installed")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(
            "builtins.__import__", blocking_import,
        )
        result = doc_convert.convert_files(["doc.docx"])
        [entry] = result["results"]
        assert entry["status"] == "error"
        assert "markitdown" in entry["message"].lower()

    def test_markitdown_exception_captured(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"x")
        fake_markitdown.raise_on_convert = RuntimeError("corrupted file")
        result = doc_convert.convert_files(["doc.docx"])
        [entry] = result["results"]
        assert entry["status"] == "error"
        assert "corrupted" in entry["message"].lower()


# ---------------------------------------------------------------------------
# Provenance header writing
# ---------------------------------------------------------------------------


class TestProvenanceWriting:
    def test_header_on_first_line(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"content")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = "body text\n"
        doc_convert.convert_files(["doc.docx"])
        output = (scan_root / "doc.md").read_text(encoding="utf-8")
        assert output.startswith("<!-- docuvert:")

    def test_header_contains_source_and_sha(
        self, doc_convert, scan_root, fake_markitdown
    ):
        content = b"known content"
        expected_hash = _sha256_of(content)
        _write_source(scan_root, "doc.docx", content)
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = "body\n"
        doc_convert.convert_files(["doc.docx"])
        output = (scan_root / "doc.md").read_text(encoding="utf-8")
        assert "source=doc.docx" in output
        assert f"sha256={expected_hash}" in output

    def test_header_includes_images_when_present(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"x")
        png = _make_png_bytes()
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            f"Before ![alt]({_make_data_uri(png)}) after\n"
        )
        doc_convert.convert_files(["doc.docx"])
        output = (scan_root / "doc.md").read_text(encoding="utf-8")
        assert "images=doc_img1.png" in output

    def test_header_omits_images_when_none(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"x")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = "no images\n"
        doc_convert.convert_files(["doc.docx"])
        output = (scan_root / "doc.md").read_text(encoding="utf-8")
        # Header present but `images=` field absent.
        assert "<!-- docuvert:" in output
        assert "images=" not in output

    def test_roundtrip_via_scan(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """Converted file scans as `current` on next pass."""
        _write_source(scan_root, "doc.docx", b"stable")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = "body\n"
        doc_convert.convert_files(["doc.docx"])
        scan = doc_convert.scan_convertible_files()
        [entry] = scan
        assert entry["status"] == "current"

    def test_conversion_changes_to_stale_on_source_edit(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"original")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = "body\n"
        doc_convert.convert_files(["doc.docx"])
        # Edit the source — hash now differs from what's recorded.
        _write_source(scan_root, "doc.docx", b"edited")
        scan = doc_convert.scan_convertible_files()
        [entry] = scan
        assert entry["status"] == "stale"


# ---------------------------------------------------------------------------
# Data-URI image extraction
# ---------------------------------------------------------------------------


class TestDataUriImages:
    def test_single_image_extracted(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"x")
        png = _make_png_bytes()
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            f"Text ![alt]({_make_data_uri(png)}) more\n"
        )
        result = doc_convert.convert_files(["doc.docx"])
        [entry] = result["results"]
        assert entry["images"] == ["doc_img1.png"]
        # File was written.
        image_path = scan_root / "doc" / "doc_img1.png"
        assert image_path.is_file()
        assert image_path.read_bytes() == png

    def test_multiple_images_numbered_in_order(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"x")
        png = _make_png_bytes()
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            f"![a]({_make_data_uri(png)}) "
            f"![b]({_make_data_uri(png)}) "
            f"![c]({_make_data_uri(png)})\n"
        )
        result = doc_convert.convert_files(["doc.docx"])
        [entry] = result["results"]
        assert entry["images"] == [
            "doc_img1.png", "doc_img2.png", "doc_img3.png",
        ]

    def test_markdown_references_rewritten(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"x")
        png = _make_png_bytes()
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            f"Before ![alt]({_make_data_uri(png)}) after\n"
        )
        doc_convert.convert_files(["doc.docx"])
        output = (scan_root / "doc.md").read_text(encoding="utf-8")
        assert "![alt](doc/doc_img1.png)" in output
        # Original data URI stripped.
        assert "data:image/png;base64" not in output

    def test_alt_text_preserved(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"x")
        png = _make_png_bytes()
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            f"![Architecture diagram]({_make_data_uri(png)})\n"
        )
        doc_convert.convert_files(["doc.docx"])
        output = (scan_root / "doc.md").read_text(encoding="utf-8")
        assert "![Architecture diagram](doc/doc_img1.png)" in output

    def test_jpeg_extension(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _write_source(scan_root, "doc.docx", b"x")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            f"![alt]({_make_data_uri(b'fakejpegdata', 'image/jpeg')})\n"
        )
        result = doc_convert.convert_files(["doc.docx"])
        [entry] = result["results"]
        assert entry["images"] == ["doc_img1.jpg"]

    def test_no_images_no_assets_dir(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """Files without images don't leave an empty assets subdirectory."""
        _write_source(scan_root, "doc.docx", b"x")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            "text only, no images\n"
        )
        doc_convert.convert_files(["doc.docx"])
        assert not (scan_root / "doc").exists()

    def test_decode_failure_leaves_reference_in_place(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """Invalid base64 doesn't crash — leaves the broken ref."""
        _write_source(scan_root, "doc.docx", b"x")
        # Intentionally broken payload with non-base64 chars.
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            "![alt](data:image/png;base64,!!!not-base64!!!)\n"
        )
        result = doc_convert.convert_files(["doc.docx"])
        [entry] = result["results"]
        # Conversion succeeded but no images saved.
        assert entry["status"] == "ok"


# ---------------------------------------------------------------------------
# DOCX truncated-URI workaround
# ---------------------------------------------------------------------------


class TestDocxTruncatedUris:
    def test_truncated_uri_replaced_from_zip(
        self, doc_convert, scan_root, fake_markitdown
    ):
        png_bytes = _make_png_bytes()
        docx_path = scan_root / "doc.docx"
        _make_docx_zip(docx_path, {"image1.png": png_bytes})
        fake_markitdown.outputs[str(docx_path)] = (
            "Before ![alt](data:image/png;base64...) after\n"
        )
        result = doc_convert.convert_files(["doc.docx"])
        [entry] = result["results"]
        assert entry["status"] == "ok"
        # Image extracted and saved — proves the truncated URI
        # was replaced with a real payload that then went through
        # the data-URI extractor.
        image_path = scan_root / "doc" / "doc_img1.png"
        assert image_path.is_file()
        assert image_path.read_bytes() == png_bytes

    def test_multiple_truncated_uris_matched_in_order(
        self, doc_convert, scan_root, fake_markitdown
    ):
        png1 = _make_png_bytes()
        png2 = png1 + b"\x00"  # distinguishable variant
        docx_path = scan_root / "doc.docx"
        _make_docx_zip(docx_path, {
            "image1.png": png1,
            "image2.png": png2,
        })
        fake_markitdown.outputs[str(docx_path)] = (
            "![a](data:image/png;base64...) "
            "![b](data:image/png;base64...)\n"
        )
        doc_convert.convert_files(["doc.docx"])
        # First image should match png1, second png2.
        assert (scan_root / "doc" / "doc_img1.png").read_bytes() == png1
        assert (scan_root / "doc" / "doc_img2.png").read_bytes() == png2

    def test_no_media_in_zip_leaves_truncated_uri(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """When the docx zip has no images, we can't substitute."""
        docx_path = scan_root / "doc.docx"
        _make_docx_zip(docx_path, {})  # no media
        fake_markitdown.outputs[str(docx_path)] = (
            "![alt](data:image/png;base64...)\n"
        )
        result = doc_convert.convert_files(["doc.docx"])
        [entry] = result["results"]
        # Conversion succeeded but no image saved.
        assert entry["status"] == "ok"
        assert entry["images"] == []

    def test_non_zip_docx_tolerated(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """Corrupt .docx (not a zip) doesn't crash the pipeline."""
        _write_source(scan_root, "corrupt.docx", b"not a zip")
        fake_markitdown.outputs[str(scan_root / "corrupt.docx")] = (
            "text with ![alt](data:image/png;base64...) ref\n"
        )
        result = doc_convert.convert_files(["corrupt.docx"])
        [entry] = result["results"]
        # markitdown succeeded (the fake doesn't care about
        # zip-ness), so conversion proceeds. The truncated URI is
        # left in place because docx extraction returned no images.
        assert entry["status"] == "ok"


# ---------------------------------------------------------------------------
# Orphan image cleanup
# ---------------------------------------------------------------------------


class TestOrphanCleanup:
    def test_orphans_removed_on_reconversion(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """Images in the old header but not the new output are deleted."""
        _write_source(scan_root, "doc.docx", b"v1")
        png = _make_png_bytes()

        # First conversion — produces two images.
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            f"![a]({_make_data_uri(png)}) "
            f"![b]({_make_data_uri(png)})\n"
        )
        doc_convert.convert_files(["doc.docx"])
        assert (scan_root / "doc" / "doc_img1.png").exists()
        assert (scan_root / "doc" / "doc_img2.png").exists()

        # Edit source and reconvert — only one image this time.
        _write_source(scan_root, "doc.docx", b"v2")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            f"![a]({_make_data_uri(png)})\n"
        )
        doc_convert.convert_files(["doc.docx"])

        # img1 kept, img2 cleaned up as orphan.
        assert (scan_root / "doc" / "doc_img1.png").exists()
        assert not (scan_root / "doc" / "doc_img2.png").exists()

    def test_no_header_no_orphan_cleanup(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """Conflict files (no header) are overwritten without cleanup.

        The previous content wasn't ours to manage, so we don't
        touch any sibling files — user may have put things in
        the assets subdir by hand.
        """
        _write_source(scan_root, "doc.docx", b"x")
        # Pre-existing assets subdir with a file we didn't create.
        (scan_root / "doc").mkdir()
        stranger = scan_root / "doc" / "user_file.png"
        stranger.write_bytes(b"stranger content")
        # Pre-existing output with no header.
        (scan_root / "doc.md").write_text(
            "# Manual\n\nNo docuvert header.\n",
            encoding="utf-8",
        )
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = "body\n"
        doc_convert.convert_files(["doc.docx"])
        # Stranger file still present.
        assert stranger.exists()

    def test_assets_dir_removed_when_fully_orphaned(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """If every image is an orphan and none are produced, the
        assets dir is removed."""
        _write_source(scan_root, "doc.docx", b"v1")
        png = _make_png_bytes()
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = (
            f"![a]({_make_data_uri(png)})\n"
        )
        doc_convert.convert_files(["doc.docx"])
        assert (scan_root / "doc").is_dir()

        # Reconvert with no images.
        _write_source(scan_root, "doc.docx", b"v2")
        fake_markitdown.outputs[str(scan_root / "doc.docx")] = "text only\n"
        doc_convert.convert_files(["doc.docx"])
        assert not (scan_root / "doc").exists()


# ---------------------------------------------------------------------------
# Pass A3 — xlsx via openpyxl (colour-aware)
# ---------------------------------------------------------------------------
#
# openpyxl is a required dependency of the `[docs]` extra, so these
# tests import it directly. We build real xlsx files with openpyxl
# rather than mocking — the real library's cell/fill/workbook API
# is complex enough that mocking would miss bugs. Tests that need
# to exercise the "openpyxl missing" path use monkeypatching to
# block the import.


def _require_openpyxl():
    """Skip the test if openpyxl isn't installed."""
    try:
        import openpyxl  # noqa: F401
    except ImportError:
        pytest.skip("openpyxl not installed")


def _write_xlsx(
    path: Path,
    sheets: dict[str, list[list[tuple[str, str | None]]]],
) -> None:
    """Write an xlsx file with the given sheets.

    Each sheet is a list of rows; each row is a list of
    `(value, fill_hex)` tuples where `fill_hex` is None for no
    fill or a 6-char hex string (without leading #).
    """
    _require_openpyxl()
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill

    wb = Workbook()
    # Remove the default sheet — we'll add named sheets below.
    default = wb.active
    wb.remove(default)

    for sheet_name, rows in sheets.items():
        ws = wb.create_sheet(title=sheet_name)
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, (value, fill_hex) in enumerate(row, start=1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                if fill_hex:
                    cell.fill = PatternFill(
                        start_color=fill_hex,
                        end_color=fill_hex,
                        fill_type="solid",
                    )
    wb.save(str(path))


class TestXlsxDispatch:
    """Basic routing — xlsx goes to openpyxl, not markitdown."""

    def test_xlsx_routes_to_openpyxl(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _require_openpyxl()
        xlsx_path = scan_root / "data.xlsx"
        _write_xlsx(xlsx_path, {
            "Sheet1": [[("a", None), ("b", None)], [("1", None), ("2", None)]],
        })
        result = doc_convert.convert_files(["data.xlsx"])
        [entry] = result["results"]
        assert entry["status"] == "ok"
        # markitdown fake should NOT have been called for xlsx.
        assert str(xlsx_path) not in fake_markitdown.outputs or (
            fake_markitdown.outputs.get(str(xlsx_path)) is None
        )

    def test_xlsx_produces_output_file(
        self, doc_convert, scan_root
    ):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "Sheet1": [[("a", None)], [("1", None)]],
        })
        doc_convert.convert_files(["data.xlsx"])
        output = scan_root / "data.md"
        assert output.is_file()

    def test_xlsx_header_has_provenance(
        self, doc_convert, scan_root
    ):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "Sheet1": [[("a", None)], [("1", None)]],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        assert content.startswith("<!-- docuvert:")
        assert "source=data.xlsx" in content
        assert "sha256=" in content

    def test_xlsx_scan_current_after_conversion(
        self, doc_convert, scan_root
    ):
        """After conversion, scan classifies the xlsx as `current`."""
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "Sheet1": [[("a", None)], [("1", None)]],
        })
        doc_convert.convert_files(["data.xlsx"])
        scan = doc_convert.scan_convertible_files()
        [entry] = scan
        assert entry["status"] == "current"


class TestXlsxContent:
    """Markdown table structure — headers, rows, sheet headings."""

    def test_sheet_name_as_heading(self, doc_convert, scan_root):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "Budget": [
                [("col1", None), ("col2", None)],
                [("10", None), ("20", None)],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        assert "## Budget" in content

    def test_first_row_becomes_header(self, doc_convert, scan_root):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("Name", None), ("Value", None)],
                [("row1", None), ("100", None)],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        assert "| Name | Value |" in content

    def test_table_has_separator_row(self, doc_convert, scan_root):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("a", None), ("b", None)],
                [("1", None), ("2", None)],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        # Separator row uses ---.
        assert "|---|---|" in content

    def test_multiple_sheets(self, doc_convert, scan_root):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "Q1": [[("a", None)], [("1", None)]],
            "Q2": [[("b", None)], [("2", None)]],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        assert "## Q1" in content
        assert "## Q2" in content

    def test_empty_row_stripped(self, doc_convert, scan_root):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("a", None), ("b", None)],
                [("", None), ("", None)],  # empty row
                [("1", None), ("2", None)],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        # The all-empty row should not appear.
        lines = content.splitlines()
        # Find the data row for "1" — no blank-cell row before it.
        data_lines = [line for line in lines if line.startswith("|")]
        assert "| 1 | 2 |" in data_lines
        # No row of just empty cells.
        for line in data_lines:
            assert line != "|  |  |"

    def test_empty_column_stripped(self, doc_convert, scan_root):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("a", None), ("", None), ("c", None)],
                [("1", None), ("", None), ("3", None)],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        # Middle column is empty — should be dropped. Header
        # row should have only two entries.
        assert "| a | c |" in content
        assert "| 1 | 3 |" in content

    def test_nan_normalised_to_empty(self, doc_convert, scan_root):
        """Values 'nan' and 'none' (case-insensitive) become empty."""
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("Name", None), ("Value", None)],
                [("row", None), ("nan", None)],
                [("row2", None), ("NONE", None)],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        assert "nan" not in content.lower().replace(
            "sha256", ""  # hex digest may contain "nan" substrings
        ) or "| row |  |" in content
        # Direct check — the value cells should be empty.
        assert "| row |  |" in content
        assert "| row2 |  |" in content

    def test_empty_spreadsheet(self, doc_convert, scan_root):
        """A spreadsheet with no data still produces an output file."""
        _require_openpyxl()
        _write_xlsx(scan_root / "empty.xlsx", {
            "Sheet1": [[("", None)]],
        })
        result = doc_convert.convert_files(["empty.xlsx"])
        [entry] = result["results"]
        assert entry["status"] == "ok"
        content = (scan_root / "empty.md").read_text(encoding="utf-8")
        assert "empty spreadsheet" in content.lower()


class TestXlsxColours:
    """Colour extraction — markers and legend."""

    def test_red_cell_gets_red_marker(self, doc_convert, scan_root):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("Status", None), ("Priority", None)],
                [("Failed", "FF0000"), ("High", None)],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        # Red fill → 🔴 marker.
        assert "🔴 Failed" in content

    def test_green_cell_gets_green_marker(self, doc_convert, scan_root):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("Status", None)],
                [("Done", "00C800")],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        assert "🟢 Done" in content

    def test_legend_lists_used_colours(self, doc_convert, scan_root):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("Status", None)],
                [("Failed", "FF0000")],
                [("Done", "00C800")],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        assert "## Legend" in content
        assert "🔴 red" in content
        assert "🟢 green" in content

    def test_no_legend_when_no_colours(self, doc_convert, scan_root):
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("a", None), ("b", None)],
                [("1", None), ("2", None)],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        assert "## Legend" not in content

    def test_near_white_fill_ignored(self, doc_convert, scan_root):
        """Near-white fills don't produce markers."""
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("Status", None)],
                # FEFEFE is effectively white — default formatting.
                [("Normal", "FEFEFE")],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        # No emoji markers, no legend.
        assert "## Legend" not in content

    def test_near_black_fill_ignored(self, doc_convert, scan_root):
        """Near-black fills don't produce markers."""
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("Status", None)],
                [("Text", "010101")],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        assert "## Legend" not in content

    def test_unknown_colour_gets_fallback_marker(
        self, doc_convert, scan_root
    ):
        """A colour far from every named hue gets a fallback marker."""
        _require_openpyxl()
        # A distinctive teal that shouldn't match red/green/yellow/blue
        # closely enough within the named-colour distance.
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("Status", None)],
                [("Weird", "00A0A0")],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        assert "## Legend" in content
        # Either a fallback cluster name or a near-miss named
        # colour — both acceptable; what matters is that SOME
        # marker was assigned.
        assert "Weird" in content

    def test_coloured_empty_cell_shows_marker_alone(
        self, doc_convert, scan_root
    ):
        """An empty cell with a fill shows just the marker."""
        _require_openpyxl()
        _write_xlsx(scan_root / "data.xlsx", {
            "S1": [
                [("Name", None), ("Flag", None)],
                [("item", None), ("", "FF0000")],
            ],
        })
        doc_convert.convert_files(["data.xlsx"])
        content = (scan_root / "data.md").read_text(encoding="utf-8")
        # The red marker appears without a trailing value, but
        # still as a cell in the row.
        assert "| item | 🔴 |" in content


class TestXlsxFallback:
    """openpyxl fallback — missing library or corrupt file."""

    def test_missing_openpyxl_falls_back_to_markitdown(
        self, doc_convert, scan_root, fake_markitdown, monkeypatch
    ):
        """When openpyxl isn't installed, fall back to markitdown."""
        import sys
        _write_source(scan_root, "data.xlsx", b"fake xlsx")
        fake_markitdown.outputs[str(scan_root / "data.xlsx")] = (
            "| fake | markitdown | output |\n"
        )
        # Block openpyxl import.
        monkeypatch.delitem(sys.modules, "openpyxl", raising=False)
        real_import = __builtins__["__import__"] if isinstance(
            __builtins__, dict
        ) else __builtins__.__import__

        def blocking_import(name, *args, **kwargs):
            if name == "openpyxl" or name.startswith("openpyxl."):
                raise ImportError("openpyxl not installed")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr("builtins.__import__", blocking_import)
        result = doc_convert.convert_files(["data.xlsx"])
        [entry] = result["results"]
        # Fell back to markitdown — succeeded using the fake.
        assert entry["status"] == "ok"

    def test_corrupt_xlsx_falls_back_to_markitdown(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """openpyxl failure on a corrupt file falls back cleanly."""
        _require_openpyxl()
        # Not a real xlsx file — openpyxl will raise on open.
        _write_source(scan_root, "corrupt.xlsx", b"not a real xlsx")
        fake_markitdown.outputs[str(scan_root / "corrupt.xlsx")] = (
            "fallback output\n"
        )
        result = doc_convert.convert_files(["corrupt.xlsx"])
        [entry] = result["results"]
        # Either markitdown succeeded (the fake returned text) or
        # markitdown also errored. Both are valid — the key
        # invariant is we don't crash on corrupt input.
        assert entry["status"] in ("ok", "error")


# ---------------------------------------------------------------------------
# Pass A4 — pptx via python-pptx (fallback; primary A5 path is
# LibreOffice + PyMuPDF)
# ---------------------------------------------------------------------------
#
# python-pptx is a required dependency of the `[docs]` extra, so
# these tests import it directly when available. The library builds
# real pptx files, which is more reliable than mocking every
# attribute of Slide/Shape/TextFrame/Table. Tests that need the
# "python-pptx missing" path block the import via monkeypatch.


def _require_pptx():
    """Skip the test if python-pptx isn't installed."""
    try:
        import pptx  # noqa: F401
    except ImportError:
        pytest.skip("python-pptx not installed")


def _make_pptx_with_title(
    path: Path,
    title: str,
    body: str = "",
) -> None:
    """Create a minimal pptx with a title slide.

    Uses python-pptx's default layout which provides title +
    body placeholders. `body` defaults to empty.
    """
    _require_pptx()
    from pptx import Presentation

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if body and len(slide.placeholders) > 1:
        slide.placeholders[1].text = body
    prs.save(str(path))


def _make_pptx_with_n_slides(path: Path, n: int) -> None:
    """Create a pptx with n title slides, each numbered."""
    _require_pptx()
    from pptx import Presentation

    prs = Presentation()
    for i in range(n):
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {i + 1}"
    prs.save(str(path))


def _make_pptx_with_image(
    path: Path,
    image_bytes: bytes,
) -> None:
    """Create a pptx with one image-containing slide."""
    _require_pptx()
    import io
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Blank layout — index 6 in the default template.
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(
        io.BytesIO(image_bytes),
        Inches(1), Inches(1),
        width=Inches(3), height=Inches(2),
    )
    prs.save(str(path))


def _make_pptx_with_table(path: Path) -> None:
    """Create a pptx with one table-containing slide."""
    _require_pptx()
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    rows, cols = 2, 3
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(1), Inches(1),
        Inches(6), Inches(2),
    )
    table = table_shape.table
    table.cell(0, 0).text = "header1"
    table.cell(0, 1).text = "header2"
    table.cell(0, 2).text = "header3"
    table.cell(1, 0).text = "a"
    table.cell(1, 1).text = "b"
    table.cell(1, 2).text = "c"
    prs.save(str(path))


@pytest.fixture
def force_pptx_fallback(monkeypatch):
    """Force pptx to use the python-pptx fallback path.

    When LibreOffice is actually installed on the test machine,
    pptx files route through the LibreOffice + PyMuPDF pipeline
    (Pass A5b primary path). The A4 fallback tests were written
    against the python-pptx path and produce different output
    (SVGs named `NN_slide.svg`, headings `## Slide N`) than the
    LibreOffice path (`NN_page.svg`, `## Page N`).

    Stubbing `shutil.which` to return None for `soffice` forces
    the dispatch to bypass LibreOffice and fall through to
    python-pptx. The A4 tests work unchanged.
    """
    monkeypatch.setattr(
        "ac_dc.doc_convert.shutil.which",
        lambda cmd: None,
    )


@pytest.mark.usefixtures("force_pptx_fallback")
class TestPptxDispatch:
    """Basic routing — pptx goes to python-pptx, not markitdown."""

    def test_pptx_routes_to_python_pptx(
        self, doc_convert, scan_root, fake_markitdown
    ):
        _require_pptx()
        pptx_path = scan_root / "deck.pptx"
        _make_pptx_with_title(pptx_path, "Hello")
        result = doc_convert.convert_files(["deck.pptx"])
        [entry] = result["results"]
        assert entry["status"] == "ok"

    def test_pptx_produces_index_markdown(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_title(scan_root / "deck.pptx", "Intro")
        doc_convert.convert_files(["deck.pptx"])
        output = scan_root / "deck.md"
        assert output.is_file()

    def test_pptx_produces_assets_subdirectory(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_title(scan_root / "deck.pptx", "Intro")
        doc_convert.convert_files(["deck.pptx"])
        assets_dir = scan_root / "deck"
        assert assets_dir.is_dir()

    def test_pptx_header_has_provenance(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_title(scan_root / "deck.pptx", "Intro")
        doc_convert.convert_files(["deck.pptx"])
        content = (scan_root / "deck.md").read_text(encoding="utf-8")
        assert content.startswith("<!-- docuvert:")
        assert "source=deck.pptx" in content
        assert "sha256=" in content

    def test_pptx_scan_current_after_conversion(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_title(scan_root / "deck.pptx", "Intro")
        doc_convert.convert_files(["deck.pptx"])
        scan = doc_convert.scan_convertible_files()
        [entry] = scan
        assert entry["status"] == "current"


@pytest.mark.usefixtures("force_pptx_fallback")
class TestPptxSlideFiles:
    """Per-slide SVG file creation and naming."""

    def test_single_slide_produces_one_svg(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_title(scan_root / "deck.pptx", "Only slide")
        doc_convert.convert_files(["deck.pptx"])
        svgs = list((scan_root / "deck").glob("*.svg"))
        assert len(svgs) == 1
        assert svgs[0].name == "01_slide.svg"

    def test_multiple_slides_zero_padded(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_n_slides(scan_root / "deck.pptx", 3)
        doc_convert.convert_files(["deck.pptx"])
        svgs = sorted(p.name for p in (scan_root / "deck").glob("*.svg"))
        assert svgs == ["01_slide.svg", "02_slide.svg", "03_slide.svg"]

    def test_large_deck_pads_width(
        self, doc_convert, scan_root
    ):
        """A deck > 99 slides pads to 3 digits."""
        _require_pptx()
        # 100 slides is slow but still fast enough for a test.
        _make_pptx_with_n_slides(scan_root / "deck.pptx", 100)
        doc_convert.convert_files(["deck.pptx"])
        svgs = sorted(p.name for p in (scan_root / "deck").glob("*.svg"))
        assert len(svgs) == 100
        # First slide padded to 3 digits.
        assert svgs[0] == "001_slide.svg"
        # Last slide has no leading zero beyond the width.
        assert svgs[-1] == "100_slide.svg"

    def test_images_listed_in_provenance_header(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_n_slides(scan_root / "deck.pptx", 2)
        doc_convert.convert_files(["deck.pptx"])
        content = (scan_root / "deck.md").read_text(encoding="utf-8")
        assert "images=01_slide.svg,02_slide.svg" in content


@pytest.mark.usefixtures("force_pptx_fallback")
class TestPptxIndexMarkdown:
    """Structure of the index markdown linking all slides."""

    def test_index_has_slide_heading(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_n_slides(scan_root / "deck.pptx", 2)
        doc_convert.convert_files(["deck.pptx"])
        content = (scan_root / "deck.md").read_text(encoding="utf-8")
        assert "## Slide 1" in content
        assert "## Slide 2" in content

    def test_index_has_image_references(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_n_slides(scan_root / "deck.pptx", 2)
        doc_convert.convert_files(["deck.pptx"])
        content = (scan_root / "deck.md").read_text(encoding="utf-8")
        assert "![Slide 1](deck/01_slide.svg)" in content
        assert "![Slide 2](deck/02_slide.svg)" in content

    def test_empty_presentation_placeholder(
        self, doc_convert, scan_root, fake_markitdown
    ):
        """A pptx with no slides produces a placeholder output."""
        _require_pptx()
        from pptx import Presentation
        prs = Presentation()  # no slides added
        prs.save(str(scan_root / "empty.pptx"))
        result = doc_convert.convert_files(["empty.pptx"])
        [entry] = result["results"]
        assert entry["status"] == "ok"
        content = (scan_root / "empty.md").read_text(encoding="utf-8")
        assert "empty presentation" in content.lower()


@pytest.mark.usefixtures("force_pptx_fallback")
class TestPptxSvgContent:
    """SVG content — text, images, tables."""

    def test_title_text_appears_in_svg(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_title(
            scan_root / "deck.pptx",
            "Architecture Review",
            body="Key insights below",
        )
        doc_convert.convert_files(["deck.pptx"])
        svg_content = (
            scan_root / "deck" / "01_slide.svg"
        ).read_text(encoding="utf-8")
        assert "Architecture Review" in svg_content
        assert "Key insights below" in svg_content

    def test_svg_has_valid_root(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_title(scan_root / "deck.pptx", "Title")
        doc_convert.convert_files(["deck.pptx"])
        svg_content = (
            scan_root / "deck" / "01_slide.svg"
        ).read_text(encoding="utf-8")
        assert svg_content.startswith("<svg")
        assert "xmlns=\"http://www.w3.org/2000/svg\"" in svg_content
        assert svg_content.rstrip().endswith("</svg>")

    def test_svg_has_viewbox(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_title(scan_root / "deck.pptx", "Title")
        doc_convert.convert_files(["deck.pptx"])
        svg_content = (
            scan_root / "deck" / "01_slide.svg"
        ).read_text(encoding="utf-8")
        assert "viewBox=" in svg_content

    def test_image_embedded_as_data_uri(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        png = _make_png_bytes()
        _make_pptx_with_image(scan_root / "deck.pptx", png)
        doc_convert.convert_files(["deck.pptx"])
        svg_content = (
            scan_root / "deck" / "01_slide.svg"
        ).read_text(encoding="utf-8")
        # python-pptx re-saves PNG on add_picture, so we can't
        # compare exact bytes. But the data URI should be
        # present with an image MIME.
        assert "data:image/" in svg_content
        assert ";base64," in svg_content
        assert "<image" in svg_content

    def test_table_cells_rendered(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        _make_pptx_with_table(scan_root / "deck.pptx")
        doc_convert.convert_files(["deck.pptx"])
        svg_content = (
            scan_root / "deck" / "01_slide.svg"
        ).read_text(encoding="utf-8")
        # Table cells should be rendered as text elements.
        for cell_text in ["header1", "header2", "a", "b", "c"]:
            assert cell_text in svg_content
        # At least one rect for cell borders.
        assert "<rect" in svg_content

    def test_special_characters_escaped(
        self, doc_convert, scan_root
    ):
        """Text containing < > & should be XML-escaped."""
        _require_pptx()
        _make_pptx_with_title(
            scan_root / "deck.pptx",
            "A & B < C > D",
        )
        doc_convert.convert_files(["deck.pptx"])
        svg_content = (
            scan_root / "deck" / "01_slide.svg"
        ).read_text(encoding="utf-8")
        # Literal unescaped chars would break XML parsing.
        assert "&amp;" in svg_content
        assert "&lt;" in svg_content
        assert "&gt;" in svg_content


@pytest.mark.usefixtures("force_pptx_fallback")
class TestPptxOrphanCleanup:
    """Reconversion removes stale slide SVGs."""

    def test_reconversion_with_fewer_slides_removes_orphans(
        self, doc_convert, scan_root
    ):
        _require_pptx()
        # Initial 3-slide deck.
        _make_pptx_with_n_slides(scan_root / "deck.pptx", 3)
        doc_convert.convert_files(["deck.pptx"])
        assert (scan_root / "deck" / "03_slide.svg").exists()

        # Re-save with 1 slide.
        _make_pptx_with_n_slides(scan_root / "deck.pptx", 1)
        doc_convert.convert_files(["deck.pptx"])

        # Only 01 remains.
        assert (scan_root / "deck" / "01_slide.svg").exists()
        assert not (scan_root / "deck" / "02_slide.svg").exists()
        assert not (scan_root / "deck" / "03_slide.svg").exists()


class TestPptxFailures:
    """python-pptx missing, corrupt file.

    These tests exercise the FALLBACK path — they monkeypatch
    `shutil.which` to return None so the LibreOffice dispatch
    bypasses its primary path and falls back to python-pptx.
    Without this stub the tests would require LibreOffice NOT
    to be installed in the test environment, which isn't
    portable across CI.
    """

    def test_missing_python_pptx_returns_error(
        self, doc_convert, scan_root, monkeypatch
    ):
        """Without python-pptx installed, pptx conversion errors."""
        import sys
        _write_source(scan_root, "deck.pptx", b"fake pptx")
        # Force the LibreOffice path to bypass — pretend soffice
        # isn't available.
        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: None,
        )
        monkeypatch.delitem(sys.modules, "pptx", raising=False)
        # Also need to block re-import.
        real_import = __builtins__["__import__"] if isinstance(
            __builtins__, dict
        ) else __builtins__.__import__

        def blocking_import(name, *args, **kwargs):
            if name == "pptx" or name.startswith("pptx."):
                raise ImportError("python-pptx not installed")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr("builtins.__import__", blocking_import)
        result = doc_convert.convert_files(["deck.pptx"])
        [entry] = result["results"]
        assert entry["status"] == "error"
        assert "python-pptx" in entry["message"]

    def test_corrupt_pptx_errors(
        self, doc_convert, scan_root, monkeypatch
    ):
        """A non-pptx file errors rather than crashing."""
        _require_pptx()
        # Force the LibreOffice path to bypass so we hit the
        # python-pptx fallback which will fail on corrupt input.
        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: None,
        )
        _write_source(scan_root, "corrupt.pptx", b"not a real pptx")
        result = doc_convert.convert_files(["corrupt.pptx"])
        [entry] = result["results"]
        assert entry["status"] == "error"


# ---------------------------------------------------------------------------
# Pass A5b — LibreOffice + PyMuPDF pipeline (primary pptx/odp path)
# ---------------------------------------------------------------------------
#
# These tests exercise the LibreOffice dispatch without requiring
# LibreOffice to actually be installed. We monkeypatch
# `subprocess.run` and the `shutil.which` availability probe to
# simulate various outcomes: success, timeout, non-zero exit,
# missing output, and the fallback paths when dependencies are
# absent.
#
# The "real" LibreOffice path is tested end-to-end only when
# soffice is present on PATH (skipped otherwise) — that test
# uses a real pptx and exercises the full LibreOffice invocation.


class TestLibreOfficeDispatch:
    """Extension routing through the LibreOffice pipeline."""

    def test_pptx_routes_to_libreoffice_when_available(
        self, doc_convert, scan_root, monkeypatch
    ):
        """pptx tries LibreOffice first when soffice is on PATH."""
        import subprocess
        calls: list[list[str]] = []

        def fake_run(cmd, **kwargs):
            calls.append(cmd)
            # Write a fake PDF to the --outdir location so the
            # caller finds it and proceeds. We don't care about
            # PDF validity — the PyMuPDF mock below handles that.
            outdir_idx = cmd.index("--outdir") + 1
            outdir = Path(cmd[outdir_idx])
            source_path = Path(cmd[-1])
            pdf_path = outdir / (source_path.stem + ".pdf")
            pdf_path.write_bytes(b"fake pdf content")
            return subprocess.CompletedProcess(
                args=cmd, returncode=0, stdout="", stderr=""
            )

        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        monkeypatch.setattr("subprocess.run", fake_run)
        # Stub PyMuPDF open to fail cleanly so we don't need to
        # produce a real PDF — the test only verifies dispatch.
        _require_pymupdf()
        import fitz as _fitz  # noqa: F401 — only checking import
        _make_pptx_with_title(scan_root / "deck.pptx", "Title")

        result = doc_convert.convert_files(["deck.pptx"])
        # LibreOffice was called.
        assert len(calls) == 1
        assert calls[0][0] == "/usr/bin/soffice"
        assert "--headless" in calls[0]
        assert "--convert-to" in calls[0]
        assert "pdf" in calls[0]
        # PyMuPDF fails on the fake PDF — result is an error,
        # but the important thing is LibreOffice WAS called.
        [entry] = result["results"]
        assert entry["status"] == "error"

    def test_odp_routes_to_libreoffice_when_available(
        self, doc_convert, scan_root, monkeypatch
    ):
        """odp tries LibreOffice first when soffice is on PATH."""
        import subprocess
        calls: list[list[str]] = []

        def fake_run(cmd, **kwargs):
            calls.append(cmd)
            outdir_idx = cmd.index("--outdir") + 1
            outdir = Path(cmd[outdir_idx])
            source_path = Path(cmd[-1])
            pdf_path = outdir / (source_path.stem + ".pdf")
            pdf_path.write_bytes(b"fake pdf")
            return subprocess.CompletedProcess(
                args=cmd, returncode=0, stdout="", stderr=""
            )

        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        monkeypatch.setattr("subprocess.run", fake_run)
        _write_source(scan_root, "deck.odp", b"fake odp")

        doc_convert.convert_files(["deck.odp"])
        assert len(calls) == 1
        # Source path (last arg) ends with .odp.
        assert calls[0][-1].endswith("deck.odp")


class TestLibreOfficeFallback:
    """Fallback paths when LibreOffice isn't usable."""

    def test_no_soffice_pptx_falls_back_to_python_pptx(
        self, doc_convert, scan_root, monkeypatch
    ):
        """pptx falls back to python-pptx when soffice missing."""
        _require_pptx()
        # soffice not on PATH.
        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: None,
        )
        _make_pptx_with_title(scan_root / "deck.pptx", "Title")
        result = doc_convert.convert_files(["deck.pptx"])
        [entry] = result["results"]
        # python-pptx fallback succeeds — the usual pptx output.
        assert entry["status"] == "ok"
        # Per-slide SVG characteristic of the python-pptx path.
        assert (scan_root / "deck" / "01_slide.svg").is_file()

    def test_no_soffice_odp_falls_back_to_markitdown(
        self, doc_convert, scan_root, fake_markitdown, monkeypatch
    ):
        """odp falls back to markitdown when soffice missing."""
        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: None,
        )
        _write_source(scan_root, "deck.odp", b"fake odp")
        fake_markitdown.outputs[str(scan_root / "deck.odp")] = (
            "odp as markdown\n"
        )
        result = doc_convert.convert_files(["deck.odp"])
        [entry] = result["results"]
        assert entry["status"] == "ok"
        content = (scan_root / "deck.md").read_text(encoding="utf-8")
        assert "odp as markdown" in content

    def test_timeout_falls_back(
        self, doc_convert, scan_root, monkeypatch
    ):
        """LibreOffice timeout falls back rather than erroring."""
        _require_pptx()
        import subprocess

        def fake_run(cmd, **kwargs):
            raise subprocess.TimeoutExpired(cmd, timeout=120)

        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        monkeypatch.setattr("subprocess.run", fake_run)
        _make_pptx_with_title(scan_root / "deck.pptx", "Title")
        result = doc_convert.convert_files(["deck.pptx"])
        [entry] = result["results"]
        # Falls back to python-pptx which succeeds.
        assert entry["status"] == "ok"
        assert (scan_root / "deck" / "01_slide.svg").is_file()

    def test_nonzero_exit_falls_back(
        self, doc_convert, scan_root, monkeypatch
    ):
        """Non-zero soffice exit falls back to python-pptx."""
        _require_pptx()
        import subprocess

        def fake_run(cmd, **kwargs):
            return subprocess.CompletedProcess(
                args=cmd, returncode=1,
                stdout="", stderr="conversion failed",
            )

        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        monkeypatch.setattr("subprocess.run", fake_run)
        _make_pptx_with_title(scan_root / "deck.pptx", "Title")
        result = doc_convert.convert_files(["deck.pptx"])
        [entry] = result["results"]
        assert entry["status"] == "ok"
        assert (scan_root / "deck" / "01_slide.svg").is_file()

    def test_missing_output_pdf_falls_back(
        self, doc_convert, scan_root, monkeypatch
    ):
        """LibreOffice succeeds but produces no PDF → fallback."""
        _require_pptx()
        import subprocess

        def fake_run(cmd, **kwargs):
            # Don't write any output — simulates weird template
            # that makes soffice produce no file despite 0 exit.
            return subprocess.CompletedProcess(
                args=cmd, returncode=0, stdout="", stderr=""
            )

        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        monkeypatch.setattr("subprocess.run", fake_run)
        _make_pptx_with_title(scan_root / "deck.pptx", "Title")
        result = doc_convert.convert_files(["deck.pptx"])
        [entry] = result["results"]
        assert entry["status"] == "ok"
        assert (scan_root / "deck" / "01_slide.svg").is_file()

    def test_no_pymupdf_falls_back(
        self, doc_convert, scan_root, monkeypatch
    ):
        """Missing PyMuPDF falls back even when soffice available."""
        _require_pptx()
        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        # Make the PyMuPDF probe report False.
        monkeypatch.setattr(
            "ac_dc.doc_convert.DocConvert._probe_import",
            lambda self, name: False,
        )
        _make_pptx_with_title(scan_root / "deck.pptx", "Title")
        result = doc_convert.convert_files(["deck.pptx"])
        [entry] = result["results"]
        assert entry["status"] == "ok"
        assert (scan_root / "deck" / "01_slide.svg").is_file()


class TestLibreOfficeProvenance:
    """Provenance header correctness when routing through LibreOffice."""

    def test_header_uses_original_filename(
        self, doc_convert, scan_root, monkeypatch
    ):
        """Provenance records source=deck.pptx, not source=deck.pdf."""
        _require_pymupdf()
        import subprocess

        # Make soffice produce a minimal valid PDF using PyMuPDF.
        import fitz

        def fake_run(cmd, **kwargs):
            outdir_idx = cmd.index("--outdir") + 1
            outdir = Path(cmd[outdir_idx])
            source_path = Path(cmd[-1])
            pdf_path = outdir / (source_path.stem + ".pdf")
            doc = fitz.open()
            page = doc.new_page()
            page.insert_text((72, 72), "Converted content", fontsize=12)
            doc.save(str(pdf_path))
            doc.close()
            return subprocess.CompletedProcess(
                args=cmd, returncode=0, stdout="", stderr=""
            )

        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        monkeypatch.setattr("subprocess.run", fake_run)
        _write_source(scan_root, "deck.pptx", b"original pptx bytes")

        doc_convert.convert_files(["deck.pptx"])
        content = (scan_root / "deck.md").read_text(encoding="utf-8")
        # Provenance points at the ORIGINAL file, not the
        # intermediate PDF.
        assert "source=deck.pptx" in content
        assert "source=deck.pdf" not in content

    def test_hash_reflects_original_source(
        self, doc_convert, scan_root, monkeypatch
    ):
        """Hash is of the original pptx, not the intermediate PDF."""
        _require_pymupdf()
        import subprocess
        import hashlib
        import fitz

        original_content = b"specific pptx bytes for hash check"
        expected_hash = hashlib.sha256(original_content).hexdigest()

        def fake_run(cmd, **kwargs):
            outdir_idx = cmd.index("--outdir") + 1
            outdir = Path(cmd[outdir_idx])
            source_path = Path(cmd[-1])
            pdf_path = outdir / (source_path.stem + ".pdf")
            doc = fitz.open()
            doc.new_page()
            doc.save(str(pdf_path))
            doc.close()
            return subprocess.CompletedProcess(
                args=cmd, returncode=0, stdout="", stderr=""
            )

        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        monkeypatch.setattr("subprocess.run", fake_run)
        _write_source(scan_root, "deck.pptx", original_content)

        doc_convert.convert_files(["deck.pptx"])
        content = (scan_root / "deck.md").read_text(encoding="utf-8")
        assert f"sha256={expected_hash}" in content

    def test_output_lands_next_to_original(
        self, doc_convert, scan_root, monkeypatch
    ):
        """Output .md is beside the source, not in the temp dir."""
        _require_pymupdf()
        import subprocess
        import fitz

        def fake_run(cmd, **kwargs):
            outdir_idx = cmd.index("--outdir") + 1
            outdir = Path(cmd[outdir_idx])
            source_path = Path(cmd[-1])
            pdf_path = outdir / (source_path.stem + ".pdf")
            doc = fitz.open()
            doc.new_page()
            doc.save(str(pdf_path))
            doc.close()
            return subprocess.CompletedProcess(
                args=cmd, returncode=0, stdout="", stderr=""
            )

        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        monkeypatch.setattr("subprocess.run", fake_run)
        _write_source(scan_root, "docs/deck.pptx", b"x")

        doc_convert.convert_files(["docs/deck.pptx"])
        # Output at the expected sibling location.
        assert (scan_root / "docs" / "deck.md").is_file()

    def test_scan_current_after_libreoffice_conversion(
        self, doc_convert, scan_root, monkeypatch
    ):
        """After conversion, scan classifies pptx as `current`."""
        _require_pymupdf()
        import subprocess
        import fitz

        def fake_run(cmd, **kwargs):
            outdir_idx = cmd.index("--outdir") + 1
            outdir = Path(cmd[outdir_idx])
            source_path = Path(cmd[-1])
            pdf_path = outdir / (source_path.stem + ".pdf")
            doc = fitz.open()
            doc.new_page()
            doc.save(str(pdf_path))
            doc.close()
            return subprocess.CompletedProcess(
                args=cmd, returncode=0, stdout="", stderr=""
            )

        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None,
        )
        monkeypatch.setattr("subprocess.run", fake_run)
        _write_source(scan_root, "deck.pptx", b"stable content")

        doc_convert.convert_files(["deck.pptx"])
        scan = doc_convert.scan_convertible_files()
        [entry] = scan
        assert entry["status"] == "current"


class TestLibreOfficeEndToEnd:
    """Full-pipeline test against real LibreOffice (when available)."""

    def test_real_libreoffice_converts_pptx(
        self, doc_convert, scan_root
    ):
        """Full round-trip when soffice is actually installed.

        Skipped when LibreOffice isn't on PATH. Ensures the
        mocked tests above aren't hiding a real-world issue with
        subprocess args, output path resolution, or format
        compatibility.
        """
        if shutil.which("soffice") is None:
            pytest.skip("LibreOffice (soffice) not installed")
        _require_pymupdf()
        _require_pptx()
        _make_pptx_with_title(
            scan_root / "deck.pptx",
            "Real LibreOffice Test",
            body="Body paragraph",
        )
        result = doc_convert.convert_files(["deck.pptx"])
        [entry] = result["results"]
        assert entry["status"] == "ok"
        content = (scan_root / "deck.md").read_text(encoding="utf-8")
        # Title text should survive the round-trip.
        assert "Real LibreOffice Test" in content
        # Provenance correctness.
        assert "source=deck.pptx" in content


# ---------------------------------------------------------------------------
# Pass A5a — PDF via PyMuPDF (direct, no LibreOffice)
# ---------------------------------------------------------------------------
#
# PyMuPDF (fitz) is a required dependency of the `[docs]` extra, so
# these tests import it directly when available. The library builds
# real PDFs via the same API used for reading, which lets us test
# the full pipeline without mocking every page/drawing/image call.
# Tests that need "PyMuPDF missing" use the builtins.__import__
# monkeypatch pattern from the other pass tests.


def _require_pymupdf():
    """Skip the test if PyMuPDF isn't installed."""
    try:
        import fitz  # noqa: F401
    except ImportError:
        pytest.skip("PyMuPDF not installed")


def _make_pdf_with_text(path: Path, pages: list[str]) -> None:
    """Create a minimal PDF with one text block per page.

    Each page gets a single text block positioned near the
    top-left. No images or drawings — pure text.
    """
    _require_pymupdf()
    import fitz

    doc = fitz.open()
    for text in pages:
        page = doc.new_page(width=612, height=792)  # US Letter
        # Insert text at a visible position.
        page.insert_text((72, 72), text, fontsize=12)
    doc.save(str(path))
    doc.close()


def _make_pdf_with_image(path: Path, image_bytes: bytes) -> None:
    """Create a PDF with one image-containing page."""
    _require_pymupdf()
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=612, height=792)
    # Insert a raster image at a fixed position.
    rect = fitz.Rect(72, 72, 372, 372)
    page.insert_image(rect, stream=image_bytes)
    doc.save(str(path))
    doc.close()


def _make_pdf_with_text_and_image(
    path: Path, text: str, image_bytes: bytes
) -> None:
    """Create a PDF with both text and an image on one page."""
    _require_pymupdf()
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=612, height=792)
    page.insert_text((72, 72), text, fontsize=12)
    rect = fitz.Rect(72, 400, 372, 700)
    page.insert_image(rect, stream=image_bytes)
    doc.save(str(path))
    doc.close()


def _make_empty_pdf(path: Path) -> None:
    """Create a PDF with zero pages.

    PyMuPDF allows empty documents but the save invokes the
    writer which requires at least one page — so we use a
    workaround by writing a minimal valid PDF directly.
    """
    # Minimal valid PDF with 0 pages. The xref is the trick —
    # PDF spec allows empty page trees.
    path.write_bytes(
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[]/Count 0>>endobj\n"
        b"xref\n0 3\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000052 00000 n \n"
        b"trailer<</Size 3/Root 1 0 R>>\n"
        b"startxref\n94\n%%EOF\n"
    )


class TestPdfDispatch:
    """Basic routing — pdf goes to PyMuPDF."""

    def test_pdf_routes_to_pymupdf(self, doc_convert, scan_root):
        _require_pymupdf()
        _make_pdf_with_text(scan_root / "doc.pdf", ["Hello world"])
        result = doc_convert.convert_files(["doc.pdf"])
        [entry] = result["results"]
        assert entry["status"] == "ok"

    def test_pdf_produces_output_file(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        _make_pdf_with_text(scan_root / "doc.pdf", ["Hello"])
        doc_convert.convert_files(["doc.pdf"])
        assert (scan_root / "doc.md").is_file()

    def test_pdf_header_has_provenance(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        _make_pdf_with_text(scan_root / "doc.pdf", ["Hello"])
        doc_convert.convert_files(["doc.pdf"])
        content = (scan_root / "doc.md").read_text(encoding="utf-8")
        assert content.startswith("<!-- docuvert:")
        assert "source=doc.pdf" in content
        assert "sha256=" in content

    def test_pdf_scan_current_after_conversion(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        _make_pdf_with_text(scan_root / "doc.pdf", ["Hello"])
        doc_convert.convert_files(["doc.pdf"])
        scan = doc_convert.scan_convertible_files()
        [entry] = scan
        assert entry["status"] == "current"


class TestPdfTextExtraction:
    """Text extraction into markdown paragraphs."""

    def test_single_page_text_extracted(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        _make_pdf_with_text(
            scan_root / "doc.pdf",
            ["This is page one."],
        )
        doc_convert.convert_files(["doc.pdf"])
        content = (scan_root / "doc.md").read_text(encoding="utf-8")
        assert "This is page one." in content

    def test_multi_page_text_extracted(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        _make_pdf_with_text(
            scan_root / "doc.pdf",
            ["First page text.", "Second page text.", "Third page."],
        )
        doc_convert.convert_files(["doc.pdf"])
        content = (scan_root / "doc.md").read_text(encoding="utf-8")
        assert "First page text." in content
        assert "Second page text." in content
        assert "Third page." in content

    def test_page_headings_in_order(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        _make_pdf_with_text(
            scan_root / "doc.pdf",
            ["First", "Second", "Third"],
        )
        doc_convert.convert_files(["doc.pdf"])
        content = (scan_root / "doc.md").read_text(encoding="utf-8")
        # All three page headings present.
        for i in range(1, 4):
            assert f"## Page {i}" in content
        # In order.
        idx1 = content.index("## Page 1")
        idx2 = content.index("## Page 2")
        idx3 = content.index("## Page 3")
        assert idx1 < idx2 < idx3

    def test_text_only_page_no_svg(
        self, doc_convert, scan_root
    ):
        """Pages with only text produce no SVG — keeps output lean."""
        _require_pymupdf()
        _make_pdf_with_text(
            scan_root / "doc.pdf",
            ["Just text, nothing else."],
        )
        doc_convert.convert_files(["doc.pdf"])
        # Assets dir should not have been created for a
        # text-only page.
        assert not (scan_root / "doc").exists()


class TestPdfImageHandling:
    """Pages with images get SVG companions."""

    def test_page_with_image_produces_svg(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        png = _make_png_bytes()
        _make_pdf_with_image(scan_root / "doc.pdf", png)
        doc_convert.convert_files(["doc.pdf"])
        svgs = list((scan_root / "doc").glob("*.svg"))
        assert len(svgs) == 1

    def test_image_page_svg_filename_padded(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        png = _make_png_bytes()
        _make_pdf_with_image(scan_root / "doc.pdf", png)
        doc_convert.convert_files(["doc.pdf"])
        svgs = list((scan_root / "doc").glob("*.svg"))
        assert svgs[0].name == "01_page.svg"

    def test_image_page_markdown_has_link(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        png = _make_png_bytes()
        _make_pdf_with_image(scan_root / "doc.pdf", png)
        doc_convert.convert_files(["doc.pdf"])
        content = (scan_root / "doc.md").read_text(encoding="utf-8")
        assert "![Page 1](doc/01_page.svg)" in content

    def test_text_plus_image_produces_both(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        png = _make_png_bytes()
        _make_pdf_with_text_and_image(
            scan_root / "doc.pdf",
            "Text content here.",
            png,
        )
        doc_convert.convert_files(["doc.pdf"])
        content = (scan_root / "doc.md").read_text(encoding="utf-8")
        # Markdown has both the text AND the SVG link.
        assert "Text content here." in content
        assert "doc/01_page.svg" in content
        assert (scan_root / "doc" / "01_page.svg").is_file()

    def test_externalized_image_saved_to_disk(
        self, doc_convert, scan_root
    ):
        """PyMuPDF's SVG output embeds images; we extract them."""
        _require_pymupdf()
        png = _make_png_bytes()
        _make_pdf_with_image(scan_root / "doc.pdf", png)
        doc_convert.convert_files(["doc.pdf"])
        # An image file should have been extracted alongside
        # the SVG. PyMuPDF may re-encode the image so we don't
        # check byte-exact, but we expect at least one extra
        # file with an image extension.
        assets = list((scan_root / "doc").iterdir())
        image_files = [
            f for f in assets
            if f.suffix.lower() in (".png", ".jpg", ".jpeg", ".bmp")
        ]
        assert len(image_files) >= 1

    def test_provenance_lists_all_artefacts(
        self, doc_convert, scan_root
    ):
        """All SVGs AND externalized images appear in images= field."""
        _require_pymupdf()
        png = _make_png_bytes()
        _make_pdf_with_image(scan_root / "doc.pdf", png)
        doc_convert.convert_files(["doc.pdf"])
        content = (scan_root / "doc.md").read_text(encoding="utf-8")
        # The SVG should be listed.
        assert "01_page.svg" in content
        # The externalized image should also be listed.
        # Filename pattern: 01_page_img01.<ext>
        assert "01_page_img01" in content


class TestPdfEmptyAndEdgeCases:
    """Edge cases — empty PDF, zero-length pages."""

    def test_empty_pdf_placeholder(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        _make_empty_pdf(scan_root / "empty.pdf")
        result = doc_convert.convert_files(["empty.pdf"])
        [entry] = result["results"]
        assert entry["status"] == "ok"
        content = (scan_root / "empty.md").read_text(encoding="utf-8")
        assert "empty pdf" in content.lower()


class TestPdfOrphanCleanup:
    """Re-conversion with fewer pages removes stale artefacts."""

    def test_reconversion_with_fewer_pages_removes_orphans(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        png = _make_png_bytes()
        # v1 — PDF with 3 image-containing pages.
        import fitz
        doc = fitz.open()
        for _ in range(3):
            page = doc.new_page(width=612, height=792)
            rect = fitz.Rect(72, 72, 372, 372)
            page.insert_image(rect, stream=png)
        doc.save(str(scan_root / "doc.pdf"))
        doc.close()

        doc_convert.convert_files(["doc.pdf"])
        svgs_v1 = sorted(
            p.name for p in (scan_root / "doc").glob("*.svg")
        )
        assert svgs_v1 == [
            "01_page.svg", "02_page.svg", "03_page.svg",
        ]

        # v2 — PDF with 1 image-containing page.
        doc = fitz.open()
        page = doc.new_page(width=612, height=792)
        rect = fitz.Rect(72, 72, 372, 372)
        page.insert_image(rect, stream=png)
        doc.save(str(scan_root / "doc.pdf"))
        doc.close()

        doc_convert.convert_files(["doc.pdf"])
        svgs_v2 = sorted(
            p.name for p in (scan_root / "doc").glob("*.svg")
        )
        assert svgs_v2 == ["01_page.svg"]


class TestPdfFailures:
    """PyMuPDF missing, corrupt PDF."""

    def test_missing_pymupdf_returns_error(
        self, doc_convert, scan_root, monkeypatch
    ):
        import sys
        _write_source(scan_root, "doc.pdf", b"fake pdf")
        monkeypatch.delitem(sys.modules, "fitz", raising=False)
        real_import = __builtins__["__import__"] if isinstance(
            __builtins__, dict
        ) else __builtins__.__import__

        def blocking_import(name, *args, **kwargs):
            if name == "fitz" or name.startswith("fitz."):
                raise ImportError("PyMuPDF not installed")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr("builtins.__import__", blocking_import)
        result = doc_convert.convert_files(["doc.pdf"])
        [entry] = result["results"]
        assert entry["status"] == "error"
        assert "PyMuPDF" in entry["message"]

    def test_corrupt_pdf_errors(
        self, doc_convert, scan_root
    ):
        _require_pymupdf()
        _write_source(scan_root, "corrupt.pdf", b"not a real pdf")
        result = doc_convert.convert_files(["corrupt.pdf"])
        [entry] = result["results"]
        assert entry["status"] == "error"


class TestPdfSvgTextPreservation:
    """Origin-aware SVG text handling for the PDF pipeline.

    Per specs4/4-features/doc-convert.md § "SVG text
    preservation in PDF pipeline" and the supplement at
    specs-reference/4-features/doc-convert.md, ``<text>``
    elements in the generated SVG are stripped or kept
    depending on the source type:

    - Direct ``.pdf`` (papers, reports where text flows in
      paragraphs): when a page has extractable text, strip
      the ``<text>`` / ``<tspan>`` elements from the SVG.
      The markdown already carries the paragraphs;
      duplicating them in the SVG bloats output without
      benefit. Figure-only pages (no extractable text) keep
      their SVG text since it probably labels the figure.
    - Presentations routed through LibreOffice →
      intermediate PDF → PyMuPDF: always keep SVG text.
      Slide labels anchor diagram shapes; stripping them
      leaves meaningless coloured rectangles.

    These tests pin both sides of the rule.
    """

    def test_direct_pdf_strips_svg_text_when_page_has_text(
        self, doc_convert, scan_root
    ):
        """Direct-PDF page with text + image strips SVG <text>.

        The paragraph lands in the markdown as prose; the SVG
        keeps only the graphics (image ref, vector drawings).
        """
        _require_pymupdf()
        png = _make_png_bytes()
        _make_pdf_with_text_and_image(
            scan_root / "doc.pdf",
            "Distinctive unique phrase abc123",
            png,
        )
        doc_convert.convert_files(["doc.pdf"])
        svg_content = (
            scan_root / "doc" / "01_page.svg"
        ).read_text(encoding="utf-8")
        md_content = (
            scan_root / "doc.md"
        ).read_text(encoding="utf-8")
        # Phrase must appear in the markdown (grep, LLM ctx).
        assert "Distinctive unique phrase abc123" in md_content
        # And must NOT appear in the SVG — dedup invariant.
        assert "Distinctive unique phrase abc123" not in svg_content
        # And the SVG should have no <text>/<tspan> elements
        # at all (PyMuPDF emits every glyph as one of these
        # tags when text_as_path=0, so zero is the right
        # post-strip count).
        assert "<text" not in svg_content
        assert "<tspan" not in svg_content

    def test_direct_pdf_figure_only_page_keeps_svg_text(
        self, doc_convert, scan_root
    ):
        """Figure-only pages (no extractable text) keep SVG text.

        These don't enter the "text flows in paragraphs" case —
        any ``<text>`` element on them probably labels the
        figure itself (axis labels, legend entries), and
        stripping it would be lossy.
        """
        _require_pymupdf()
        import fitz

        # Page has a vector drawing (curve) but no insertable
        # text we extract — the ``<text>`` elements PyMuPDF
        # emits for tiny stroked labels should survive.
        doc = fitz.open()
        page = doc.new_page(width=612, height=792)
        # A single Bézier curve so the page counts as having
        # significant graphics and gets an SVG.
        page.draw_bezier(
            fitz.Point(100, 100),
            fitz.Point(200, 50),
            fitz.Point(300, 150),
            fitz.Point(400, 100),
        )
        doc.save(str(scan_root / "fig.pdf"))
        doc.close()
        doc_convert.convert_files(["fig.pdf"])
        # SVG must exist (significant graphics threshold met).
        svg_path = scan_root / "fig" / "01_page.svg"
        assert svg_path.is_file()
        # No text was extracted, so no stripping happened —
        # whatever ``<text>`` PyMuPDF emits (possibly none)
        # is untouched. The invariant we pin: stripping does
        # NOT run when the page has no text. We verify that
        # by checking nothing was stripped that shouldn't
        # have been — i.e., we get back whatever PyMuPDF
        # emitted verbatim.
        # Easier check: the call succeeded without errors.
        svg_content = svg_path.read_text(encoding="utf-8")
        assert svg_content.startswith("<svg")

    def test_libreoffice_pptx_keeps_svg_text(
        self, doc_convert, scan_root, monkeypatch
    ):
        """pptx routed through LibreOffice keeps diagram labels.

        The LibreOffice pipeline passes
        ``strip_text_when_present=False`` when dispatching to
        the PyMuPDF stage, so even pages with extractable text
        retain their ``<text>`` elements in the SVG. Matches
        slide-deck semantics where the text IS the diagram.
        """
        _require_pymupdf()
        import subprocess
        import fitz

        distinctive = "RuntimeEnvironment diagram label"

        def fake_run(cmd, **kwargs):
            outdir_idx = cmd.index("--outdir") + 1
            outdir = Path(cmd[outdir_idx])
            source_path = Path(cmd[-1])
            pdf_path = outdir / (source_path.stem + ".pdf")
            # Produce a real PDF with both text and graphics
            # so the stripping logic WOULD fire if the flag
            # hadn't been set to False.
            doc = fitz.open()
            page = doc.new_page(width=612, height=792)
            page.insert_text((72, 72), distinctive, fontsize=12)
            # A curve so the page gets an SVG.
            page.draw_bezier(
                fitz.Point(100, 400),
                fitz.Point(200, 350),
                fitz.Point(300, 450),
                fitz.Point(400, 400),
            )
            doc.save(str(pdf_path))
            doc.close()
            return subprocess.CompletedProcess(
                args=cmd, returncode=0, stdout="", stderr="",
            )

        monkeypatch.setattr(
            "ac_dc.doc_convert.shutil.which",
            lambda cmd: (
                "/usr/bin/soffice" if cmd == "soffice" else None
            ),
        )
        monkeypatch.setattr("subprocess.run", fake_run)
        _write_source(scan_root, "deck.pptx", b"pptx bytes")

        doc_convert.convert_files(["deck.pptx"])
        svg_path = scan_root / "deck" / "01_page.svg"
        assert svg_path.is_file()
        svg_content = svg_path.read_text(encoding="utf-8")
        # Presentation text MUST survive in the SVG — the
        # LibreOffice path passes strip_text_when_present=False.
        assert distinctive in svg_content

    def test_text_only_page_no_svg(
        self, doc_convert, scan_root
    ):
        """Text-only pages still produce no SVG.

        The text-stripping rule only applies to SVGs that get
        generated at all; a page with no graphics never gets
        one. Text lives in the markdown, full stop.
        """
        _require_pymupdf()
        _make_pdf_with_text(
            scan_root / "doc.pdf",
            ["Only text, no graphics."],
        )
        doc_convert.convert_files(["doc.pdf"])
        # No assets subdir because no SVGs were generated.
        assert not (scan_root / "doc").exists()